
import pandas as pd
import numpy as np
import xlwings as xw
from openpyxl.utils import range_boundaries
from pptx import Presentation
import ast
import datetime as dt
import time
import traceback
import warnings
import subprocess
import cv2
from pdf2image import convert_from_path
from PIL import Image
from io import BytesIO
import xlrd
import os
import sys
import re
import shutil
import pathlib
import sqlite3

#################################################################################################################################
#TEMPLATES DEL APP
#################################################################################################################################

template_ico_app_tapar_pluma_tkinter = os.path.join(sys._MEIPASS, "ico_app_tapar_pluma_tkinter.ico") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "ico_app_tapar_pluma_tkinter.ico")
template_plantilla_xls_config = os.path.join(sys._MEIPASS, "PLANTILLA_CONFIG.xlsx") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "PLANTILLA_CONFIG.xlsx")
template_pdf_guia_usuario = os.path.join(sys._MEIPASS, "pdf_guia_usuario.pdf") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "pdf_guia_usuario.pdf")

template_img_guia_usuario = os.path.join(sys._MEIPASS, "img_guia_usuario.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_guia_usuario.png")
template_img_config = os.path.join(sys._MEIPASS, "img_config.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_config.png")
template_img_boton_ver = os.path.join(sys._MEIPASS, "img_boton_ver.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_boton_ver.png")
template_img_accion_pptx = os.path.join(sys._MEIPASS, "img_accion_pptx.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_accion_pptx.png")
template_img_boton_add = os.path.join(sys._MEIPASS, "img_boton_add.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_boton_add.png")
template_img_boton_clear = os.path.join(sys._MEIPASS, "img_boton_clear.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_boton_clear.png")
template_img_guardar = os.path.join(sys._MEIPASS, "img_guardar.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_guardar.png")
template_img_abrir_fichero = os.path.join(sys._MEIPASS, "img_abrir_fichero.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_abrir_fichero.png")
template_img_update_datos_config = os.path.join(sys._MEIPASS, "img_update_datos_config.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_update_datos_config.png")

template_img_clean_rango_celdas = os.path.join(sys._MEIPASS, "img_clean_rango_celdas.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_clean_rango_celdas.png")

template_img_guardar_id = os.path.join(sys._MEIPASS, "img_guardar_id.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_guardar_id.png")
template_img_eliminar_id = os.path.join(sys._MEIPASS, "img_eliminar_id.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_eliminar_id.png")
template_img_guardar_ruta = os.path.join(sys._MEIPASS, "img_guardar_ruta.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_guardar_ruta.png")
template_img_eliminar_ruta = os.path.join(sys._MEIPASS, "img_eliminar_ruta.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_eliminar_ruta.png")
template_img_update_ruta = os.path.join(sys._MEIPASS, "img_update_ruta.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_update_ruta.png")

template_img_messagebox_askokcancel = os.path.join(sys._MEIPASS, "img_messagebox_askokcancel.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_messagebox_askokcancel.png")
template_img_messagebox_showwarning = os.path.join(sys._MEIPASS, "img_messagebox_showwarning.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_messagebox_showwarning.png")
template_img_messagebox_showerror = os.path.join(sys._MEIPASS, "img_messagebox_showerror.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_messagebox_showerror.png")
template_img_messagebox_showinfo = os.path.join(sys._MEIPASS, "img_messagebox_showinfo.png") if getattr(sys, 'frozen', False) else os.path.join(pathlib.Path(__file__).parent.absolute(), "img_messagebox_showinfo.png")


#################################################################################################################################
#VARIABLES GLOBALES VARIAS
#################################################################################################################################


nombre_app = "SEMI-AUTOMATIZADOR POWERPOINT"
nombre_sistema_sqlite = "SISTEMA_CONFIG_APP_PPTX"
resolucion_pantalla_recomendada = "1920x1080"

opcion_gui_guia_usuario = "GUIA_USUARIO"
opcion_gui_config_id_pptx = "CONFIG_ID_PPTX"
nombre_guia_usuario = "APP_SEMI_AUTOMATIZADOR_POWERPOINT_GUIA_USUARIO"


nombre_id_pptx = "PPTX"
nombre_id_xls = "XLS"

nomenclatura_nombres_screenshots = "APP_PPTX_NOMBRE_PANTALLAZOS"
nombre_carpeta_screenshots = "APP_PTX_PANTALLAZOS"
nombre_carpeta_screenshots_muestra = "APP_PTX_PANTALLAZOS_MUESTRA"
nombre_log_warning_errores_procesos_app = "LOG_WARNINGS_Y_ERRORES"
nombre_pptx_tras_config = "APP_PTX_PANTALLAZOS_CONFIG_PASO_1"
nombre_pptx_tras_ejecucion = "APP_PTX_PANTALLAZOS"
dpi_pdf2image_convert_from_path = 300
opencv_recortar_png_casi_blancos = 245
time_sleep_prudencial_poppler = 0.1

nombre_xls_descarga_parametrica = "APP_PPTX_PARAMETRICA"


lista_opciones_id_xls_combobox_actualizar_vinculos = ["Si", "No"]
lista_opciones_rangos_celdas_xls_combobox_seleccion_1 = ["TODO", "hoja xls", "slide pptx"]


global_proceso_en_ejecucion = "NO"
global_path_sistema_sqlite = None
global_df_parametrica = None
global_df_treeview_id_pptx = None
global_dicc_datos_id_pptx = None
global_poppler_path = None
global_dicc_tablas_y_campos_sistema = None
global_ruta_local_config_sistema_sqlite = None
global_ruta_local_screenshots_png_muestra = None
global_dicc_check_screenshots_si_todo_blanco = {}
global_lista_dicc_errores = []
global_lista_dicc_warning = []



tupla_coordenadas_colocacion_pptx_por_defecto = (278000, 302000, 968000, 1087000)
tupla_resize_pixeles_screenshot_muestra = (600, 600)


dicc_modelo_warning_y_errores_procesos_app = {"ERROR": 
                                                        {"TIPO_LOG": "ERROR"
                                                        , "PROCESO": None
                                                        , "FASE_PROCESO": None
                                                        , "ID_PPTX": None
                                                        , "ID_XLS": None
                                                        , "RUTA_FICHERO": None
                                                        , "HOJA_XLS": None
                                                        , "RANGO_CELDAS": None
                                                        , "SLIDE_PPTX": None
                                                        , "RESUMEN_ERROR": None
                                                        , "MODULO_PYTHON": None
                                                        , "RUTINA_PYTHON": None
                                                        , "LINEA_CODIGO_PYTHON": None
                                                        , "TRACEBACK": None
                                                        }

                                                , "WARNING": {"TIPO_LOG": "WARNING"
                                                            , "PROCESO": None
                                                            , "FASE_PROCESO": None
                                                            , "ID_PPTX": None
                                                            , "ID_XLS": None
                                                            , "HOJA_XLS": None
                                                            , "RANGO_CELDAS": None
                                                            , "SLIDE_PPTX": None
                                                            , "COMENTARIO": None
                                                            }
                                                }



dicc_tabla_config_sistema = {"PARAMETRICA":
                                            {"TABLA": "T_PARAMETRICA"
                                            , "DICC_CAMPOS":
                                                            {"ID_PPTX":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": True
                                                                        , "IMPORT_SISTEMA": True
                                                                        , "TREEVIEW_ID_PPTX": True
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": True
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": False
                                                                        }
                                                            
                                                            , "DESC_ID_PPTX":
                                                                            {"TIPO_DATO": "VARCHAR(100)"
                                                                            , "EXPORT_EXCEL": True
                                                                            , "IMPORT_SISTEMA": True
                                                                            , "TREEVIEW_ID_PPTX": True
                                                                            , "TREEVIEW_ID_XLS": False
                                                                            , "TREEVIEW_RANGOS_CELDAS": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": True
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                            , "CAMPO_JOIN": False
                                                                            }
                                                            
                                                            , "RUTA_PLANTILLA_PPTX":
                                                                                    {"TIPO_DATO": "VARCHAR(100)"
                                                                                    , "EXPORT_EXCEL": True
                                                                                    , "IMPORT_SISTEMA": True
                                                                                    , "TREEVIEW_ID_PPTX": False
                                                                                    , "TREEVIEW_ID_XLS": False
                                                                                    , "TREEVIEW_RANGOS_CELDAS": False
                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": True
                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                                    , "CAMPO_JOIN": False
                                                                                    }

                                                            , "TIEMPO_ESPERA_MAX_APERTURA_PPTX":
                                                                                                {"TIPO_DATO": "INT"
                                                                                                , "EXPORT_EXCEL": True
                                                                                                , "IMPORT_SISTEMA": True
                                                                                                , "TREEVIEW_ID_PPTX": False
                                                                                                , "TREEVIEW_ID_XLS": False
                                                                                                , "TREEVIEW_RANGOS_CELDAS": False
                                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": True
                                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                                                , "CAMPO_JOIN": False
                                                                                                }

                                                            , "NUMERO_TOTAL_SLIDES":
                                                                                    {"TIPO_DATO": "INT"
                                                                                    , "EXPORT_EXCEL": False
                                                                                    , "IMPORT_SISTEMA": True
                                                                                    , "TREEVIEW_ID_PPTX": False
                                                                                    , "TREEVIEW_ID_XLS": False
                                                                                    , "TREEVIEW_RANGOS_CELDAS": False
                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": True
                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                                    , "CAMPO_JOIN": False
                                                                                    }

                                                            , "ID_XLS":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": True
                                                                        , "IMPORT_SISTEMA": True
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": True
                                                                        , "TREEVIEW_RANGOS_CELDAS": True
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                        , "CAMPO_JOIN": False
                                                                        }

                                                            , "DESC_ID_XLS":
                                                                            {"TIPO_DATO": "VARCHAR(100)"
                                                                            , "EXPORT_EXCEL": True
                                                                            , "IMPORT_SISTEMA": True
                                                                            , "TREEVIEW_ID_PPTX": False
                                                                            , "TREEVIEW_ID_XLS": True
                                                                            , "TREEVIEW_RANGOS_CELDAS": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                            , "CAMPO_JOIN": False
                                                                            }

                                                            , "RUTA_XLS_ORIGEN":
                                                                                {"TIPO_DATO": "VARCHAR(100)"
                                                                                , "EXPORT_EXCEL": True
                                                                                , "IMPORT_SISTEMA": True
                                                                                , "TREEVIEW_ID_PPTX": False
                                                                                , "TREEVIEW_ID_XLS": False
                                                                                , "TREEVIEW_RANGOS_CELDAS": False
                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                                , "CAMPO_JOIN": False
                                                                                }

                                                            , "TIEMPO_ESPERA_MAX_APERTURA_EXCEL":
                                                                                                {"TIPO_DATO": "INT"
                                                                                                , "EXPORT_EXCEL": True
                                                                                                , "IMPORT_SISTEMA": True
                                                                                                , "TREEVIEW_ID_PPTX": False
                                                                                                , "TREEVIEW_ID_XLS": False
                                                                                                , "TREEVIEW_RANGOS_CELDAS": False
                                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                                                , "CAMPO_JOIN": False
                                                                                                }

                                                            , "ACTUALIZAR_VINCULOS_OTROS_EXCELS":
                                                                                                {"TIPO_DATO": "VARCHAR(100)"
                                                                                                , "EXPORT_EXCEL": True
                                                                                                , "IMPORT_SISTEMA": True
                                                                                                , "TREEVIEW_ID_PPTX": False
                                                                                                , "TREEVIEW_ID_XLS": False
                                                                                                , "TREEVIEW_RANGOS_CELDAS": False
                                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                                                , "CAMPO_JOIN": False
                                                                                                }

                                                            , "LISTA_HOJAS_XLS_FICHERO":
                                                                                    {"TIPO_DATO": "VARCHAR(1000)"
                                                                                    , "EXPORT_EXCEL": False
                                                                                    , "IMPORT_SISTEMA": True
                                                                                    , "TREEVIEW_ID_PPTX": False
                                                                                    , "TREEVIEW_ID_XLS": False
                                                                                    , "TREEVIEW_RANGOS_CELDAS": False
                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                                    , "CAMPO_JOIN": False
                                                                                    }

                                                            , "HOJA_XLS":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": True
                                                                        , "IMPORT_SISTEMA": True
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": True 
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                        , "CAMPO_JOIN": False
                                                                        }

                                                            , "RANGO_XLS":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": True
                                                                        , "IMPORT_SISTEMA": True
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": True
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                        , "CAMPO_JOIN": False
                                                                        }

                                                            , "SLIDE_PPTX":
                                                                        {"TIPO_DATO": "INT"
                                                                        , "EXPORT_EXCEL": True
                                                                        , "IMPORT_SISTEMA": True
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": True
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                        , "CAMPO_JOIN": False
                                                                        }

                                                            , "NOMBRE_SCREENSHOT":
                                                                                {"TIPO_DATO": "VARCHAR(100)"
                                                                                , "EXPORT_EXCEL": True
                                                                                , "IMPORT_SISTEMA": True
                                                                                , "TREEVIEW_ID_PPTX": False
                                                                                , "TREEVIEW_ID_XLS": False
                                                                                , "TREEVIEW_RANGOS_CELDAS": True
                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                                , "CAMPO_JOIN": False
                                                                                }

                                                            , "SCREENSHOT_PNG":
                                                                                {"TIPO_DATO": "BLOB"
                                                                                , "EXPORT_EXCEL": False
                                                                                , "IMPORT_SISTEMA": True
                                                                                , "TREEVIEW_ID_PPTX": False
                                                                                , "TREEVIEW_ID_XLS": False
                                                                                , "TREEVIEW_RANGOS_CELDAS": True
                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                                , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                                , "CAMPO_JOIN": False
                                                                                }

                                                            , "COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX":
                                                                                                    {"TIPO_DATO": "VARCHAR(100)"
                                                                                                    , "EXPORT_EXCEL": False
                                                                                                    , "IMPORT_SISTEMA": True
                                                                                                    , "TREEVIEW_ID_PPTX": False
                                                                                                    , "TREEVIEW_ID_XLS": False
                                                                                                    , "TREEVIEW_RANGOS_CELDAS": True
                                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                                                    , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": True
                                                                                                    , "CAMPO_JOIN": False
                                                                                                    }
                                                            }
                                            }

                            , "OTROS_DATOS":
                                            {"TABLA": "T_OTROS_DATOS"
                                            , "DICC_CAMPOS":
                                                            {"CONCEPTO":
                                                                            {"TIPO_DATO": "VARCHAR(100)"
                                                                            , "EXPORT_EXCEL": True
                                                                            , "IMPORT_SISTEMA": True
                                                                            , "TREEVIEW_ID_PPTX": False
                                                                            , "TREEVIEW_ID_XLS": False
                                                                            , "TREEVIEW_RANGOS_CELDAS": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                            , "CAMPO_JOIN": False
                                                                            }

                                                            , "DATO":
                                                                            {"TIPO_DATO": "VARCHAR(1000)"
                                                                            , "EXPORT_EXCEL": True
                                                                            , "IMPORT_SISTEMA": True
                                                                            , "TREEVIEW_ID_PPTX": False
                                                                            , "TREEVIEW_ID_XLS": False
                                                                            , "TREEVIEW_RANGOS_CELDAS": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                            , "CAMPO_JOIN": False
                                                                            }

                                                            , "ORDEN":
                                                                            {"TIPO_DATO": "INT"
                                                                            , "EXPORT_EXCEL": True
                                                                            , "IMPORT_SISTEMA": True
                                                                            , "TREEVIEW_ID_PPTX": False
                                                                            , "TREEVIEW_ID_XLS": False
                                                                            , "TREEVIEW_RANGOS_CELDAS": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                            , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                            , "CAMPO_JOIN": False
                                                                            }
                                                            }
                                            }


                            , "UPDATE_SCREENSHOT_BINARIO":
                                            {"TABLA": "T_UPDATE_SCREENSHOT_BINARIO"
                                            , "DICC_CAMPOS":
                                                            {"ID_PPTX":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": False
                                                                        , "IMPORT_SISTEMA": False
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": True
                                                                        }
                                                            , "ID_XLS":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": False
                                                                        , "IMPORT_SISTEMA": False
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": True
                                                                        }
                                                            , "HOJA_XLS":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": False
                                                                        , "IMPORT_SISTEMA": False
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": True
                                                                        }
                                                            , "RANGO_XLS":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": False
                                                                        , "IMPORT_SISTEMA": False
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": True
                                                                        }
                                                            , "SLIDE_PPTX":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": False
                                                                        , "IMPORT_SISTEMA": False
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": True
                                                                        }
                                                            , "NOMBRE_SCREENSHOT":
                                                                        {"TIPO_DATO": "VARCHAR(100)"
                                                                        , "EXPORT_EXCEL": False
                                                                        , "IMPORT_SISTEMA": False
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": True
                                                                        }
                                                            , "SCREENSHOT_PNG":
                                                                        {"TIPO_DATO": "BLOB"
                                                                        , "EXPORT_EXCEL": False
                                                                        , "IMPORT_SISTEMA": False
                                                                        , "TREEVIEW_ID_PPTX": False
                                                                        , "TREEVIEW_ID_XLS": False
                                                                        , "TREEVIEW_RANGOS_CELDAS": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1": False
                                                                        , "CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2": False
                                                                        , "CAMPO_JOIN": False
                                                                        }
                                                            }
                                            }
                            }


dicc_gui_combobox_procesos = {"COMBOBOX_SISTEMA":
                                            {"IMPORTAR_SISTEMA": 
                                                                {"OPCION": "Importar sistema"
                                                                , "LISTA_PARA_MSG": ["Se importara el sistema sqlite en la ram del pc\n\nTendras que indicar la ruta de la base datos sqlite.\n\n"
                                                                                    , "Si se localizan errores, se generara un log de errores (en formato .txt) que te detallara el porque de los fallos.\n\n"
                                                                                    , "Tendras que indicar en que directorio deseas guardar este posible log.\n\n"
                                                                                    , "Deseas continuar?"]
                                                                }
                                            
                                            , "CREAR_SISTEMA": {"OPCION": "Crear sistema"
                                                                , "LISTA_PARA_MSG": ["Se creara el sistema de base de datos sqlite en la ruta que indiques a continuación.\n\n"
                                                                                    , "Tendrás asimismo que indicar:\n\n--> ruta de Poppler\n\n--> una ruta local q1ue se usara para realizar las capturas de rangos de celdas.\n\n"
                                                                                    , "Si se localizan errores, se generara un log de errores (en formato .txt) que te detallara el porque de los fallos y "
                                                                                    , "se guardara en el mismo directorio donde quieres guardar el sistema sqlite.\n\n"
                                                                                    , "Deseas continuar?"]
                                                                }



                                            , "AGREGAR_RUTA_POPPLER": {"OPCION": "Agregar ruta Poppler"
                                                                , "LISTA_PARA_MSG": ["Agrega una ruta a Poppler.\n\nPoppler es el ejecutable que permite realizar los pantallazos de rangos de celdas "
                                                                                    , "excel con destino a un powerpoint y que permite realizar el proceso en segundo plano.\n\n"
                                                                                    , "Se pueden agregar varias rutas por lo que si usuarios de un mismo equipo prefieren tener Poppler instalado "
                                                                                    , "en su disco duro en vez de estar en una unidad compartida es por lo tanto viable.\n\n"
                                                                                    , "El aplicativo buscara la primera que sea valida y se quedara con esta el tiempo que el aplicativo este abierto.\n\n"
                                                                                    , "IMPORTANTE: se tiene que seleccionar la ruta de la carpeta 'bin' de Poppler:\n\n"
                                                                                    , r"ejemplo: 'C:\Users\user\Poppler\poppler-25.12.0\Library\bin'.\n\n"
                                                                                    , "Deseas continuar?"]
                                                                }

                                            , "AGREGAR_RUTA_LOCAL": {"OPCION": "Agregar ruta local"
                                                                , "LISTA_PARA_MSG": ["Agrega una ruta local que el aplicativo usara para realizar las conversiones de capturas de rangos excel en ficheros png "
                                                                                    , "para a posteriori colocar los pantallazos en el powerpoint y para guardar las imagenes binarias de estos png y guardarlos en el sistema sqlite.\n\n"
                                                                                    , "Se pueden agregar varias rutas por lo que si usuarios de un mismo equipo prefieren tener esta ruta "
                                                                                    , "en su disco duro en vez de estar en una unidad compartida es por lo tanto viable.\n\n"
                                                                                    , "El aplicativo buscara la primera que sea valida y se quedara con esta el tiempo que el aplicativo este abierto.\n\n"
                                                                                    , "Deseas continuar?"]
                                                                }

                                            , "DESCARGAR_SISTEMA_A_XLS": 
                                                                {"OPCION": "Descargar sistema a Excel"
                                                                , "LISTA_PARA_MSG": ["Se exportara la tabla parametrica sqlite a un fichero Excel donde podras configurar tus pptx.\n\n"
                                                                                    , "Tendras que indicar un directorio donde deseas guardar dicho fichero Excel.\n\n"
                                                                                    , "Si se localizan errores, se generara un log de errores (en formato .txt) que te detallara el porque de los fallos y "
                                                                                    , "se guardara en el mismo directorio donde deseas guardar el excel.\n\n"

                                                                                    , "Deseas continuar?"]
                                                                }

                                            }

                            , "COMBOBOX_PPTX":
                                                {"CONFIG_PASO_1":
                                                            {"OPCION": "Configurar pptx (paso 1)"
                                                            , "OPCION_LISTA_DESC": ["Se copiaran en la presentación pptx asociada al 'ID PPTX' seleccionado todos los screenshots correspondientes"
                                                                              , "a los rangos de celda excel configurados en el sistema sqlite:\n\n"
                                                                              , "--> COORD (sup-izq)\n"
                                                                              , "Se colocan todos los pantallazos de rangos de celdas en las slides de destino en el powerpoint asociado "
                                                                              , "en la esquina superior  izquierda. Ess la opción que hay que escojer cuando se trata de la 1era configuración de un powerpoint.\n\n"
                                                                              , "--> COORD (hibrido)\n"
                                                                              , "Se aplica en el caso de que se realizaron colocaciones anteriores (coordenadas y dimensiones) de pantallazos en el powerpoint asociado "
                                                                              , "a las que se le suman nuevos rangos de celdas configurados en el aplicativo pendientes de guardar sus coordenadas y dimensiones.\n\n"
                                                                              , "Al optar por esta opción, los pantallazos nuevos se colocan en la esquina superior izquierda "
                                                                              , "y el resto se colocan con las coordenadas y dimensiones ya guardadas.\n\n"
                                                                              , "*" * 73
                                                                              , "\n\n01\nSe abrira una ventana de dialogo para que el usuario seleccione donde quiere guardar el powerpoint "
                                                                              , "que ha de configurar manualmente (ver pasos 05 a 11 de este descriptivo). "
                                                                              , "Esta misma ruta sirve para guardar un log de errores (en caso de haberlos) en formato .txt que indique rango a rango el motivo del error."
                                                                              , "\n\n02\nSe creara una carpeta temporal en la misma ubicación que en el paso 01 de este descriptivo."
                                                                              , "\n\n03\nSe abriran uno a uno los distintos excel configurados (si es que hay varios) aplicando "
                                                                              , "el tiempo de espera maximo para abrir los ficheros (si la apertura de los excels excede este tiempo "
                                                                              , "de espera se informara de ello), pero no aplicando la opción 'ACTUALIZAR_VINCULOS_OTROS_EXCELS' este o no habilitada."
                                                                              , "\n\n04\nSe realizaran los screenshots de los rangos de celdas configurados para cada hoja del / los fichero(s) excel "
                                                                              , "y se guardaran en la carpeta temporal mencionada en el paso 02. "
                                                                              , "Los screenshots, asimismo, se guardaran en el sistema sqlite (mediante su imagen binaria) para poder usarlos como muestra"
                                                                              , "en la ventana de configuración de los ID PPTX en el apartado rangos de celdas."
                                                                              , "\n\n05\nSe abrira el fichero powerpoint configurado y se colocaran copias de los screenshots mencionados en el powerpoint "
                                                                              , "en las slides configuradas para cada rango de celda excel."
                                                                              , "\n\n06\nLas copias de los screenshots se colocaran en el powerpoint en la esquina superior izquierda de cada slide configurada."
                                                                              , "\n\n07\nSe eliminara la carpeta temporal mencionada en el punto 02 de este descriptivo."
                                                                              , "\n\n08\nSe guardara el powerpoint pero se dejara abierto."
                                                                              , "\n\n09\nUna vez el proceso finalizado y mediante aviso, el usuario tendra que colocar los screenshots "
                                                                              , "en el powerpoint donde quiera y con las dimensiones que desee."
                                                                              , "\n\n10\nUna vez realizado el ajuste de coordenadas y dimensiones, el usuario tendra que guardar el powerpoint y cerrarlo."
                                                                              , "\n\n11\nEl usuario tendra ahora que ejecutar el paso 'Configurar pptx (paso 2)'."
                                                                              ]
                                                            }

                                                , "CONFIG_PASO_2":
                                                            {"OPCION": "Configurar pptx (paso 2)"
                                                            , "OPCION_LISTA_DESC": ["Este proceso es automatico, el usuario no tiene que hacer nada pasado el punto 01 de este descriptivo."
                                                                                    , "\n\n01\nSe abrira una ventana de dialiogo para que el usuario seleccione el powerpoint que configuro "
                                                                                    , "manualmente en el proceso 'Configurar pptx (paso 1)'. "
                                                                                    , "En la misma ruta donde se ubica el powerpoint que se ha de seleccionar en el paso 01, se generara un log de errores (en caso de haberlos) "
                                                                                    , "en formato .txt que indique rango a rango el motivo del error.\n"
                                                                                    , "\n\n02\nEl app recuperara en cada slide del fichero powerpoint las coordenadas y dimensiones de cada screenshot y "
                                                                                    , "pasara a almacenarlas en el sistema sqlite."
                                                                                    , "\n\n03\nUna vez el proceso finalizado, se cerrara el powerpoint y se avisara al usuario de que el proceso esta finalizado."
                                                                                    , "\n\n" + "*" * 50
                                                                                    , "\nLlegado a este paso el 'ID PPTX' ya esta configurado y listo para reusarse cuando los distintos ficheros excels asociados "
                                                                                    , "tengan cambios en sus gráficas y/o cuadros etc etc y asi tener el powerpoint con los screenshots actualizados (ver paso 'Ejecutar pptx')."
                                                                                    ]
                                                            }

                                                , "EJECUCION":
                                                            {"OPCION": "Ejecutar pptx"
                                                            , "OPCION_LISTA_DESC": ["Este proceso ejecuta el powerpoint con los screenshots actualizados de los distintos excels asociados al 'ID PPTX' seleccionado."
                                                                                    , "\n\n01\nSe abrira una ventana de dialogo para que el usuario seleccione donde quiere guardar el powerpoint actualizado. "
                                                                                    , "Esta misma ruta sirve para guardar un log de errores (en caso de haberlos) en formato .txt que indique rango a rango el motivo del error."
                                                                                    ,"\n\n02\nSe creara una carpeta temporal en la misma ubicación donde se situa el app."
                                                                                    , "\n\n03\nSe abriran uno a uno los distintos excel configurados (si es que hay varios) aplicando "
                                                                                    , "el tiempo de espera maximo para abrir los ficheros (si la apertura de los excels excede este tiempo "
                                                                                    , "de espera se informara de ello), se aplicara la opción 'ACTUALIZAR_VINCULOS_OTROS_EXCELS' si está habilitada."
                                                                                    , "\n\n04\nSe realizaran los screenshots de los rangos de celda configurados para cada hoja del / los fichero(s) excel "
                                                                                    , "y se guardaran en la carpeta temporal mencionada en el paso 02."
                                                                                    , "\n\n05\nSe abrira el fichero powerpoint configurado y se colocaran copias de los screenshots mencionados en el powerpoint "
                                                                                    , "en las slides configuradas para cada rango de celda excel en las coordenadas y dimensiones registradas en el sistema sqlite."
                                                                                    , "\n\n06\nUna vez el proceso finalizado y mediante aviso (pop-up) desde la app, se borrara la carpeta del paso 02 de este descriptivo y "
                                                                                    , "se guardara el powerpoint que se dejara abierto."
                                                                                    , "\n\n07\nEl usuario, llegado a este punto, tendra el powerpoint actualizado y solo tendrá que adaptar sus comentarios."
                                                                                    ]
                                                            }
                                                }
                            
                            }


dicc_kwargs_config_gui_messagebox = {"showinfo":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_showinfo, "tupla_pixeles_imagen": (40, 40)}
                                                }

                                    , "showerror":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_showerror, "tupla_pixeles_imagen": (40, 40)}
                                                }

                                    , "showwarning":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_showwarning, "tupla_pixeles_imagen": (40, 40)}
                                                }

                                    , "askokcancel":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_askokcancel, "tupla_pixeles_imagen": (40, 40)}
                                                }

                                    , "askokcancel_warning":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_showwarning, "tupla_pixeles_imagen": (40, 40)}
                                                }

                                    , "askokcancel_con_opciones_colocacion_pantallazos":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_showwarning, "tupla_pixeles_imagen": (40, 40)}
                                                }

                                    , "askokcancel_con_opciones_guardado_config_id_pptx":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_showwarning, "tupla_pixeles_imagen": (40, 40)}
                                                }

                                    , "askokcancel_con_o_sin_muestra_screenshots":
                                                {"iconbitmap": template_ico_app_tapar_pluma_tkinter
                                                , "titulo": nombre_app
                                                , "font_messagebox": ("Calibri", 9)
                                                , "bg_messagebox":  "#FFFFFF"
                                                , "fg_messagebox": "black"
                                                , "bg_boton": "black"
                                                , "fg_boton": "white"
                                                , "font_boton": ("Calibri", 9, "bold")
                                                , "dicc_imagen": {"png_imagen_messagebox": template_img_messagebox_showwarning, "tupla_pixeles_imagen": (40, 40)}
                                                }
                                    }




def def_varios(opcion, **kwargs):
    #rutina que funciona a su vez como funcion y realiza calculos generales que se usan en otras rutinas

    global global_dicc_datos_id_pptx
    global global_ruta_local_config_sistema_sqlite
    global global_dicc_tablas_y_campos_sistema
    global global_lista_dicc_errores
    global global_lista_dicc_warning



    resultado_funcion = None

    #parametros kwargs
    kwargs_messagebox = kwargs.get("kwargs_messagebox", None)
    valor_check_tipo_dato = kwargs.get("valor_check_tipo_dato", None)
    id_xls_path = kwargs.get("id_xls_path", None)

    opcion_warning_errores = kwargs.get("opcion_warning_errores", None)
    opcion_proceso = kwargs.get("opcion_proceso", None)
    id_pptx_selecc = kwargs.get("id_pptx_selecc", None)

    directorio_log_errores_warning = kwargs.get("directorio_log_errores_warning", None)



    ####################################################################################################################################################
    # DICC_DATOS_ID_PPTX
    ####################################################################################################################################################
    if opcion == "DICC_DATOS_ID_PPTX":
        #aqui funciona como rutina y permnite crear la variable global global_dicc_datos_id_pptx
        #es un diccionario (es diccionario vacio si no hay configuraciones todavia guardadas en el sistema sqlite)
        #donde cada key corresponde con un id pptx y el valor de estas keys son diccionarios con todos los datos asociados al id pptx seleccionado 
        #prepara los datos en formatos que faciliten su incorporacion en la GUI en la ventana que se genera con la clase gui_config_id_pptx)
        #cada diccionario asociado a un id pptx contiene las keys siguientes:
        # --> id_pptx_desc                            descripcion del id pptx
        # --> ruta_plantilla_pptx                     ruta donde se ubica el fichero pptx en el cual se realizan als acciones pptx
        # --> tiempo_espera_max_apertura_pptx         tiempo de espera maximo que se espera para abrit el fichero pptx
        # --> numero_total_slides                     numero de slides que contiene el pptx
        # --> df_widget_treeview_id_xls               df con los id xls asociados al id pptx seleccionado (permite informar el treeview de la clase gui_config_id_pptx)
        #                                             (el df tiene las columnas siguientes ID_XLS y DESC_ID_XLS)
        #
        # --> lista_dicc_id_xls                       lista de diccionarios (cada diccionario corresponde a un id xls asociado al id pptx)
        #                                             cada diccionario contiene las keys:
        #                                                     --> id_xls                                     id xls asociado al id pptx
        #                                                     --> desc_id_xls                                descripcion del id xls
        #                                                     --> ruta_xls_origen                            ruta donde se ubica el fichero excel de origen correspondiente al id xls
        #                                                     --> tiempo_espera_max_apertura_xls             tiempo de espera maximo que se espera para abrit el fichero excel de origen
        #                                                     --> actualizar_vinculos_otros_xls              actualizar los vinculos hacia otros excels externos al abrir el fichero excel de origen
        #
        #                                                     --> lista_hojas_xls                            lista con las hojas disponibles en el excel asociado al id xls
        #                                                                                                    (se usa para actualizar el combobox de hojas xls en el frame de rangos de celdas)
        #
        #                                                     --> df_widget_treeview_rangos_celdas           df que sirve para informar el treeview por rangos de celda nada mas seleccionar un id xls en el treeview por id xls
        #                                                                                                    en la clase gui_config_id_pptx (sirve tambien para informar el treeview mencionado cuando se filtra el combobox 
        #                                                                                                    por la opcion TODOS para el treview por rangos de celda)
        #                                                                                                    (el df se crea con las columnas ID_XLS, HOJA_XLS, RANGO_XLS, SLIDE_PPTX, NOMBRE_SCREENSHOT,
        #                                                                                                     COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX y SCREENSHOT_PNG)


        #se extraen los datos de los id pptx en la lista sin duplicados lista_id_pptx
        lista_campos_df_treeview_id_xls = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_LISTA_CAMPOS"]["TREEVIEW_ID_XLS"]
        lista_campos_df_treeview_rangos_celdas = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_LISTA_CAMPOS"]["TREEVIEW_RANGOS_CELDAS"]
        lista_campos_calculo_1 = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_LISTA_CAMPOS"]["CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_1"]
        lista_campos_calculo_2 = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_LISTA_CAMPOS"]["CALCULO_GLOBAL_DICC_DATOS_ID_PPTX_2"]

        #se crea lista de datos en base al df global global_df_parametrica donde tan solo se extraen los datos de las columnas
        #ID_PPTX, DESC_ID_PPTX, RUTA_PLANTILLA_PPTX, TIEMPO_ESPERA_MAX_APERTURA_PPTX y NUMERO_TOTAL_SLIDES
        lista_id_pptx = global_df_parametrica[lista_campos_calculo_1].values.tolist() if len(global_df_parametrica) != 0 else []
        lista_id_pptx = [sublista for ind, sublista in enumerate(lista_id_pptx) if sublista not in lista_id_pptx[:ind]] if len(lista_id_pptx) != 0 else []


        #se calcula la variable global global_dicc_datos_id_pptx
        #siempre y cuando la bbdd del sistema sqlite importado no este vacia (en caso contrario devuelve un diccionario vacio)
        global_dicc_datos_id_pptx = {}
        if len(lista_id_pptx) != 0:

            #se crea el diccionario resultado de la funcion y mediante bucle sobre la lista lista_id_pptx se actualiza este diccionario
            for id_pptx, id_pptx_desc, ruta_plantilla_pptx, tiempo_espera_max_apertura_pptx, numero_total_slides in lista_id_pptx:


                #se extrae de la variable global global_df_parametrica un df filtrado por el id pptx de la iteracion y por idxls no vacios
                df_parametrica_id_pptx = (global_df_parametrica.loc[(global_df_parametrica["ID_PPTX"] == id_pptx) 
                                                                    & (global_df_parametrica["ID_XLS"].isnull() == False)
                                                                    , [col for col in global_df_parametrica.columns]])
                

                #se crea el df df_widget_treeview_id_xls asociado al id_pptx de la iteracion
                #puede estar vacio si no hay id xls configurados todavia
                df_widget_treeview_id_xls = pd.DataFrame(columns = lista_campos_df_treeview_id_xls)
                if len(df_parametrica_id_pptx) != 0:
                    df_widget_treeview_id_xls = df_parametrica_id_pptx[lista_campos_df_treeview_id_xls]
                    df_widget_treeview_id_xls.drop_duplicates(subset = lista_campos_df_treeview_id_xls, keep = "last", inplace = True)
                    df_widget_treeview_id_xls = df_widget_treeview_id_xls[lista_campos_df_treeview_id_xls].sort_values(["ID_XLS"], ascending = [True])
                    df_widget_treeview_id_xls.reset_index(drop = True, inplace = True)


                #se crea la lista de diccionarios lista_dicc_id_xls mediante bucle sobre el indice del df df_widget_treeview_id_xls
                #siempre y cuando df_widget_treeview_id_xls no este vacio sino se crea como lista vacia
                lista_dicc_id_xls = []
                if len(df_widget_treeview_id_xls) != 0:

                    for ind in df_widget_treeview_id_xls.index:
                        id_xls_iter = df_widget_treeview_id_xls.iloc[ind, df_widget_treeview_id_xls.columns.get_loc("ID_XLS")]


                        #se extrae un df de global_df_parametrica asociado al id pptx seleccionado y correspondiente al id xls de la iteracion del bucle y se re-ordena el df filtrado
                        #se exraeen las columnas DESC_ID_XLS, RUTA_XLS_ORIGEN, TIEMPO_ESPERA_MAX_APERTURA_EXCEL, ACTUALIZAR_VINCULOS_OTROS_EXCELS, LISTA_HOJAS_XLS_FICHERO, HOJA_XLS, RANGO_XLS, SLIDE_PPTX, 
                        #NOMBRE_SCREENSHOT, SCREENSHOT_PNG y COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX
                        df_id_xls_iter = (global_df_parametrica.loc[(global_df_parametrica["ID_PPTX"] == id_pptx) 
                                                                    & (global_df_parametrica["ID_XLS"] == id_xls_iter)
                                                                    , lista_campos_calculo_2])
                        
                        df_id_xls_iter.reset_index(drop = True, inplace = True)


                        #se calculan los valores desc_id_xls, tiempo_espera_max_apertura_xls y actualizar_vinculos_otros_xls
                        #(se extrae el 1er registro de df_id_xls_iter (son todos iguales para un mismo id_xls)
                        desc_id_xls = df_id_xls_iter.iloc[0, df_id_xls_iter.columns.get_loc("DESC_ID_XLS")] if len(df_id_xls_iter) != 0 else None
                        tiempo_espera_max_apertura_xls = df_id_xls_iter.iloc[0, df_id_xls_iter.columns.get_loc("TIEMPO_ESPERA_MAX_APERTURA_EXCEL")] if len(df_id_xls_iter) != 0 else None
                        actualizar_vinculos_otros_xls = df_id_xls_iter.iloc[0, df_id_xls_iter.columns.get_loc("ACTUALIZAR_VINCULOS_OTROS_EXCELS")] if len(df_id_xls_iter) != 0 else None


                        #se calcula el valor de la key lista_hojas_xls se extrae el 1er registro de df_id_xls_iter (son todos iguales para un mismo id_xls)
                        #y se convierte este valor de string a lista mediante la funcion def_varios (opcion = CHECK_TIPO_DATO_PYTHON)
                        #de la cual se extrae la 2nda variable que calcula de ahi el [1] en "def_varios("CHECK_TIPO_DATO_PYTHON", valor_check_tipo_dato = str_lista_hojas)[1]"
                        str_lista_hojas = df_id_xls_iter.iloc[0, df_id_xls_iter.columns.get_loc("LISTA_HOJAS_XLS_FICHERO")] if len(df_id_xls_iter) != 0 else None             

                        lista_hojas_xls_calc = def_varios("CHECK_TIPO_DATO_PYTHON", valor_check_tipo_dato = str_lista_hojas)[1] if str_lista_hojas is not None else []
                        lista_hojas_xls = lista_hojas_xls_calc if isinstance(lista_hojas_xls_calc, list) else []


                        #se calcula el valor de la key df_widget_treeview_rangos_celdas
                        #se exraeen las columnas HOJA_XLS, RANGO_XLS, SLIDE_PPTX, NOMBRE_SCREENSHOT, COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX y SCREENSHOT_PNG
                        df_widget_treeview_rangos_celdas = df_id_xls_iter.loc[df_id_xls_iter["HOJA_XLS"].isnull() == False, lista_campos_df_treeview_rangos_celdas]

                        df_widget_treeview_rangos_celdas["RANGO_XLS"] = df_widget_treeview_rangos_celdas["RANGO_XLS"].apply(lambda x: x.strip())
                        df_widget_treeview_rangos_celdas["SLIDE_PPTX"] = df_widget_treeview_rangos_celdas["SLIDE_PPTX"].apply(lambda x: int(float(x)) if x is not None else None)


                        if len(df_widget_treeview_rangos_celdas) != 0:
                            df_widget_treeview_rangos_celdas = df_widget_treeview_rangos_celdas[lista_campos_df_treeview_rangos_celdas].sort_values(["HOJA_XLS", "SLIDE_PPTX", "RANGO_XLS"], ascending = [True, True, True])
                            df_widget_treeview_rangos_celdas.reset_index(drop = True, inplace = True)


                        #se calcula el valor de la key ruta_xls_origen
                        #(se extrae el 1er registro de df_id_xls_iter (son todos iguales para un mismo id_xls)
                        ruta_xls_origen_ajust = os.path.normpath(df_id_xls_iter.iloc[0, df_id_xls_iter.columns.get_loc("RUTA_XLS_ORIGEN")]) if len(df_id_xls_iter) != 0 else None


                        #se crea diccionario temporal con las keys calculadas y se agrega a lista_dicc_id_xls
                        dicc_temp = {"id_xls": id_xls_iter
                                    , "desc_id_xls": desc_id_xls
                                    , "ruta_xls_origen": ruta_xls_origen_ajust
                                    , "tiempo_espera_max_apertura_xls": tiempo_espera_max_apertura_xls
                                    , "actualizar_vinculos_otros_xls": actualizar_vinculos_otros_xls
                                    , "lista_hojas_xls": lista_hojas_xls
                                    , "df_widget_treeview_rangos_celdas": df_widget_treeview_rangos_celdas
                                    }
                        
                        lista_dicc_id_xls.append(dicc_temp)
                    

                #se actualiza el diccionario global_dicc_datos_id_pptx
                ruta_plantilla_pptx_ajust = os.path.normpath(ruta_plantilla_pptx)

                dicc_id_pptx = {str(id_pptx):
                                            {"id_pptx_desc": id_pptx_desc
                                            , "ruta_plantilla_pptx": ruta_plantilla_pptx_ajust
                                            , "tiempo_espera_max_apertura_pptx": tiempo_espera_max_apertura_pptx
                                            , "numero_total_slides": int(numero_total_slides)#aqui es int
                                            , "df_widget_treeview_id_xls": df_widget_treeview_id_xls
                                            , "lista_dicc_id_xls": lista_dicc_id_xls
                                            }
                                    }
                
                global_dicc_datos_id_pptx.update(dicc_id_pptx)




    ####################################################################################################################################################
    # DICC_TABLAS_Y_CAMPOS_SISTEMA
    ####################################################################################################################################################
    elif opcion == "DICC_TABLAS_Y_CAMPOS_SISTEMA":
        #permite crear la variable global global_dicc_tablas_y_campos_sistema que es un diccionario basado en el diccionario dicc_tabla_config_sistema
        #pero cuya estructura se modifica para poder usar sus keys y valores de forma directa en distintas rutinas del app
        #tiene 3 keys (key_1 del diccionario dicc_tabla_config_sistema), cada key contiene un diccionario:
        # --> NOMBRE_TABLA              nombre de la tabla
        # --> DICC_SENTENCIAS_SQL       diccionario con 3 keys
        #                                  --> DROP   string con la sentencia sql para eliminar (si existe) la tabla asociada a la key_1
        #                                  --> CREATE string con la sentencia sql para crear la tabla asociada a la key_1
        #                                  --> DELETE string con la sentencia sql para eliminar los datos de la tabla asociada a la key_1
        #                                  --> INSERT string con la sentencia sql para insertar datos en la tabla asociada a la key_1

        # --> DICC_LISTA_CAMPOS        diccionario con las keys siguientes
        #                              (cada una almacena la lista de campos que tienen por valor True en la subkey_4 de dicc_tabla_config_sistema)
    

        global_dicc_tablas_y_campos_sistema = {}
        lista_tipo_tablas = list(dicc_tabla_config_sistema.keys())

        for tipo_tabla in lista_tipo_tablas:

            nombre_tabla = dicc_tabla_config_sistema[tipo_tabla]["TABLA"]
            dicc_campos = dicc_tabla_config_sistema[tipo_tabla]["DICC_CAMPOS"]

            lista_campos = list(dicc_campos.keys())
            lista_campos_con_tipo_dato = [f"{campo} " + dicc["TIPO_DATO"] for campo, dicc in dicc_campos.items()]


            #se calculan las sentencias drop, create, delete e insert
            sentencia_drop = f"DROP TABLE IF EXISTS {nombre_tabla}"
            sentencia_create = f"CREATE TABLE {nombre_tabla} (" + ", ".join(lista_campos_con_tipo_dato) + ")"
            sentencia_delete = f"DELETE FROM {nombre_tabla}"
            sentencia_insert = f"INSERT INTO {nombre_tabla} (" + ", ".join(lista_campos) + ") VALUES (" + ", ".join(["?" for _ in lista_campos]) + ")"


            #se calcula la sentencia_update (aplica solo para tipo_tabla = PARAMETRICA)
            # --> PARAMETRICA    es una sentencia para actualizar el campo SCREENSHOT_PNG desde la tabla UPDATE_SCREENSHOT_BINARIO
            #                    cuando se ejecuta el proceso de config pptx (paso 1) y asi dejar una muestra del screenshot en el sistema sqlite
            #                    para usarlo en la GUI de config de id pptx cuando se accede a la muestra de un rango de celdas
            sentencia_update = None
            if tipo_tabla == "PARAMETRICA":

                tabla_parametrica = dicc_tabla_config_sistema["PARAMETRICA"]["TABLA"]
                tabla_screenshots_blob = dicc_tabla_config_sistema["UPDATE_SCREENSHOT_BINARIO"]["TABLA"]

                concat_campos_update = ", ".join([f"{campo} = T1.{campo}" for campo, dicc in dicc_tabla_config_sistema["UPDATE_SCREENSHOT_BINARIO"]["DICC_CAMPOS"].items() if not dicc["CAMPO_JOIN"]])
                concat_campos_join = " AND ".join([f"{tabla_parametrica}.{campo} = T1.{campo}" for campo, dicc in dicc_tabla_config_sistema["UPDATE_SCREENSHOT_BINARIO"]["DICC_CAMPOS"].items() if dicc["CAMPO_JOIN"]])

                sentencia_update = f"UPDATE {tabla_parametrica} SET {concat_campos_update} FROM {tabla_screenshots_blob} T1 WHERE {concat_campos_join}"
                                                                    

            #se calcula dicc_lista_campos
            lista_opciones_campos = [tipo_opcion 
                                    for _, dicc in dicc_campos.items()
                                    for tipo_opcion in dicc.keys()
                                    if tipo_opcion != "TIPO_DATO"
                                    ]
            
            lista_opciones_campos = list(dict.fromkeys(lista_opciones_campos))

            dicc_lista_campos = {}
            for tipo_opcion in lista_opciones_campos:

                lista_campos = [campo for campo, dicc in dicc_campos.items() if dicc.get(tipo_opcion, False)]
                dicc_lista_campos.update({tipo_opcion: lista_campos})


            #se actualiza global_dicc_tablas_y_campos_sistema
            dicc_temp = {tipo_tabla:
                                    {"NOMBRE_TABLA": nombre_tabla
                                    , "DICC_SENTENCIAS_SQL": {"DROP": sentencia_drop
                                                            , "CREATE": sentencia_create
                                                            , "DELETE": sentencia_delete
                                                            , "INSERT": sentencia_insert
                                                            , "UPDATE": sentencia_update
                                                            }

                                    , "DICC_LISTA_CAMPOS": dicc_lista_campos 
                                    } 
                        }
            
            global_dicc_tablas_y_campos_sistema.update(dicc_temp)



    ####################################################################################################################################################
    # KWARGS_PARA_MESSAGEBOX_PROPIO
    ####################################################################################################################################################
    elif opcion == "KWARGS_PARA_MESSAGEBOX_PROPIO":
        #construye los kwargs para poder generar messagebox propios
        #mediante busqueda dentro del diccionario dicc_kwargs_config_gui_messagebox del tipo de messagebox
        #configurado en el kwargs (kwargs_messagebox) de la presente funcion (en la key tipo_messagebox)
        #el diccionario localizado en dicc_kwargs_config_gui_messagebox se le suma el mensaje configurado

        tipo_messagebox = kwargs_messagebox["tipo_messagebox"]
        mensaje = kwargs_messagebox["mensaje"]

        resultado_funcion = dicc_kwargs_config_gui_messagebox[tipo_messagebox]

        resultado_funcion.update({"tipo_messagebox": tipo_messagebox
                                  , "mensaje": mensaje})


    ####################################################################################################################################################
    # CHECK_TIPO_DATO_PYTHON
    ####################################################################################################################################################
    elif opcion == "CHECK_TIPO_DATO_PYTHON":
        #devuelve el tipo de dato (type) y el valor convertido al tipo de dato
        #se hace con literal_eval (libreria ast) en vez de usar eval pq literal_eval
        #chequea el tipo de dato sin ejecutar lo que aparece en el parametro que se le pasa a la funcion
        #y asi evita ejecutar codigo malicioso que podria encontrarse en el parametro kwargs valor_check_tipo_dato

        try:
            valor = ast.literal_eval(valor_check_tipo_dato)
            resultado_funcion = type(valor), valor

        
        except (ValueError, SyntaxError):
            resultado_funcion = str, valor_check_tipo_dato



    ####################################################################################################################################################
    # LISTA_HOJAS_XLS_DESDE_PATH_XLS
    ####################################################################################################################################################
    elif opcion == "LISTA_HOJAS_XLS_DESDE_PATH_XLS":
        #permite recuperar desde un path de un excel la lista de hojas que el excel continene 
        #se usa el metodo ExcelFile (pandas) pero segun el tipo de fichero excel
        #requiere agregarle un "engine" si el excdel es binario (xlsb)
        #se usa una funcion lambda lambda_engine_excelfile

        lambda_engine_excelfile = lambda extension: "pyxlsb" if extension.lower() == ".xlsb" else None

        extension_xls = pathlib.Path(id_xls_path).suffix.lower()
        engine_xls = lambda_engine_excelfile(extension_xls)
        xls_objeto = pd.ExcelFile(id_xls_path) if engine_xls is None else pd.ExcelFile(id_xls_path, engine = engine_xls)

        resultado_funcion = xls_objeto.sheet_names


      


    ####################################################################################################################################################
    # GENERAR_LOG_WARNING_ERRORES_PROCESOS_APP
    ####################################################################################################################################################
    elif opcion == "GENERAR_LOG_WARNING_ERRORES_PROCESOS_APP":
        #funciona como rutina y genera el log de warnings y/o errores a partir de las variables globales (se guarda en la ruta local configurada en el sistema sqlite)
        #se calcula con las variables globales global_lista_dicc_errores y global_lista_dicc_warning

        now = dt.datetime.now()



        ruta_fichero_log_warning_errores = os.path.join(directorio_log_errores_warning, nombre_log_warning_errores_procesos_app + "_" + str(re.sub("[^0-9a-zA-Z]+", "_", str(now))))
        ruta_fichero_log_warning_errores = os.path.normpath(ruta_fichero_log_warning_errores)


        lista_dicc_warning_errores_por_usar = (global_lista_dicc_warning
                                               if opcion_warning_errores == "WARNING"
                                               else
                                               global_lista_dicc_errores
                                               if opcion_warning_errores == "ERRORES"
                                               else
                                               global_lista_dicc_errores + global_lista_dicc_warning
                                               if opcion_warning_errores == "WARNING_ERRORES"
                                               else
                                               None
                                               )
        

        if len(lista_dicc_warning_errores_por_usar) != 0:

            string_log = ""
            for dicc in lista_dicc_warning_errores_por_usar:

                string_log_iter = "".join([f"{key}: {valor}\n"
                                            for key, valor in dicc.items()
                                            if valor is not None
                                            ])
                
                string_log_iter = string_log_iter + "*" * 100 + "\n\n"
                string_log = string_log + string_log_iter



        with open(ruta_fichero_log_warning_errores, 'w') as fich_log:
            fich_log.write(string_log)

        #se abre el fichero de logs (no se hace con os.startfile pq requiere que el que los ficheros txt tengan asociados siempre el bloc de notas para abrirlos)
        #usando subprocess se puede asignar por codigo usar el bloc de notas
        subprocess.run(["notepad.exe", ruta_fichero_log_warning_errores])




    ####################################################################################################################################################
    # MESSAGEBOX_PROCESOS_APP
    ####################################################################################################################################################
    elif opcion == "MESSAGEBOX_PROCESOS_APP":
        #funciona como funcion y devuelve tipo de messagebox y el mensaje asociado que se muestra en la GUI cuando el usuario ejecuta los procesos:
        # --> sistema
        # --> interaccion pptx
        # --> acceso a la ventana de configuracion de los id pptx
        # --> guardar configuraciones de los id pptx



        ####################################################################
        #IMPORTAR_SISTEMA
        #CREAR_SISTEMA
        #AGREGAR_RUTA_POPPLER
        #AGREGAR_RUTA_LOCAL
        #CREAR_DESCARGAR_SISTEMA_A_XLSSISTEMA
        ####################################################################
        if (opcion_proceso in [dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["IMPORTAR_SISTEMA"]["OPCION"]
                               , dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["CREAR_SISTEMA"]["OPCION"]
                               , dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["AGREGAR_RUTA_POPPLER"]["OPCION"]
                               , dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["AGREGAR_RUTA_LOCAL"]["OPCION"]
                               , dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["DESCARGAR_SISTEMA_A_XLS"]["OPCION"]
                               ]):

            mensaje_gui = ["".join(dicc["LISTA_PARA_MSG"]) for _, dicc in dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"].items() if opcion_proceso == dicc["OPCION"]][0]


            if (opcion_proceso in [dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["AGREGAR_RUTA_POPPLER"]["OPCION"]
                                    , dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["AGREGAR_RUTA_LOCAL"]["OPCION"]
                                    , dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"]["DESCARGAR_SISTEMA_A_XLS"]["OPCION"]
                                    ]
                and global_path_sistema_sqlite is None):

                resultado_funcion = "showerror", "No has importado el sistema sqlite todavia en memoria."

            else:
                resultado_funcion = "askokcancel", mensaje_gui



        ####################################################################
        #VENTANA_CONFIG_ID_PPTX
        ####################################################################       
        elif opcion_proceso == "VENTANA_CONFIG_ID_PPTX":

            if global_path_sistema_sqlite is None or not os.path.exists(global_path_sistema_sqlite):

                mensaje_gui_1 = "No se ha podido establecer conexión con el sistema sqlite para la descarga de los pantallazos de rangos de celdas para la muestras.\n\n"
                mensaje_gui_2 = "Comprueba que el sistema existe (es posible que alguien lo haya eliminado, movido a otra carpeta o cambiado el nombre)."
                mensaje_gui = mensaje_gui_1 + mensaje_gui_2

                resultado_funcion = "showerror", mensaje_gui
                
            else:

                mensaje_gui_1 = "Se abrira una nueva ventana donde podras crear y configurar id pptx definiendo rutas de ficheros excel, rangos de celda y slides pptx de destino.\n\n"
                mensaje_gui_2 = "Tienes la opción de habilitar o no la muestra de pantallazos asociados a cada rango de celdas (si ya previamente se ejecuto "
                mensaje_gui_3 = "la configuración de las coordenadas y dimensiones de los pantallazos en powerpoint).\n\n"
                mensaje_gui_4 = "Según el tamaño de tu sistema sqlite la apertura de la ventana de configuración puede demorarse.\n\n"
                mensaje_gui = mensaje_gui_1 + mensaje_gui_2 + mensaje_gui_3 + mensaje_gui_4

                resultado_funcion = "askokcancel_con_o_sin_muestra_screenshots", mensaje_gui





        ####################################################################
        #CONFIG_PASO_1
        #CONFIG_PASO_2
        #EJECUCION
        ####################################################################  

        elif (opcion_proceso in [dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_1"]["OPCION"]
                                    , dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_2"]["OPCION"]
                                    , dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["EJECUCION"]["OPCION"]
                                    ]):
            
            proceso_config_1 = dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_1"]["OPCION"]
            proceso_config_2 = dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_2"]["OPCION"]
            proceso_ejecucion =  dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["EJECUCION"]["OPCION"]


            if opcion_proceso is None and id_pptx_selecc is None:
                resultado_funcion = "showerror", "No has seleccionado ningún id pptx ni tampoco seleccionado ni ninguna acción a realizar sobre el mismo."


            elif opcion_proceso is not None and id_pptx_selecc is None:
                resultado_funcion = "showerror", "No has seleccionado ningún id pptx."


            elif opcion_proceso is None and id_pptx_selecc is not None:
                resultado_funcion = "showerror", f"No has seleccionado ninguna acción a realizar sobre '{id_pptx_selecc}'."


            else:
                proceso_ejecutable = True

                #se chequea para el proceso EJECUCION si hay coordenadas guardaddas en el sistema sqlite
                #en caso contrario no se puede ejecutar el proceso
                #(se chequea que el valor en el campo COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX de la tabla sqlite esta informado
                #y en caso de estarlo tras conversion de tipo de dato del string devuelve una tupla y que esta se compone de 4 elementos positivos
                #se usa la funcion def_varios, opcion = CHECK_TIPO_DATO_PYTHON)
                if opcion_proceso == proceso_ejecucion:

                    df_temp = global_df_parametrica.loc[global_df_parametrica["ID_PPTX"] == id_pptx_selecc, ["COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX"]]
                    lista_coordenadas_screenshots = df_temp.values.tolist()


                    lista_coordenadas_screenshots = [def_varios("CHECK_TIPO_DATO_PYTHON", valor_check_tipo_dato = item[0])
                                                        for item in lista_coordenadas_screenshots]

                    lista_coordenadas_screenshots = [tupla_coordenadas for _, tupla_coordenadas in lista_coordenadas_screenshots if tupla_coordenadas is not None]

                    proceso_ejecutable = (True
                                            if len(lista_coordenadas_screenshots) != 0
                                                and sum(1 if isinstance(tupla_coordenadas, tuple)
                                                                and len(tupla_coordenadas) == 4
                                                                and sum(1 if isinstance(item, (int, float)) and item >= 0 else 0
                                                                        for item in tupla_coordenadas) == 4
                                                            else 0 for tupla_coordenadas in lista_coordenadas_screenshots
                                                        ) == len(lista_coordenadas_screenshots)
                                            else False)
                        
                
                if not proceso_ejecutable:
                    resultado_funcion = "showerror", f"No se puede ejecutar el proceso '{opcion_proceso}' sobre '{id_pptx_selecc}' porque la configuración del paso 2 no esta realizada."

                else:
                    #se calcula el tiempo minimo de ejecucion del proceso
                    suma_tiempo_minimo_pptx = int(float(global_dicc_datos_id_pptx[id_pptx_selecc]["tiempo_espera_max_apertura_pptx"]))
                    suma_tiempo_minimo_xls = 0

                    lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx_selecc]["lista_dicc_id_xls"]


                    check_si_xls_config = 0
                    check_si_rangos_celdas_config = 0

                    if len(lista_dicc_id_xls) != 0:

                        check_si_xls_config = 1

                        for dicc in lista_dicc_id_xls:

                            tiempo_espera_max_apertura_xls = int(float(dicc["tiempo_espera_max_apertura_xls"]))
                            df_widget_treeview_rangos_celdas = dicc["df_widget_treeview_rangos_celdas"]

                            if len(df_widget_treeview_rangos_celdas) != 0:

                                check_si_rangos_celdas_config += 1
                                suma_tiempo_minimo_xls = suma_tiempo_minimo_xls + tiempo_espera_max_apertura_xls

                    suma_tiempo_minimo_pptx_y_xls = suma_tiempo_minimo_pptx + suma_tiempo_minimo_xls



                    #mensaje tiempo ejecucion
                    str_tiempo_pptx = (str(round(suma_tiempo_minimo_pptx / 60, 2)) + f" mn ({suma_tiempo_minimo_pptx} segundos)"
                                    if suma_tiempo_minimo_pptx / 60 >= 1 else f"{suma_tiempo_minimo_pptx} segundos")
                    

                    str_tiempo_pptx_y_xls = (str(round(suma_tiempo_minimo_pptx_y_xls / 60, 2)) + f" mn ({suma_tiempo_minimo_pptx_y_xls} segundos)"
                                            if suma_tiempo_minimo_pptx_y_xls / 60 >= 1 else f"{suma_tiempo_minimo_pptx_y_xls} segundos")


                    str_suma_tiempo_minimo = (str_tiempo_pptx_y_xls if opcion_proceso in [proceso_config_1, proceso_ejecucion]
                                                else
                                                str_tiempo_pptx if opcion_proceso == proceso_config_2
                                                else
                                                None
                                                )
                    
                    #mensaje por powerpoint sin excel o excel sin rangos de celdas
                    str_config_xls_y_rangos = (f"No se puede ejecutar el proceso para el id pptx ({id_pptx_selecc}) porque no tiene ningún excel asociado todavia."
                                            if check_si_xls_config == 0
                                            else
                                            f"No se puede ejecutar el proceso para el id pptx ({id_pptx_selecc}) porque ninguno de sus excel asociados tiene rangos de celdas configurados"
                                            if check_si_rangos_celdas_config == 0
                                            else
                                            None
                                            )
                    

                    if str_config_xls_y_rangos is not None:
                        resultado_funcion = "showerror", str_config_xls_y_rangos

                    else:
                        mensaje_gui_1 = f"El proceso tardara en ejecutarse un minimo de {str_suma_tiempo_minimo}.\n\n"
                        mensaje_gui_2 = "Es el tiempo de apertura maximo configurado para el id pptx seleccionado sumado a los de los id xls configurados (con rangos de celdas configurados).\n\n"


                        mensaje_gui_3 = ("Tendras que indicar a continuación en que directorio deseas guardar el powerpoint que se generara con los pantallazos de las capturas de rangos de celdas.\n\n"
                                        if opcion_proceso == proceso_config_1
                                        else
                                        "Tendras que indicar a continuación el powerpoint que usaste para recolocar y redimensionar los pantallazos de rangos de celdas.\n\n"
                                        if opcion_proceso == proceso_config_2
                                        else
                                        "Tendras que indicar a continuación en que directorio deseas guardar el powerpoint que se generara con los pantallazos de las capturas de rangos actualizados y correctamente colocados.\n\n"
                                        if opcion_proceso == proceso_ejecucion
                                        else
                                        ""
                                        )

                        mensaje_gui_4 = "Si se localizan errores en el proceso, se generara al final del proceso un log de errores (formato .txt) "
                        mensaje_gui_5 = "que te detallaran los errores y que se guardara en el mismo directorio comentado en el parrafo anterior.\n\nDeseas continuar?"

                        tipo_messagebox = "askokcancel" if opcion_proceso != proceso_config_1 else "askokcancel_con_opciones_colocacion_pantallazos"
                        resultado_funcion = tipo_messagebox, mensaje_gui_1 + mensaje_gui_2 + mensaje_gui_3 + mensaje_gui_4 + mensaje_gui_5



        ####################################################################
        #GUARDAR_CONFIGURACIONES_EN_SISTEMA_SQLITE
        ####################################################################       
        elif opcion_proceso == "GUARDAR_CONFIGURACIONES_EN_SISTEMA_SQLITE":

            #se reinician a listas vacias las variables globales global_lista_dicc_warning y global_lista_dicc_errores
            global_lista_dicc_warning = []
            global_lista_dicc_errores = []


            #se calcula el tiempo minimo de ejecucion del proceso
            suma_tiempo_minimo = 0
            for _, dicc_1 in global_dicc_datos_id_pptx.items():

                tiempo_espera_max_apertura_pptx = int(float(dicc_1["tiempo_espera_max_apertura_pptx"]))

                lista_dicc_id_xls = dicc_1["lista_dicc_id_xls"]


                if len(lista_dicc_id_xls) != 0:

                    for dicc_2 in lista_dicc_id_xls:

                        tiempo_espera_max_apertura_xls = int(float(dicc_2["tiempo_espera_max_apertura_xls"]))
                        df_widget_treeview_rangos_celdas = dicc_2["df_widget_treeview_rangos_celdas"]

                        if len(df_widget_treeview_rangos_celdas) != 0:
                            suma_tiempo_minimo = suma_tiempo_minimo + tiempo_espera_max_apertura_pptx + tiempo_espera_max_apertura_xls

            str_tiempo_minimo = str(round(suma_tiempo_minimo / 60, 2)) + f" mn ({suma_tiempo_minimo} segundos)"



            mensaje_gui_1 = "Se guardaran las configuraciones en el sistema sqlite.\n\n"
            mensaje_gui_2 = "Tienes la opción de actualizar tambien el número total de slides de cada powerpoint configurado y asimismo la lista de hojas disponibles actuales de cada excel de origen.\n\n"
            mensaje_gui_3 = "Cuando optas por esta actualización, permite al final del proceso localizar los rangos de celdas desfasados:\n\n"
            mensaje_gui_4 = "--> aquellos cuya slide powerpoint excede el número actual de slides del mismo.\n"
            mensaje_gui_5 = "--> aquellos cuyas hojas excel de origen ya no se encuentran dispobibles en el excel.\n\n"
            mensaje_gui_6 = f"Tiempo minimo de ejecución: {str_tiempo_minimo}.\n\n"
            mensaje_gui_7 = "Deseas continuar?"
            mensaje_gui_temp = mensaje_gui_1 + mensaje_gui_2 + mensaje_gui_3 + mensaje_gui_4 + mensaje_gui_5 + mensaje_gui_6 + mensaje_gui_7

            
            mensaje_gui = ("No hay ningún id pptx configurado por lo que no hay nada que guardar."
                           if len(global_dicc_datos_id_pptx) == 0 and len(global_df_parametrica) == 0
                           else
                           mensaje_gui_temp
                           )

            resultado_funcion = "askokcancel_con_opciones_guardado_config_id_pptx", mensaje_gui



    #resultado funcion
    return resultado_funcion



#################################################################################################################################
#################################################################################################################################
#################################################################################################################################
# RUTINAS / FUNCIONES - GUI VENTANA INICIO
#################################################################################################################################
#################################################################################################################################
#################################################################################################################################

def def_varios_gui_ventana_inicio(opcion, **kwargs):
    #rutina que funciona como funcion segun las opciones
    #ejecuta procesos varios de back end asociados con la GUI de la ventana de inicio

    global global_path_sistema_sqlite
    global global_df_parametrica
    global global_ruta_local_config_sistema_sqlite
    global global_df_treeview_id_pptx
    global global_proceso_en_ejecucion

    resultado_funcion = None

    #parametros kwargs
    opcion_combobox_sistema = kwargs.get("opcion_combobox_sistema", None)
    ruta_carpeta_guia_usuario = kwargs.get("ruta_carpeta_guia_usuario", None)
    id_pptx = kwargs.get("id_pptx", None)




    ####################################################################################################################################################
    # DESCARGA_GUIA_USUARIO
    ####################################################################################################################################################
    if opcion == "DESCARGA_GUIA_USUARIO":
        #funciona como rutina y permite descargar la guia de usuario

        now = dt.datetime.now()
        path_guia_usuario = os.path.join(ruta_carpeta_guia_usuario, nombre_guia_usuario + str(re.sub("[^0-9a-zA-Z]+", "_", str(now))) + ".pdf")
        shutil.copyfile(template_pdf_guia_usuario, path_guia_usuario)

        os.startfile(path_guia_usuario)





    ####################################################################################################################################################
    # DF_TREEVIEW_ID_PPTX
    ####################################################################################################################################################
    elif opcion == "DF_TREEVIEW_ID_PPTX":
        #devuelve un df que se usa para informar el treeview de id pptx

        lista_campos_treeview = [campo for campo, dicc_campo in dicc_tabla_config_sistema["PARAMETRICA"]["DICC_CAMPOS"].items() if dicc_campo["TREEVIEW_ID_PPTX"]]

        resultado_funcion = global_df_parametrica[lista_campos_treeview]
        resultado_funcion.drop_duplicates(subset = lista_campos_treeview, keep = "last", inplace = True)
        resultado_funcion = resultado_funcion[lista_campos_treeview].sort_values(["ID_PPTX"], ascending = [True])
        resultado_funcion.reset_index(drop = True, inplace = True)




    ####################################################################################################################################################
    # KEY_PROCESO_SISTEMA
    ####################################################################################################################################################
    elif opcion == "KEY_PROCESO_SISTEMA":
        #devuelve la key del proceso de sistema (esta key es fija internamente en el codigo, su valor es configurable
        #para que salga el texto que quiera el usuario en el combobox de sistema en la GUI, es para evitar por si
        #se cambian los literales tener que retocar el codigo del app)

        resultado_funcion = next((key for key, value in dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"].items() if value["OPCION"] == opcion_combobox_sistema), None)



    elif opcion == "LISTA_OPCIONES_COMBOBOX_GUI_SISTEMA":
        #devuelve la lista de opciones que sirve para el combobox de sistema en la GUI

        resultado_funcion = [dicc["OPCION"] for _, dicc in dicc_gui_combobox_procesos["COMBOBOX_SISTEMA"].items()]




    elif opcion == "LISTA_OPCIONES_COMBOBOX_GUI_ACCIONES_PPTX":
        #devuelve la lista de opciones que sirve para el combobox de acciones pptx en la GUI

        resultado_funcion = [dicc["OPCION"] for _, dicc in dicc_gui_combobox_procesos["COMBOBOX_PPTX"].items()]



    elif opcion == "DESC_ASOCIADO_ACCION_PPTX":
        #devuelve el texto asociado a la accion pptx seleccionada en el combobox y que sirve para informar el scrolledtext de descripcion de la accion

        resultado_funcion = "".join([dicc["OPCION_LISTA_DESC"] for _, dicc in dicc_gui_combobox_procesos["COMBOBOX_PPTX"].items() if opcion_combobox_sistema == dicc["OPCION"]][0])


    
    #resultado de la funcion
    return resultado_funcion



def def_config_sistema_sqlite(opcion, **kwargs):
    #rutina que funciona como funcion segun la opcion seleccionada
    #que permite interactuar con la bbdd sqlite donde se parametrizan todas las configuraciones de pptx

    global global_dicc_tablas_y_campos_sistema
    global global_path_sistema_sqlite
    global global_df_parametrica
    global global_poppler_path
    global global_ruta_local_config_sistema_sqlite
    global global_ruta_local_screenshots_png_muestra
    global global_df_treeview_id_pptx
    global global_proceso_en_ejecucion
    global global_lista_dicc_errores
    global global_lista_dicc_warning


    resultado_funcion = None

    warnings.filterwarnings("ignore")


    #parametros kwargs
    id_pptx_selecc = kwargs.get("id_pptx_selecc", None)
    lista_screenshots_con_coordenadas = kwargs.get("lista_screenshots_con_coordenadas", None)
    lista_datos_screenshots_binarios = kwargs.get("lista_datos_screenshots_binarios", None)
    opcion_interaccion_pptx = kwargs.get("opcion_interaccion_pptx", None)
    directorio_sqlite = kwargs.get("directorio_sqlite", None)
    path_sqlite = kwargs.get("path_sqlite", None)
    directorio_excel = kwargs.get("directorio_excel", None)
    directorio_poppler = kwargs.get("directorio_poppler", None)
    directorio_local_screenshots = kwargs.get("directorio_local_screenshots", None)
    ruta_add = kwargs.get("ruta_add", None)


    #AL INICIAR EL PROCESO: se inicializa la variable global global_lista_dicc_errores como lista vacia solo para los procesos independientes
    #UPDATE_SCREENSHOTS_COORDENADAS y UPDATE_SCREENSHOTS_IMAGENES_BINARIAS se usan en la rutina def_interaccion_pptx
    #que tambien necesita el uso de global_lista_dicc_errores para almacenar los posiblers errores de ahi que para estas opciones
    #no se ha de reiniciar la variable global
    #afecta tambien la variable global global_proceso_en_ejecucion que se usa para impedir ejecutar en la GUI otro proceso hasta que acabe el que esta en curso
    if (def_varios_gui_ventana_inicio("KEY_PROCESO_SISTEMA", opcion_combobox_sistema = opcion) in ["IMPORTAR_SISTEMA", "CREAR_SISTEMA", "AGREGAR_RUTA_POPPLER", "AGREGAR_RUTA_LOCAL", "DESCARGAR_SISTEMA_A_XLS"] 
        or opcion == "DESCARGA_SCREENSHOTS_PNG_PARA_MUESTRA_GUI_CONFIG_ID_PPTX"):


        global_lista_dicc_errores = []

        #se fija la variable global para la rutina de los threads en la GUI y asi bloquear la ejecucion de cualquier otro proceso 
        #hasta que el actual en curso no haya finalizado
        global_proceso_en_ejecucion = "SI"



    #####################################################################################################################################
    # CREAR_SISTEMA --> crea la bbdd sqlite con las tablas necesarias (configuradas en el diccionario dicc_tabla_config_sistema)
    #####################################################################################################################################
    if def_varios_gui_ventana_inicio("KEY_PROCESO_SISTEMA", opcion_combobox_sistema = opcion) == "CREAR_SISTEMA":

        #se comprueba que la ruta que se quiere agregar de Poppler tiene la ultima carpeta llamada 'bin' y que dentro de ella se encuentran los exe_
        # --> pdftoppm.exe
        # --> pdfinfo.exe
        check_directorio_poppler = False
        if os.path.exists(directorio_poppler):
            
            directorio_poppler_ajust = os.path.normpath(directorio_poppler)
            ultima_carpeta = os.path.basename(directorio_poppler_ajust)

            if ultima_carpeta.lower() == "bin":
                ruta_pdfinfo = os.path.join(directorio_poppler, "pdfinfo.exe")
                ruta_pdftoppm = os.path.join(directorio_poppler, "pdftoppm.exe")

                ruta_pdfinfo_ajust = os.path.normpath(ruta_pdfinfo)
                ruta_pdftoppm_ajust = os.path.normpath(ruta_pdftoppm)

                if os.path.exists(ruta_pdfinfo_ajust) and os.path.exists(ruta_pdftoppm_ajust):
                    check_directorio_poppler = True


            if not check_directorio_poppler:
                mensaje_gui_1 = f"La ruta que intentas agregar no es una ruta valida de Poppler:\n\n{directorio_poppler}\n\nLa ultima carpeta debe llamarse 'bin' y dentro de la misma debes tener los exe:\n"
                mensaje_gui_2 = "--> pdfinfo.exe\n--> pdftoppm.exe"
                resultado_funcion = mensaje_gui_1 + mensaje_gui_2

            else:

                now = dt.datetime.now()
                ruta_bbdd_config = os.path.join(directorio_sqlite, nombre_sistema_sqlite + "_" + str(re.sub("[^0-9a-zA-Z]+", "_", str(now))) + ".db")


                #se borra la bbdd sqlite si ya se creo anteriormente
                if os.path.exists(ruta_bbdd_config):
                    os.remove(ruta_bbdd_config)



                #se crea la bbdd sqlite con las tablas correspondientes configuradas en el diccionario dicc_tabla_config_sistema
                #(se usa la variable global global_dicc_tablas_y_campos_sistema donde se han almacenado las sentencias sql)
                #no es necesario poner un bloqwue try except pq la bbdd se crea siempre
                #(solo para la conexion a la bbdd, para las operaciones SCUD si es necesario)
                MiConexion = sqlite3.connect(ruta_bbdd_config)
                cursor = MiConexion.cursor()

                try:

                    lista_sentencias_sql = [sentencia_sql
                                            for _, dicc in global_dicc_tablas_y_campos_sistema.items()
                                            for tipo_sentencia_sql, sentencia_sql in dicc["DICC_SENTENCIAS_SQL"].items()
                                            if tipo_sentencia_sql in ["DROP", "CREATE"]]

                    for sentencia_sql in lista_sentencias_sql:
                        cursor.execute(sentencia_sql)
                        MiConexion.commit()

                    #se actualiza el directorio poppler (se recupera la sentencia insert de la variable global global_dicc_tablas_y_campos_sistema)
                    sentencia_sql_insert = global_dicc_tablas_y_campos_sistema["OTROS_DATOS"]["DICC_SENTENCIAS_SQL"]["INSERT"]
                    lista_ruta = [["RUTA_POPPLER", directorio_poppler_ajust, 1]]
                    cursor.executemany(sentencia_sql_insert, lista_ruta)


                    #se actualiza el path local (se recupera la sentencia insert de la variable global global_dicc_tablas_y_campos_sistema)
                    directorio_local_screenshots_ajust = os.path.normpath(directorio_local_screenshots)

                    sentencia_sql_insert = global_dicc_tablas_y_campos_sistema["OTROS_DATOS"]["DICC_SENTENCIAS_SQL"]["INSERT"]
                    lista_ruta = [["RUTA_LOCAL", directorio_local_screenshots_ajust, 1]]
                    cursor.executemany(sentencia_sql_insert, lista_ruta)

                    MiConexion.commit()


                except Exception as err_1:
                    #se registra el posible error de guardado de las coordenadas del screenshot de la iteracion

                    traceback_error = traceback.extract_tb(err_1.__traceback__)
                    modulo_python = os.path.basename(traceback_error[0].filename)
                    rutina_python = traceback_error[0].name
                    linea_error = traceback_error[0].lineno

                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                    mensaje_error_1 = f"No se han podido crear en el sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                    mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                    mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                    dicc_error["PROCESO"] = opcion
                    dicc_error["FASE_PROCESO"] = "Creación de las tablas."
                    dicc_error["ID_PPTX"] = None
                    dicc_error["ID_XLS"] = None
                    dicc_error["RUTA_FICHERO"] = None
                    dicc_error["HOJA_XLS"] = None
                    dicc_error["RANGO_CELDAS"] = None
                    dicc_error["SLIDE_PPTX"] = None
                    dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                    dicc_error["MODULO_PYTHON"] = modulo_python
                    dicc_error["RUTINA_PYTHON"] = rutina_python
                    dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                    dicc_error["TRACEBACK"] = str(err_1)

                    global_lista_dicc_errores.append(dicc_error)

                    pass

                finally:
                    cursor.close()
                    MiConexion.close()




    #####################################################################################################################################
    # IMPORTAR_SISTEMA --> genera un df con la parametrica (se almacena en una variable global para usarla cuando se ejecuta 
    #                      esta misma rutina con la opcion DESCARGAR_SISTEMA_BBDD_A_XLS o cuando se ejecutan las rutinas 
    #                      de config y ejecucion de los pptx)
    #
    # --> se almacena el path del sistema sqlite en la variable global global_path_sistema_sqlite pq se usa en la misma rutina
    #     pero con opciones distintas en otros procesos del app
    #
    # --> se almacena en la variable global global_ruta_local_config_sistema_sqlite la ruta local donde se guardan temporalmente
    #     los screenshots de las acciones pptx y donde se guardan los de errores / warnings
    #
    #  --> se crea tambien el df para informar el treeview id pptx y se informa en la variable global global_df_treeview_id_pptx
    #####################################################################################################################################
    elif def_varios_gui_ventana_inicio("KEY_PROCESO_SISTEMA", opcion_combobox_sistema = opcion) == "IMPORTAR_SISTEMA":

        
        global_path_sistema_sqlite = os.path.normpath(path_sqlite)

        if not os.path.exists(global_path_sistema_sqlite):
            #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
            #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
            #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n"
            mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n"
            mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

            dicc_error["PROCESO"] = opcion
            dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite."
            dicc_error["ID_PPTX"] = None
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:

            MiConexion = sqlite3.connect(global_path_sistema_sqlite)

            try:
                sentencia_sql_1 = "SELECT * FROM " + dicc_tabla_config_sistema["PARAMETRICA"]["TABLA"] + " ORDER BY ID_PPTX"
                sentencia_sql_2 = "SELECT DATO, ORDEN FROM " + dicc_tabla_config_sistema["OTROS_DATOS"]["TABLA"] + " WHERE CONCEPTO = 'RUTA_LOCAL'"
                sentencia_sql_3 = "SELECT DATO, ORDEN FROM " + dicc_tabla_config_sistema["OTROS_DATOS"]["TABLA"] + " WHERE CONCEPTO = 'RUTA_POPPLER'"


                global_df_parametrica = pd.read_sql(sentencia_sql_1, MiConexion)
                global_df_parametrica = global_df_parametrica.replace({np.nan: None})
                global_df_parametrica.reset_index(drop = True, inplace = True)

                df_ruta_local = pd.read_sql(sentencia_sql_2, MiConexion)
                df_ruta_local = df_ruta_local.replace({np.nan: None})
                df_ruta_local.reset_index(drop = True, inplace = True)

                df_poppler_path = pd.read_sql(sentencia_sql_3, MiConexion)
                df_poppler_path = df_poppler_path.replace({np.nan: None})
                df_poppler_path.reset_index(drop = True, inplace = True)


                #se busca la 1era ruta valida para el usuario y se almacena
                #en la variable global global_ruta_local_config_sistema_sqlite
                if len(df_ruta_local) == 0:
                    global_ruta_local_config_sistema_sqlite = None

                else:

                    check_ruta = False
                    for ind in df_ruta_local.index:
                        ruta_config = df_ruta_local.iloc[ind, df_ruta_local.columns.get_loc("DATO")]

                        if os.path.exists(ruta_config):
                            
                            check_ruta = True
                            global_ruta_local_config_sistema_sqlite = os.path.normpath(ruta_config)
                            break


                    if not check_ruta:
                        #se registra el error de que no se ha podido localizar una ruta local valida
                        dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                        mensaje_error_1 = f"No se ha podido localizar ninguna ruta local configurada en el sistema sqlite.\n\n"
                        mensaje_error_2 = f"Descarga a Excel mediante la opción 'Descargar sistema a Excel' y chequea que las rutas locales "
                        mensaje_error_3 = "existen y que tienes acceso a almenos una de ellas.."

                        dicc_error["PROCESO"] = opcion
                        dicc_error["FASE_PROCESO"] = "Localización de la ruta local para realizar los cálculos de capturas de rangos de celdas."
                        dicc_error["ID_PPTX"] = None
                        dicc_error["ID_XLS"] = None
                        dicc_error["RUTA_FICHERO"] = None
                        dicc_error["HOJA_XLS"] = None
                        dicc_error["RANGO_CELDAS"] = None
                        dicc_error["SLIDE_PPTX"] = None
                        dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                        dicc_error["MODULO_PYTHON"] = None
                        dicc_error["RUTINA_PYTHON"] = None
                        dicc_error["LINEA_CODIGO_PYTHON"] = None
                        dicc_error["TRACEBACK"] = None

                        global_lista_dicc_errores.append(dicc_error)


                #e busca la 1era ruta valida para el usuario y se almacena
                #en la variable global global_poppler_path
                #
                #para que una ruta sea valida no basta con que el usuario pueda acceder:
                # --> su ultima carpeta tiene que llamarse 'bin'
                # --> dentro de esta carpeta tiene que haber los exe siguientes:
                #                                                                --> pdftoppm.exe
                #                                                                --> pdfinfo.exe

                if len(df_poppler_path) == 0:
                    global_poppler_path = None

                else:

                    check_ruta = False
                    for ind in df_poppler_path.index:
                        ruta_config = df_poppler_path.iloc[ind, df_poppler_path.columns.get_loc("DATO")]
                        
                        if os.path.exists(ruta_config):
                            
                            ruta_config_ajust = os.path.normpath(ruta_config)
                            ultima_carpeta = os.path.basename(ruta_config_ajust)

                            if ultima_carpeta.lower() == "bin":
                                ruta_pdfinfo = os.path.join(ruta_config_ajust, "pdfinfo.exe")
                                ruta_pdftoppm = os.path.join(ruta_config_ajust, "pdftoppm.exe")

                                ruta_pdfinfo_ajust = os.path.normpath(ruta_pdfinfo)
                                ruta_pdftoppm_ajust = os.path.normpath(ruta_pdftoppm)

                                if os.path.exists(ruta_pdfinfo_ajust) and os.path.exists(ruta_pdftoppm_ajust):
                                    check_ruta = True
                                    global_poppler_path = ruta_config_ajust

                                    break


                    if not check_ruta:
                        #se registra el error de que no se ha podido localizar una ruta local valida
                        dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                        mensaje_error_1 = f"No se ha podido localizar ninguna ruta valida a Poppler.\n\n"
                        mensaje_error_2 = f"Descarga a Excel mediante la opción 'Descargar sistema a Excel' y chequea que las rutas a Poppler "
                        mensaje_error_3 = "existen y que tienes acceso a almenos una de ellas.\n\n"
                        mensaje_error_4 = "La ultima carpeta de dichas rutas ha de llamarse 'bin' y dentre de esta carpeta debe haber estos 2 exe:\n"
                        mensaje_error_5 = "--> pdfinfo.exe\n--> pdftoppm.exe"

                        dicc_error["PROCESO"] = opcion
                        dicc_error["FASE_PROCESO"] = "Localización de la ruta Poppler para realizar los cálculos de capturas de rangos de celdas."
                        dicc_error["ID_PPTX"] = None
                        dicc_error["ID_XLS"] = None
                        dicc_error["RUTA_FICHERO"] = None
                        dicc_error["HOJA_XLS"] = None
                        dicc_error["RANGO_CELDAS"] = None
                        dicc_error["SLIDE_PPTX"] = None
                        dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3 + mensaje_error_4 + mensaje_error_5
                        dicc_error["MODULO_PYTHON"] = None
                        dicc_error["RUTINA_PYTHON"] = None
                        dicc_error["LINEA_CODIGO_PYTHON"] = None
                        dicc_error["TRACEBACK"] = None

                        global_lista_dicc_errores.append(dicc_error)


                #se calula la variable global global_df_treeview_id_pptx
                #permite incorporar datos en el treeview de la ventana de inicio
                global_df_treeview_id_pptx = def_varios_gui_ventana_inicio("DF_TREEVIEW_ID_PPTX")


                #se calula la variable global global_dicc_datos_id_pptx para qu el usuario pueda nada mas importado el sitema sqlite ejecutar acciones pptx
                def_varios("DICC_DATOS_ID_PPTX")


            except Exception as err_1:
                #se registra el posible error de carga en memoria de los datos del sistema sqlite

                traceback_error = traceback.extract_tb(err_1.__traceback__)
                modulo_python = os.path.basename(traceback_error[0].filename)
                rutina_python = traceback_error[0].name
                linea_error = traceback_error[0].lineno

                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_error_1 = f"No se ha podido guardar en memoria RAM los datos del sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                dicc_error["PROCESO"] = opcion
                dicc_error["FASE_PROCESO"] = "Descarga en memoria RAM del sistema sqlite."
                dicc_error["ID_PPTX"] = None
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                dicc_error["MODULO_PYTHON"] = modulo_python
                dicc_error["RUTINA_PYTHON"] = rutina_python
                dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                dicc_error["TRACEBACK"] = str(err_1)

                global_lista_dicc_errores.append(dicc_error)

                pass


            finally:
                MiConexion.close()



    #####################################################################################################################################
    # AGREGAR_RUTA_POPPLER --> agrega una ruta a poppler
    #                          se almacena en el campo DATO de la tabla T_OTROS_DATOS (donde el campo CONCEPTO = RUTA_POPPLER)
    #                          se agrega cada nueva ruta a una lista que se convierte a string para almacenarla en el sistema sqlite
    #####################################################################################################################################
    elif def_varios_gui_ventana_inicio("KEY_PROCESO_SISTEMA", opcion_combobox_sistema = opcion) == "AGREGAR_RUTA_POPPLER":

        if not os.path.exists(global_path_sistema_sqlite):
            #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
            #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
            #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n"
            mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n"
            mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

            dicc_error["PROCESO"] = opcion
            dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite para agregar una nueva ruta a Poppler."
            dicc_error["ID_PPTX"] = None
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:
            #se comprueba que la ruta que se quiere agregar tiene la ultima carpeta llamada 'bin' y que dentro de ella se encuentran los exe_
            # --> pdftoppm.exe
            # --> pdfinfo.exe
            check_ruta_add = False
            if os.path.exists(ruta_add):
                
                ruta_add_ajust = os.path.normpath(ruta_add)
                ultima_carpeta = os.path.basename(ruta_add_ajust)

                if ultima_carpeta.lower() == "bin":
                    ruta_pdfinfo = os.path.join(ruta_add_ajust, "pdfinfo.exe")
                    ruta_pdftoppm = os.path.join(ruta_add_ajust, "pdftoppm.exe")

                    ruta_pdfinfo_ajust = os.path.normpath(ruta_pdfinfo)
                    ruta_pdftoppm_ajust = os.path.normpath(ruta_pdftoppm)

                    if os.path.exists(ruta_pdfinfo_ajust) and os.path.exists(ruta_pdftoppm_ajust):
                        check_ruta_add = True


            if not check_ruta_add:
                mensaje_gui_1 = f"La ruta que intentas agregar no es una ruta valida de Poppler:\n\n{ruta_add}\n\nLa ultima carpeta debe llamarse 'bin' y dentro de la misma debes tener los exe:\n"
                mensaje_gui_2 = "--> pdfinfo.exe\n--> pdftoppm.exe"
                resultado_funcion = mensaje_gui_1 + mensaje_gui_2

            else:

                try:
                    ruta_add_ajust = os.path.normpath(ruta_add)

                    MiConexion = sqlite3.connect(global_path_sistema_sqlite)
                    cursor = MiConexion.cursor()

                    sentencia_sql_1 = "SELECT DATO FROM " + dicc_tabla_config_sistema["OTROS_DATOS"]["TABLA"] + f" WHERE CONCEPTO = 'RUTA_POPPLER' AND DATO = '{ruta_add_ajust}'"
                    df_check_ruta_ya_guardada = pd.read_sql(sentencia_sql_1, MiConexion)


                    if len(df_check_ruta_ya_guardada) == 0:

                        sentencia_sql_2 = "SELECT DATO FROM " + dicc_tabla_config_sistema["OTROS_DATOS"]["TABLA"] + " WHERE CONCEPTO = 'RUTA_POPPLER'"
                        
                        df_ruta = pd.read_sql(sentencia_sql_2, MiConexion)

                        lista_para_sqlite =  [["RUTA_POPPLER", ruta_add_ajust, 1]] if len(df_ruta) == 0 else [["RUTA_POPPLER", ruta_add_ajust, len(df_ruta) + 1]]
                        sentencia_sql_insert = global_dicc_tablas_y_campos_sistema["OTROS_DATOS"]["DICC_SENTENCIAS_SQL"]["INSERT"]

                        cursor.executemany(sentencia_sql_insert, lista_para_sqlite)

                        MiConexion.commit()


                except Exception as err_1:
                    #se registra el posible error de guardado de las coordenadas del screenshot de la iteracion

                    traceback_error = traceback.extract_tb(err_1.__traceback__)
                    modulo_python = os.path.basename(traceback_error[0].filename)
                    rutina_python = traceback_error[0].name
                    linea_error = traceback_error[0].lineno

                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                    mensaje_error_1 = f"No se han podido agregar la nueva ruta a Poppler en el sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                    mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                    mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                    dicc_error["PROCESO"] = opcion
                    dicc_error["FASE_PROCESO"] = "Agregación nueva ruta a Poppler."
                    dicc_error["ID_PPTX"] = None
                    dicc_error["ID_XLS"] = None
                    dicc_error["RUTA_FICHERO"] = None
                    dicc_error["HOJA_XLS"] = None
                    dicc_error["RANGO_CELDAS"] = None
                    dicc_error["SLIDE_PPTX"] = None
                    dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                    dicc_error["MODULO_PYTHON"] = modulo_python
                    dicc_error["RUTINA_PYTHON"] = rutina_python
                    dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                    dicc_error["TRACEBACK"] = str(err_1)

                    global_lista_dicc_errores.append(dicc_error)

                    pass

                finally:
                    cursor.close()
                    MiConexion.close()



    #####################################################################################################################################
    # AGREGAR_RUTA_LOCAL --> agrega una ruta a poppler
    #                        se almacena en el campo DATO de la tabla T_OTROS_DATOS (donde el campo CONCEPTO = RUTA_POPPLER)
    #                        se agrega cada nueva ruta a una lista que se convierte a string para almacenarla en el sistema sqlite
    #####################################################################################################################################
    elif def_varios_gui_ventana_inicio("KEY_PROCESO_SISTEMA", opcion_combobox_sistema = opcion) == "AGREGAR_RUTA_LOCAL":


        if not os.path.exists(global_path_sistema_sqlite):
            #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
            #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
            #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n"
            mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n"
            mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

            dicc_error["PROCESO"] = opcion
            dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite para agregar una nueva ruta local."
            dicc_error["ID_PPTX"] = None
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:

            try:

                ruta_add_ajust = os.path.normpath(ruta_add)

                MiConexion = sqlite3.connect(global_path_sistema_sqlite)
                cursor = MiConexion.cursor()

                sentencia_sql_1 = "SELECT DATO FROM " + dicc_tabla_config_sistema["OTROS_DATOS"]["TABLA"] + f" WHERE CONCEPTO = 'RUTA_LOCAL' AND DATO = '{ruta_add_ajust}'"
                df_check_ruta_ya_guardada = pd.read_sql(sentencia_sql_1, MiConexion)


                if len(df_check_ruta_ya_guardada) == 0:

                    sentencia_sql_2 = "SELECT DATO FROM " + dicc_tabla_config_sistema["OTROS_DATOS"]["TABLA"] + " WHERE CONCEPTO = 'RUTA_LOCAL'"
                    
                    df_ruta = pd.read_sql(sentencia_sql_2, MiConexion)

                    lista_para_sqlite =  [["RUTA_LOCAL", ruta_add_ajust, 1]] if len(df_ruta) == 0 else [["RUTA_LOCAL", ruta_add_ajust, len(df_ruta) + 1]]
                    sentencia_sql_insert = global_dicc_tablas_y_campos_sistema["OTROS_DATOS"]["DICC_SENTENCIAS_SQL"]["INSERT"]

                    cursor.executemany(sentencia_sql_insert, lista_para_sqlite)

                    MiConexion.commit()



            except Exception as err_1:
                #se registra el posible error de guardado de las coordenadas del screenshot de la iteracion

                traceback_error = traceback.extract_tb(err_1.__traceback__)
                modulo_python = os.path.basename(traceback_error[0].filename)
                rutina_python = traceback_error[0].name
                linea_error = traceback_error[0].lineno

                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_error_1 = f"No se han podido agregar la nueva ruta a Poppler en el sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                dicc_error["PROCESO"] = opcion
                dicc_error["FASE_PROCESO"] = "Agregación nueva ruta local."
                dicc_error["ID_PPTX"] = None
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                dicc_error["MODULO_PYTHON"] = modulo_python
                dicc_error["RUTINA_PYTHON"] = rutina_python
                dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                dicc_error["TRACEBACK"] = str(err_1)

                global_lista_dicc_errores.append(dicc_error)

                pass

            finally:
                cursor.close()
                MiConexion.close()



    #####################################################################################################################################
    # DESCARGAR_SISTEMA_A_XLS --> descarga la tabla sqlite de configuracion en un excel (usando la plantilla template_plantilla_xls_config)
    #                             en una ruta seleecionada por el usuario
    #####################################################################################################################################
    elif def_varios_gui_ventana_inicio("KEY_PROCESO_SISTEMA", opcion_combobox_sistema = opcion) == "DESCARGAR_SISTEMA_A_XLS":


        if not os.path.exists(global_path_sistema_sqlite):
            #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
            #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
            #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
            mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
            mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

            dicc_error["PROCESO"] = opcion
            dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite."
            dicc_error["ID_PPTX"] = None
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:

            now = dt.datetime.now()
            ruta_xls_sistema_bbdd = os.path.join(directorio_excel, nombre_xls_descarga_parametrica + "_" + str(re.sub("[^0-9a-zA-Z]+", "_", str(now))) + ".xlsx")


            #se crea la lista lista_tablas_y_campos que es lista de listas donde cada sublista contiene:
            # --> nombre de la tabla
            # --> lista de campos para exportar a excel
            #se realiza tan solo para los tipos de tabla PARAMETRICA y RUTA_LOCAL
            #
            #se crea la lista lista_export_xls (importante crearla vacia antes de la sentencia try) que es lista de listas donde cada sublista contiene:
            # --> nombre hoja xls
            # --> df con los datos a exportar a excel

            lista_export_xls = []
            MiConexion = sqlite3.connect(global_path_sistema_sqlite)

            try:

                #PARAMETRICA
                nombre_tabla = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["NOMBRE_TABLA"]
                lista_campos_export_xls = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_LISTA_CAMPOS"]["EXPORT_EXCEL"]
                str_campos_export_xls = ", ".join(lista_campos_export_xls)

                sentencia_sql = f"SELECT {str_campos_export_xls} FROM {nombre_tabla}"

                df_temp = pd.read_sql(sentencia_sql, MiConexion)
                df_temp = df_temp.replace({np.nan: None})

                lista_export_xls.append([nombre_tabla, df_temp[lista_campos_export_xls]])


                #OTROS_DATOS
                #se localiza para cada ruta si el usuario tiene acceso y se informa en el excel
                nombre_tabla = global_dicc_tablas_y_campos_sistema["OTROS_DATOS"]["NOMBRE_TABLA"]
                lista_campos_export_xls = global_dicc_tablas_y_campos_sistema["OTROS_DATOS"]["DICC_LISTA_CAMPOS"]["EXPORT_EXCEL"]
                str_campos_export_xls = ", ".join(lista_campos_export_xls)

                sentencia_sql= f"SELECT {str_campos_export_xls} FROM {nombre_tabla} ORDER BY CONCEPTO DESC, ORDEN"

                df_temp = pd.read_sql(sentencia_sql, MiConexion)
                lista_temp = df_temp.values.tolist()

                lista_df_export = []
                for concepto, ruta, _ in lista_temp:
                    ruta_ajust = os.path.normpath(ruta)
                    acceso_ruta = "Si" if os.path.exists(ruta_ajust) else "No"

                    lista_df_export.append([concepto, ruta, acceso_ruta])

                df_temp_ajust = pd.DataFrame(lista_df_export, columns = lista_campos_export_xls)
                lista_export_xls.append([nombre_tabla, df_temp_ajust])



            except Exception as err_1:
                #se registra el posible error de guardado de las coordenadas del screenshot de la iteracion

                traceback_error = traceback.extract_tb(err_1.__traceback__)
                modulo_python = os.path.basename(traceback_error[0].filename)
                rutina_python = traceback_error[0].name
                linea_error = traceback_error[0].lineno

                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_error_1 = f"No se ha podido acceder a la las tablas del sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                dicc_error["PROCESO"] = opcion
                dicc_error["FASE_PROCESO"] = "Acceso a las tablas."
                dicc_error["ID_PPTX"] = None
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                dicc_error["MODULO_PYTHON"] = modulo_python
                dicc_error["RUTINA_PYTHON"] = rutina_python
                dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                dicc_error["TRACEBACK"] = str(err_1)

                global_lista_dicc_errores.append(dicc_error)

                pass

            else:
                #se exportan los datos a excel si lista_export_xls no es vacia
                if len(lista_export_xls) != 0:
                    app = xw.App(visible = False)
                    wb = app.books.open(template_plantilla_xls_config, update_links = False)

                    for hoja_xls, df_export in lista_export_xls:
                        ws = wb.sheets[hoja_xls]

                        if len(df_export) != 0:
                            ws["A2"].options(pd.DataFrame, header = 0, index = False, expand = "table").value = df_export

                    wb.save(ruta_xls_sistema_bbdd)
                    wb.close()
                    
                    app = xw.App(visible = True)
                    app.quit()

                    os.startfile(ruta_xls_sistema_bbdd)


            finally:
                MiConexion.close()




    #####################################################################################################################################
    # UPDATE_SCREENSHOTS_COORDENADAS
    # 
    #actualiza las coordenadas de los screenshots re-dimensionados en el sistema sqlite
    #se usa el path del sistema sqlite almacenado en la variable global global_path_sistema_sqlite cuando se importa el sistema en el app
    #####################################################################################################################################
    elif opcion == "UPDATE_SCREENSHOTS_COORDENADAS":

        tabla_sql = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["NOMBRE_TABLA"]


        if not os.path.exists(global_path_sistema_sqlite):
            #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
            #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
            #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
            mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
            mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

            dicc_error["PROCESO"] = opcion_interaccion_pptx
            dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite para actualizar las coordenadas y dimensiones de los pantallazos."
            dicc_error["ID_PPTX"] = id_pptx_selecc
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:
            MiConexion = sqlite3.connect(global_path_sistema_sqlite)
            cursor = MiConexion.cursor()


            #una vez la conexion establecida al sistema sqlite se pasa a realizar el guardado de las coordenadas de los screenshots
            #mediante bucle sobre la lista lista_screenshots_con_coordenadas (kwargs)
            #a modo recordatorio la lista lista_screenshots_con_coordenadas es lista de listas donde cada sublista contiene:
            # --> nombre screenshot
            # --> id xls
            # --> hoja xls
            # --> rango celdas xls
            # --> slide pptx
            # --> coordenadas del screensshot en el pptx (tupla convertida a string)
            #
            #las sublistas de lista_screenshots_con_coordenadas esta filtrada de tal forma que todas tienen el item coordenadas informado
            for nombre_screenshot, id_xls_iter, hoja_xls_iter, rango_celdas_iter, slide_pptx_iter, str_tupla_coordenadas in lista_screenshots_con_coordenadas:

                try:

                    sentencia_sql = f"""UPDATE {tabla_sql} SET COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX = '{str_tupla_coordenadas}'
                                    WHERE ID_PPTX = '{id_pptx_selecc}' AND NOMBRE_SCREENSHOT = '{nombre_screenshot}'"""
                    
                    cursor.execute(sentencia_sql)
                    MiConexion.commit()


                except Exception as err_1:
                    #se registra el posible error de guardado de las coordenadas del screenshot de la iteracion

                    traceback_error = traceback.extract_tb(err_1.__traceback__)
                    modulo_python = os.path.basename(traceback_error[0].filename)
                    rutina_python = traceback_error[0].name
                    linea_error = traceback_error[0].lineno

                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                    mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                    mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                    mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                    dicc_error["PROCESO"] = opcion_interaccion_pptx
                    dicc_error["FASE_PROCESO"] = "Proceso de guardado de las coordenadas en el sistema sqlite."
                    dicc_error["ID_PPTX"] = id_pptx_selecc
                    dicc_error["ID_XLS"] = id_xls_iter
                    dicc_error["RUTA_FICHERO"] = None
                    dicc_error["HOJA_XLS"] = hoja_xls_iter
                    dicc_error["RANGO_CELDAS"] = rango_celdas_iter
                    dicc_error["SLIDE_PPTX"] = slide_pptx_iter
                    dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                    dicc_error["MODULO_PYTHON"] = modulo_python
                    dicc_error["RUTINA_PYTHON"] = rutina_python
                    dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                    dicc_error["TRACEBACK"] = str(err_1)

                    global_lista_dicc_errores.append(dicc_error)

                    pass



            #se reactualiza la variable global global_df_parametrica
            try:
                sentencia_sql = f"""SELECT * FROM {tabla_sql}
                                        ORDER BY ID_PPTX"""
                
                global_df_parametrica = pd.read_sql(sentencia_sql, MiConexion)
                global_df_parametrica = global_df_parametrica.replace({np.nan: None})
                global_df_parametrica.reset_index(drop = True, inplace = True)

            except Exception as err_2:
                #se registra el posible error de extraccion en memoria de los datos guardados

                traceback_error = traceback.extract_tb(err_2.__traceback__)
                modulo_python = os.path.basename(traceback_error[0].filename)
                rutina_python = traceback_error[0].name
                linea_error = traceback_error[0].lineno

                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_log = f"No se han podido descargar en memoria Rde nuevo los datos tras el proceso de guardado de las coordeandas en el sistema sqlite.\n\nPrueba de nuevo de aqui un momento."

                dicc_error["PROCESO"] = opcion_interaccion_pptx
                dicc_error["FASE_PROCESO"] = "Proceso de guardado de las coordenadas en el sistema sqlite."
                dicc_error["ID_PPTX"] = id_pptx_selecc
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_log
                dicc_error["MODULO_PYTHON"] = modulo_python
                dicc_error["RUTINA_PYTHON"] = rutina_python
                dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                dicc_error["TRACEBACK"] = str(err_2)

                global_lista_dicc_errores.append(dicc_error)

                pass       


            #se cierra la conexion al sistema sqlite
            cursor.close()
            MiConexion.close()






    #####################################################################################################################################
    # UPDATE_SCREENSHOTS_IMAGENES_BINARIAS
    #
    #actualiza las imagenes binarias en el campo blob SCREENSHOT_PNG del sistema sqlite
    #se usa el path del sistema sqlite almacenado en la variable global global_path_sistema_sqlite cuando se importa el sistema en el app
    #####################################################################################################################################
    elif opcion == "UPDATE_SCREENSHOTS_IMAGENES_BINARIAS":


        if not os.path.exists(global_path_sistema_sqlite):
            #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
            #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
            #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
            mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
            mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

            dicc_error["PROCESO"] = opcion_interaccion_pptx
            dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite para actualizar las imagenes binarias de los pantallazos."
            dicc_error["ID_PPTX"] = id_pptx_selecc
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:
            MiConexion = sqlite3.connect(global_path_sistema_sqlite)
            cursor = MiConexion.cursor()


            #se recuperan las sentencias sql de la variable global global_dicc_tablas_y_campos_sistema
            try:
                sentencia_sql_delete = global_dicc_tablas_y_campos_sistema["UPDATE_SCREENSHOT_BINARIO"]["DICC_SENTENCIAS_SQL"]["DELETE"]
                sentencia_sql_insert = global_dicc_tablas_y_campos_sistema["UPDATE_SCREENSHOT_BINARIO"]["DICC_SENTENCIAS_SQL"]["INSERT"]
                sentencia_sql_update = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_SENTENCIAS_SQL"]["UPDATE"]

                cursor.execute(sentencia_sql_delete)
                cursor.executemany(sentencia_sql_insert, lista_datos_screenshots_binarios)
                cursor.execute(sentencia_sql_update)
                cursor.execute(sentencia_sql_delete)

                MiConexion.commit()


            except Exception as err_1:
                #se registra el posible error de guardado de las imagenes binarias de los scrrenshots en el sistema sqlite

                traceback_error = traceback.extract_tb(err_1.__traceback__)
                modulo_python = os.path.basename(traceback_error[0].filename)
                rutina_python = traceback_error[0].name
                linea_error = traceback_error[0].lineno

                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_log_1 = f"No se han podido guardar las imagenes binarias de los screenshots en el sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n"
                mensaje_log_2 = "Por lo tanto, en la ventana de configuración de los id pptx, cuando pulses en el apartado rangos de celdas, en el botón MUESTRA no saldra la muestra actualizada.\nPrueba de nuevo en otro momento."

                dicc_error["PROCESO"] = opcion_interaccion_pptx
                dicc_error["FASE_PROCESO"] = "Proceso de guardado de las imagenes binarias de los screenshots en el sistema sqlite."
                dicc_error["ID_PPTX"] = id_pptx_selecc
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2
                dicc_error["MODULO_PYTHON"] = modulo_python
                dicc_error["RUTINA_PYTHON"] = rutina_python
                dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                dicc_error["TRACEBACK"] = str(err_1)

                global_lista_dicc_errores.append(dicc_error)

                pass

            #se cierra la conexion al sistema sqlite
            cursor.close()
            MiConexion.close()




    #####################################################################################################################################
    # DESCARGA_SCREENSHOTS_PNG_PARA_MUESTRA_GUI_CONFIG_ID_PPTX
    #
    #descarga en la ruta local configurada todos los png guardados en el sistema sqlite
    #(se convierten sus imagenes binarias a png) se realiza para todos los id pptx del sistema
    # se usa en la gui de configuracion de id pptx para mostrar muestra de los screenshots configurados
    #####################################################################################################################################
    elif opcion == "DESCARGA_SCREENSHOTS_PNG_PARA_MUESTRA_GUI_CONFIG_ID_PPTX":

        if not os.path.exists(global_path_sistema_sqlite):
            #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
            #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
            #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_log_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
            mensaje_log_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
            mensaje_log_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

            dicc_error["PROCESO"] = "Descarga en la ruta local configurada de los pantallazos en formato .png para acceso a muestras en el aplicativo."
            dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite para poder descargar los pantallazos en la ruta local que configuraste en el sitema SQlite."
            dicc_error["ID_PPTX"] = None
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2 + mensaje_log_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:

            MiConexion = sqlite3.connect(global_path_sistema_sqlite)

            try:
                tabla_sql = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["NOMBRE_TABLA"]

                sentencia_sql = f"""SELECT NOMBRE_SCREENSHOT, SCREENSHOT_PNG
                                    FROM {tabla_sql}
                                    WHERE SCREENSHOT_PNG IS NOT NULL"""
                
                df_temp = pd.read_sql(sentencia_sql, MiConexion)
                df_temp.reset_index(drop = True, inplace = True)


            except Exception as err_1:
                #se registra error en el acceso a las imagenes binarias de los screenshots en el sistema sqlite

                traceback_error = traceback.extract_tb(err_1.__traceback__)
                modulo_python = os.path.basename(traceback_error[0].filename)
                rutina_python = traceback_error[0].name
                linea_error = traceback_error[0].lineno

                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_log_1 = f"No se ha podido guardar en memoria los datos necesarios del sistema sqlite para generar los ficheroes .png de los screenshots.\n"
                mensaje_log_2 = "Por lo tanto, en la ventana de configuración de los id pptx, cuando pulses en el apartado rangos de celdas, en el botón MUESTRA no saldra la muestra actualizada."

                dicc_error["PROCESO"] = "Descarga en la ruta local configurada de los pantallazos en formato .png para acceso a muestras en el aplicativo"
                dicc_error["FASE_PROCESO"] = "Acceso a las tablas."
                dicc_error["ID_PPTX"] = None
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2
                dicc_error["MODULO_PYTHON"] = modulo_python
                dicc_error["RUTINA_PYTHON"] = rutina_python
                dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                dicc_error["TRACEBACK"] = str(err_1)

                global_lista_dicc_errores.append(dicc_error)

                pass

            else:

                if not os.path.exists(global_ruta_local_config_sistema_sqlite):

                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                    mensaje_log_1 = f"La ruta local siguiente no existe (han eliminado el directorio o lo han movido de carpeta o le han cambiado el nombre).\n\n{global_ruta_local_config_sistema_sqlite}"
                    mensaje_log_2 = "Por lo tanto, en la ventana de configuración de los id pptx, cuando pulses en el apartado rangos de celdas, en el botón MUESTRA no saldra la muestra actualizada."

                    dicc_error["PROCESO"] = "Descarga en la ruta local configurada de los pantallazos en formato .png para acceso a muestras en el aplicativo"
                    dicc_error["FASE_PROCESO"] = "Acceso a la ruta del directorio local."
                    dicc_error["ID_PPTX"] = None
                    dicc_error["ID_XLS"] = None
                    dicc_error["RUTA_FICHERO"] = None
                    dicc_error["HOJA_XLS"] = None
                    dicc_error["RANGO_CELDAS"] = None
                    dicc_error["SLIDE_PPTX"] = None
                    dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2
                    dicc_error["MODULO_PYTHON"] = None
                    dicc_error["RUTINA_PYTHON"] = None
                    dicc_error["LINEA_CODIGO_PYTHON"] = None
                    dicc_error["TRACEBACK"] = None

                    global_lista_dicc_errores.append(dicc_error)

                else:

                    #se crea la carpeta donde almacenar los screenshots en formato png como muestra para poder usar en la GUI de config id pptx
                    #(se usa la ruta almacenada en la variable global global_ruta_local_config_sistema_sqlite y se actualiza con la variable global global_ruta_local_screenshots_png_muestra)
                    now = dt.datetime.now()

                    directorio_screenshots_muestra = os.path.join(global_ruta_local_config_sistema_sqlite, nombre_carpeta_screenshots_muestra + "_" + str(re.sub("[^0-9a-zA-Z]+", "_", str(now))))
                    directorio_screenshots_muestra = os.path.normpath(directorio_screenshots_muestra)

                    if not os.path.exists(directorio_screenshots_muestra):
                        os.makedirs(directorio_screenshots_muestra)

                    global_ruta_local_screenshots_png_muestra = directorio_screenshots_muestra


                    #se convierten las imagenes binarias de los screenshots a png en la ruta almacernada
                    #en la variable global global_ruta_local_screenshots_png_muestra
                    if len(df_temp) != 0:

                        for ind in df_temp.index:
                            saveas = os.path.join(global_ruta_local_screenshots_png_muestra, df_temp.iloc[ind, df_temp.columns.get_loc("NOMBRE_SCREENSHOT")] + ".PNG")

                            screenshot_binario = BytesIO(df_temp.iloc[ind, df_temp.columns.get_loc("SCREENSHOT_PNG")])
                            screenshot_png_objeto = Image.open(screenshot_binario)
                            screenshot_png_objeto.save(saveas)

                MiConexion.close()



    #AL FINALIZAR EL PROCESO: se inicializa la variable global global_proceso_en_ejecucion que se usa
    #para impedir ejecutar en la GUI otro proceso hasta que acabe el que esta en curso
    #NO HAY QUE REINICIAR global_lista_dicc_errores sino no se genera el log de posibles errores en la GUI
    if (def_varios_gui_ventana_inicio("KEY_PROCESO_SISTEMA", opcion_combobox_sistema = opcion) in ["IMPORTAR_SISTEMA", "CREAR_SISTEMA", "AGREGAR_RUTA_POPPLER", "AGREGAR_RUTA_LOCAL", "DESCARGAR_SISTEMA_A_XLS"] 
        or opcion == "DESCARGA_SCREENSHOTS_PNG_PARA_MUESTRA_GUI_CONFIG_ID_PPTX"):


        #se fija la variable global para la rutina de los threads en la GUI y asi desbloquear la ejecucion de cualquier otro proceso 
        global_proceso_en_ejecucion = "NO"

    


    #resultado funcion
    return resultado_funcion




#################################################################################################################################
#################################################################################################################################
#################################################################################################################################
# RUTINAS / FUNCIONES - GUI CONFIG ID PPTX
#################################################################################################################################
#################################################################################################################################
#################################################################################################################################


def func_check_rango_xls_correcto(rango_xls):
    #funcion que valida si el rango informado corresponde a un rango de celdas excel
    #no requiere tener ningun excel abierto (se hace madiante range_boundaries de openpyxl.utils)

    resultado_funcion = None

    if not isinstance(rango_xls, str) or not rango_xls.strip():
        resultado_funcion = True

    try:
        range_boundaries(rango_xls)
        resultado_funcion = True

    except ValueError:
        resultado_funcion = False

    return resultado_funcion



def def_varios_gui_config_id_pptx(opcion, **kwargs):
    #rutina que funciona como funcion (segun la opcion) que permite realizar diversas accciones asociadasa la GUI de configuracion de id pptx
    #(ver los comentarios al inicio de cada opcion)

    global global_df_parametrica
    global global_df_treeview_id_pptx
    global global_path_sistema_sqlite
    global global_dicc_datos_id_pptx
    global global_ruta_local_screenshots_png_muestra
    global global_dicc_tablas_y_campos_sistema
    global global_proceso_en_ejecucion
    global global_lista_dicc_errores
    global global_lista_dicc_warning


    resultado_funcion = None


    #parametros kwargs
    id_pptx = kwargs.get("id_pptx", None)
    id_pptx_path = kwargs.get("id_pptx_path", None)
    id_pptx_desc = kwargs.get("id_pptx_desc", None)
    id_pptx_tiempo_espera_max_apertura = kwargs.get("id_pptx_tiempo_espera_max_apertura", None)
    id_pptx_numero_total_slides = kwargs.get("id_pptx_numero_total_slides", None)

    id_pptx_excluir_check_path = kwargs.get("id_pptx_excluir_check_path", None)
    id_xls_excluir_check_path = kwargs.get("id_xls_excluir_check_path", None)

    id_xls = kwargs.get("id_xls", None)
    id_xls_path = kwargs.get("id_xls_path", None)
    id_xls_desc = kwargs.get("id_xls_desc", None)
    id_xls_tiempo_espera_max_apertura = kwargs.get("id_xls_tiempo_espera_max_apertura", None)
    id_xls_actualizar_vinculos_otros_xls = kwargs.get("id_xls_actualizar_vinculos_otros_xls", None)
    lista_datos_items_selecc_update_desc = kwargs.get("lista_datos_items_selecc_update_desc", None)

    hoja_xls = kwargs.get("hoja_xls", None)
    rango_celdas = kwargs.get("rango_celdas", None)
    slide_pptx = kwargs.get("slide_pptx", None)

    nombre_screenshot = kwargs.get("nombre_screenshot", None)

    ruta_fichero = kwargs.get("ruta_fichero", None)

    lista_rangos_celdas_antes_update = kwargs.get("lista_rangos_celdas_antes_update", None)
    lista_rangos_celdas_despues_update = kwargs.get("lista_rangos_celdas_despues_update", None)


    #IMPORTANTE: se reinician las variables globales global_lista_dicc_errores y global_lista_dicc_warning
    #directamente en la GUI al pulsar el boton GUARDAR



    ####################################################################################################################################################
    # SELECCION_ID_PPTX
    ####################################################################################################################################################
    if opcion == "SELECCION_ID_PPTX":
        #devuelve una serie de datos tras seleccionar un un idpptx en el el combobox de la GUI y tras pulsar el botón VER:
        # --> id_pptx
        # --> id_pptx_desc
        # --> id_pptx_path
        # --> id_pptx_tiempo_espera_max_apertura
        # --> numero_total_slides 
        # --> df_widget_treeview_id_xls

        id_pptx_desc = global_dicc_datos_id_pptx[id_pptx]["id_pptx_desc"]
        id_pptx_path = global_dicc_datos_id_pptx[id_pptx]["ruta_plantilla_pptx"]
        tiempo_espera_max_apertura = global_dicc_datos_id_pptx[id_pptx]["tiempo_espera_max_apertura_pptx"]
        numero_total_slides = global_dicc_datos_id_pptx[id_pptx]["numero_total_slides"]
        df_widget_treeview_id_xls = global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"]

        tiempo_espera_max_apertura_ajust = int(float(tiempo_espera_max_apertura))

        resultado_funcion = id_pptx, id_pptx_desc, id_pptx_path, tiempo_espera_max_apertura_ajust, numero_total_slides, df_widget_treeview_id_xls






    #############################################################################################################
    # UPDATE_EN_MEMORIA_DATOS_EN_PANTALLA_ID_PPTX_Y_ID_XLS
    #############################################################################################################
    elif opcion == "UPDATE_EN_MEMORIA_DATOS_EN_PANTALLA_ID_PPTX_Y_ID_XLS":
        #permite guardar en la variable global global_dicc_datos_id_pptx los datos en pantalla del id pptx seleccionado
        #y asimismo los datos generales del id xls seleccionado en pantalla (si lo esta)
        #sino se realiza este ajuste estos datos si se han modificado en pantalla no se guarda ewn el sistema sqlite
        #es necesario cuando o bien se selecciona otro id pptx o bien se pulsa el boton GUARDAR
        #
        #los datos de id_pptx_desc y tiempo_espera_max_apertura_pptx se establecen a "-" y 0 respectivamente si el usuario no los tiene informados
        #los datos de desc_id_xls y tiempo_espera_max_apertura_xls se establecen a "-" y 0 respectivamente si el usuario no los tiene informados


        #se guardan los datos del id pptx en pantalla antes de actualizar el combobox de seleccion de id pptx y pulsar el boton VER
        if id_pptx is not None:

            id_pptx_desc_ajust = (id_pptx_desc if len(id_pptx_desc.replace(" ", "").replace("\t", "")) != 0 else "-")
            id_pptx_tiempo_espera_max_apertura_ajust = int(float(id_pptx_tiempo_espera_max_apertura)) if len(id_pptx_tiempo_espera_max_apertura) != 0 else 0

            global_dicc_datos_id_pptx[id_pptx]["id_pptx_desc"] = id_pptx_desc_ajust
            global_dicc_datos_id_pptx[id_pptx]["ruta_plantilla_pptx"] = id_pptx_path
            global_dicc_datos_id_pptx[id_pptx]["tiempo_espera_max_apertura_pptx"] = id_pptx_tiempo_espera_max_apertura_ajust
            global_dicc_datos_id_pptx[id_pptx]["numero_total_slides"] = id_pptx_numero_total_slides


        #se guardan los datos del id xls seleccionado (si lo esta) en el treeview asociado al id pptx en pantalla
        #antes de actualizar el combobox de seleccion de id pptx y pulsar el boton VER
        if id_xls is not None:

            for ind, dicc in enumerate(global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]):
                if dicc["id_xls"] == id_xls:

                    id_xls_desc_ajust = (id_xls_desc if len(id_xls_desc.replace(" ", "").replace("\t", "")) != 0 else "-")
                    id_xls_tiempo_espera_max_apertura_ajust = int(float(id_xls_tiempo_espera_max_apertura)) if len(id_xls_tiempo_espera_max_apertura) != 0 else 0

                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["desc_id_xls"] = id_xls_desc_ajust
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["ruta_xls_origen"] = id_xls_path
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["tiempo_espera_max_apertura_xls"] = id_xls_tiempo_espera_max_apertura_ajust
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["actualizar_vinculos_otros_xls"] = id_xls_actualizar_vinculos_otros_xls

                    break


    ####################################################################################################################################################
    # GUARDAR_CONFIGURACIONES_EN_SISTEMA_SQLITE
    ####################################################################################################################################################
    elif opcion == "GUARDAR_CONFIGURACIONES_EN_SISTEMA_SQLITE":
        #guarda las configuraciones en el sistema sqlite basandose en los datos de la variable global global_dicc_datos_id_pptx
        #recalcula tambien todos los nombres de los screenshots
        #funciona a su vez como funcion para generar en la GUI un mensaje de warning que indica:
        # --> si existen id pptx sin id xls asignados
        # --> si existen id xls sin rangos de celdas configurados


        #se prepara la lista de datos lista_datos_config_id_pptx para poder actualizar la tabla del sistema sqlite
        #siempre y cuando global_dicc_datos_id_pptx noeste vacio



        #AL INICIAR EL PROCESO: se inicializa la variable global global_lista_dicc_errores como lista vacia
        #se fija la variable global global_proceso_en_ejecucion para la rutina de los threads en la GUI y asi bloquear la ejecucion de cualquier otro proceso 
        #hasta que el actual en curso no haya finalizado
        global_proceso_en_ejecucion = "SI"


        ###########################################################################################################
        #CASO 1 - se eliminan todas las configuraciones en memoria para dejar el sistema sqlite vacio
        #         se realiza directame el borrado en sqlite
        ###########################################################################################################
        if len(global_dicc_datos_id_pptx) == 0 and len(global_df_parametrica) != 0:


            if not os.path.exists(global_path_sistema_sqlite):
                #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
                #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
                #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                dicc_error["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite."
                dicc_error["ID_PPTX"] = None
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                dicc_error["MODULO_PYTHON"] = None
                dicc_error["RUTINA_PYTHON"] = None
                dicc_error["LINEA_CODIGO_PYTHON"] = None
                dicc_error["TRACEBACK"] = None

                global_lista_dicc_errores.append(dicc_error)

            else:

                MiConexion = sqlite3.connect(global_path_sistema_sqlite)
                cursor = MiConexion.cursor()

                try:

                    sentencia_sql_delete = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_SENTENCIAS_SQL"]["DELETE"]

                    cursor.execute(sentencia_sql_delete)
                    MiConexion.commit()


                    #se actualiza en memoria el df global_df_parametrica
                    tabla_sql = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["NOMBRE_TABLA"]
                    sentencia_sql_select = f"SELECT * FROM {tabla_sql} ORDER BY ID_PPTX, ID_XLS, HOJA_XLS, SLIDE_PPTX"

                    global_df_parametrica = pd.read_sql(sentencia_sql_select, MiConexion)
                    global_df_parametrica = global_df_parametrica.replace({np.nan: None})
                    global_df_parametrica.reset_index(drop = True, inplace = True)

                    cursor.close()
                    MiConexion.close()

                    #se actualiza global_df_treeview_id_pptx para poder actualizar el treeview en la gui de ventana de inicio
                    global_df_treeview_id_pptx = def_varios_gui_ventana_inicio("DF_TREEVIEW_ID_PPTX")


                except Exception as err_1:
                    #se registra el posible error de guardado

                    traceback_error = traceback.extract_tb(err_1.__traceback__)
                    modulo_python = os.path.basename(traceback_error[0].filename)
                    rutina_python = traceback_error[0].name
                    linea_error = traceback_error[0].lineno

                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                    mensaje_error_1 = f"No se han podido guardar las configuraciones de los rangos de celdas en el sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                    mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                    mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                    dicc_error["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                    dicc_error["FASE_PROCESO"] = "Proceso de guardado de los datos en las tablas del sistema sqlite."
                    dicc_error["ID_PPTX"] = None
                    dicc_error["ID_XLS"] = None
                    dicc_error["RUTA_FICHERO"] = None
                    dicc_error["HOJA_XLS"] = None
                    dicc_error["RANGO_CELDAS"] = None
                    dicc_error["SLIDE_PPTX"] = None
                    dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                    dicc_error["MODULO_PYTHON"] = modulo_python
                    dicc_error["RUTINA_PYTHON"] = rutina_python
                    dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                    dicc_error["TRACEBACK"] = str(err_1)

                    global_lista_dicc_errores.append(dicc_error)

                    pass



        ###########################################################################################################
        #CASO 2 - las configuraciones no son vacias
        ###########################################################################################################
        elif len(global_dicc_datos_id_pptx) != 0:

            lista_datos_config_id_pptx = []

            for id_pptx, dicc_id_pptx in global_dicc_datos_id_pptx.items():

                id_pptx_desc = dicc_id_pptx["id_pptx_desc"]
                ruta_plantilla_pptx = dicc_id_pptx["ruta_plantilla_pptx"]
                tiempo_espera_max_apertura_pptx = dicc_id_pptx["tiempo_espera_max_apertura_pptx"]
                numero_total_slides = dicc_id_pptx["numero_total_slides"]

                lista_dicc_id_xls = dicc_id_pptx["lista_dicc_id_xls"]

                if isinstance(lista_dicc_id_xls, list):
                    

                    #si lista_dicc_id_xls es vacia se insertan los datos en lista_datos_config_id_pptx tan solo del id_pptx de la iteracion
                    #los datos de id xls y rangos de celdas se agregan vacios
                    if len(lista_dicc_id_xls) == 0:

                        lista_datos_config_id_pptx.append([id_pptx
                                                            , id_pptx_desc
                                                            , ruta_plantilla_pptx
                                                            , tiempo_espera_max_apertura_pptx
                                                            , numero_total_slides
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            , None
                                                            ])
                        
                        #se almacena el warning en la variable global global_lista_dicc_warning para avisar de que hay id pptx sin id xls asignados
                        dicc_warning = dicc_modelo_warning_y_errores_procesos_app["WARNING"].copy()

                        dicc_warning["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                        dicc_warning["ID_PPTX"] = id_pptx
                        dicc_warning["ID_XLS"] = None
                        dicc_warning["HOJA_XLS"] = None
                        dicc_warning["RANGO_CELDAS"] = None
                        dicc_warning["SLIDE_PPTX"] = None
                        dicc_warning["COMENTARIO"] = "El id pptx no tiene ningún id xls asignado."

                        global_lista_dicc_warning.append(dicc_warning)


    
                    #si lista_dicc_id_xls NO es vacia se insertan los datos en lista_datos_config_id_pptx mediante bucle sobre lista_dicc_id_xls   
                    elif len(lista_dicc_id_xls) != 0:

                        cont_screenshot = 0
                        for dicc_id_xls in lista_dicc_id_xls:

                            id_xls = dicc_id_xls["id_xls"]
                            desc_id_xls = dicc_id_xls["desc_id_xls"]
                            ruta_xls_origen = dicc_id_xls["ruta_xls_origen"]
                            tiempo_espera_max_apertura_xls = int(dicc_id_xls["tiempo_espera_max_apertura_xls"])
                            actualizar_vinculos_otros_xls = dicc_id_xls["actualizar_vinculos_otros_xls"]
                            lista_hojas_xls = str(dicc_id_xls["lista_hojas_xls"]) #se convierte la lista en string

                            df_widget_treeview_rangos_celdas = dicc_id_xls["df_widget_treeview_rangos_celdas"]

                            if isinstance(df_widget_treeview_rangos_celdas, pd.DataFrame):

                                #se agregan los datos a lista_datos_config_id_pptx mediante bucle sobre los datos de df_widget_treeview_rangos_celdas
                                #convertidos a lista (solo si df_widget_treeview_rangos_celdas no es vacio)
                                if len(df_widget_treeview_rangos_celdas) != 0:

                                    #se recalculan los nombres de los screenshots (solo si df_widget_treeview_rangos_celdas no esta vacio)
                                    #(se usa la funcion lambda lambda_nombre_screenshots que normaliza en 3 digitos el indice del df + 1)
                                    if len(df_widget_treeview_rangos_celdas) != 0:
                                        df_widget_treeview_rangos_celdas.reset_index(drop = True, inplace = True)

                                        lambda_nombre_screenshots = lambda cont: f"{cont:03}"
                                        df_widget_treeview_rangos_celdas["NOMBRE_SCREENSHOT"] = df_widget_treeview_rangos_celdas.apply(lambda x: f"{nomenclatura_nombres_screenshots}_{id_pptx}_{lambda_nombre_screenshots(x.name + 1 + cont_screenshot)}", axis = 1)

                                        cont_screenshot = cont_screenshot + len(df_widget_treeview_rangos_celdas)

                                    #se convierte el df en lista de listas df_widget_treeview_rangos_celdas
                                    #y mediante bucle sobre esta lista se alimenta la lista lista_datos_config_id_pptx
                                    #a modo recordatorio las columnas de lista_datos_df_rangos_celdas son ID_XLS, HOJA_XLS, RANGO_XLS, SLIDE_PPTX, NOMBRE_SCREENSHOT, SCREENSHOT_PNG y COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX
                                    lista_datos_df_rangos_celdas = df_widget_treeview_rangos_celdas.values.tolist()

                                    for _, hoja_xls, rango_celdas, slide_pptx, nombre_screenshot, screenshot_png, coordenadas_screenshot in lista_datos_df_rangos_celdas:

                                        rango_celdas_strip = rango_celdas.strip()#por si el usuario ha puesto espacios en blanco antes o despues del rango

                                        lista_datos_config_id_pptx.append([id_pptx
                                                                            , id_pptx_desc
                                                                            , ruta_plantilla_pptx
                                                                            , tiempo_espera_max_apertura_pptx
                                                                            , numero_total_slides
                                                                            , id_xls
                                                                            , desc_id_xls
                                                                            , ruta_xls_origen
                                                                            , tiempo_espera_max_apertura_xls
                                                                            , actualizar_vinculos_otros_xls
                                                                            , lista_hojas_xls
                                                                            , hoja_xls
                                                                            , rango_celdas_strip
                                                                            , slide_pptx
                                                                            , nombre_screenshot
                                                                            , screenshot_png
                                                                            , coordenadas_screenshot
                                                                            ])

                                
                                #se agregan los datos a lista_datos_config_id_pptx con los datos de rangos de celdas a None
                                #(solo si df_widget_treeview_rangos_celdas es vacio)     
                                #asimismo se informa mensaje_check_id_xls_sin_rangos_celdas     
                                elif len(df_widget_treeview_rangos_celdas) == 0:

                                    lista_datos_config_id_pptx.append([id_pptx
                                                                        , id_pptx_desc
                                                                        , ruta_plantilla_pptx
                                                                        , tiempo_espera_max_apertura_pptx
                                                                        , numero_total_slides
                                                                        , id_xls
                                                                        , desc_id_xls
                                                                        , ruta_xls_origen
                                                                        , tiempo_espera_max_apertura_xls
                                                                        , actualizar_vinculos_otros_xls
                                                                        , lista_hojas_xls
                                                                        , None
                                                                        , None
                                                                        , None
                                                                        , None
                                                                        , None
                                                                        , None
                                                                        ])

                                    #se almacena el warning en la variable global global_lista_dicc_warning para avisar de que hay id pptx
                                    #con id xls que no tienen aun rangos de celdas configurados
                                    dicc_warning = dicc_modelo_warning_y_errores_procesos_app["WARNING"].copy()

                                    dicc_warning["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                                    dicc_warning["ID_PPTX"] = id_pptx
                                    dicc_warning["ID_XLS"] = id_xls
                                    dicc_warning["HOJA_XLS"] = None
                                    dicc_warning["RANGO_CELDAS"] = None
                                    dicc_warning["SLIDE_PPTX"] = None
                                    dicc_warning["COMENTARIO"] = "El id xls asociado al id pptx no tiene ningún rango de celdas configurado."

                                    global_lista_dicc_warning.append(dicc_warning)



            #se guardan los datos en el sistema sqlite
            #(tan solo si la lista lista_datos_config_id_pptx no esta vacia)
            #se recuperan las sentencias sql de la variable global global_dicc_tablas_y_campos_sistema
            if len(lista_datos_config_id_pptx) != 0:


                if not os.path.exists(global_path_sistema_sqlite):
                    #el chequeo de la conexion al sistema sqlite no se hace por bloque de try except pq si la ruta global_path_sistema_sqlite no existe
                    #(si ha habido algun movimiento entre carpetas de la bbdd o eliminacion del fichero de la bbdd) sqlite3.connect crea una bbdd vacia
                    #con la misma ubicacion global_path_sistema_sqlite por lo que hay que realizar el check usando os.path.exists

                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                    mensaje_error_1 = f"No se ha podido crear conexion al sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                    mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                    mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                    dicc_error["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                    dicc_error["FASE_PROCESO"] = "Conexión al sistema sqlite."
                    dicc_error["ID_PPTX"] = None
                    dicc_error["ID_XLS"] = None
                    dicc_error["RUTA_FICHERO"] = None
                    dicc_error["HOJA_XLS"] = None
                    dicc_error["RANGO_CELDAS"] = None
                    dicc_error["SLIDE_PPTX"] = None
                    dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                    dicc_error["MODULO_PYTHON"] = None
                    dicc_error["RUTINA_PYTHON"] = None
                    dicc_error["LINEA_CODIGO_PYTHON"] = None
                    dicc_error["TRACEBACK"] = None

                    global_lista_dicc_errores.append(dicc_error)

                else:

                    MiConexion = sqlite3.connect(global_path_sistema_sqlite)
                    cursor = MiConexion.cursor()

                    try:
                        sentencia_sql_delete = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_SENTENCIAS_SQL"]["DELETE"]
                        sentencia_sql_insert = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_SENTENCIAS_SQL"]["INSERT"]

                        cursor.execute(sentencia_sql_delete)
                        cursor.executemany(sentencia_sql_insert, lista_datos_config_id_pptx)
                        MiConexion.commit()


                        #se actualiza en memoria el df global_df_parametrica
                        tabla_sql = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["NOMBRE_TABLA"]
                        sentencia_sql_select = f"SELECT * FROM {tabla_sql} ORDER BY ID_PPTX, ID_XLS, HOJA_XLS, SLIDE_PPTX"

                        global_df_parametrica = pd.read_sql(sentencia_sql_select, MiConexion)
                        global_df_parametrica = global_df_parametrica.replace({np.nan: None})
                        global_df_parametrica.reset_index(drop = True, inplace = True)

                        cursor.close()
                        MiConexion.close()

                        #se actualiza global_df_treeview_id_pptx para poder actualizar el treeview en la gui de ventana de inicio
                        global_df_treeview_id_pptx = def_varios_gui_ventana_inicio("DF_TREEVIEW_ID_PPTX")


                    except Exception as err_1:
                        #se registra el posible error de guardado

                        traceback_error = traceback.extract_tb(err_1.__traceback__)
                        modulo_python = os.path.basename(traceback_error[0].filename)
                        rutina_python = traceback_error[0].name
                        linea_error = traceback_error[0].lineno

                        dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                        mensaje_error_1 = f"No se han podido guardar las configuraciones de los rangos de celdas en el sistema sqlite ubicado en:\n{global_path_sistema_sqlite}.\n\n"
                        mensaje_error_2 = f"Comprueba que el sistema existe en la ruta indicada o que tienes acceso al directorio donde esta guardado.\n\n"
                        mensaje_error_3 = "Si tienes acceso a dicho directorio y que el sistema existe, prueba más adelante."

                        dicc_error["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                        dicc_error["FASE_PROCESO"] = "Proceso de guardado de los datos en las tablas del sistema sqlite."
                        dicc_error["ID_PPTX"] = None
                        dicc_error["ID_XLS"] = None
                        dicc_error["RUTA_FICHERO"] = None
                        dicc_error["HOJA_XLS"] = None
                        dicc_error["RANGO_CELDAS"] = None
                        dicc_error["SLIDE_PPTX"] = None
                        dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2 + mensaje_error_3
                        dicc_error["MODULO_PYTHON"] = modulo_python
                        dicc_error["RUTINA_PYTHON"] = rutina_python
                        dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                        dicc_error["TRACEBACK"] = str(err_1)

                        global_lista_dicc_errores.append(dicc_error)

                        pass

        #AL FINALIZAR EL PROCESO: se inicializa la variable global global_lista_dicc_errores como lista vacia
        #se fija la variable global global_proceso_en_ejecucion para la rutina de los threads en la GUI y asi desbloquear la ejecucion de cualquier otro proceso 
        global_proceso_en_ejecucion = "NO"



    ####################################################################################################################################################
    # UPDATE_NUMERO_SLIDES_PPTX_Y_LISTA_HOJAS_XLS
    ####################################################################################################################################################
    elif opcion == "UPDATE_NUMERO_SLIDES_PPTX_Y_LISTA_HOJAS_XLS":
        #permite actualizar en todos los pptx configurados el numero total de slides que contiene cada uno
        #asimismo, actualiza en cada excel de origen configurado (sea cual sea su id pptx asociado) la lista de hojas que contiene
        #al final del proceso, genera un log en formato .txt en caso de que haya rangos de celdas con número de slide pptx por encima
        #del numero total de slides que contiene el id pptx asociado
        #este log tambien informa de los casos en los que haya rangos de celdas
        #con una hoja excel configurada que ya no existe en el excel de origen asociado
        #
        #los desfases se almacenan en la variable global global_lista_dicc_warning


        #se actualiza en la variable global global_dicc_datos_id_pptx los numeros de slides de los pptx y las listas de hojas de cada id xls
        #se crea la lista lista_check_config que sirve para informar, es lista de listas donde cadea sublista contiene:
        # --> id_pptx
        # --> numero_total_slides
        # --> id_xls
        # --> lista_hojas_xls
        # --> hoja_xls (rangos celdas)
        # --> rango celdas
        # --> slide pptx (rangos celdas)
        #
        # --> mensaje para log que avisa si el numero de slide del rango celdas esta por encima del numero total de slides del pptx y/o
        #     la hoja xls del rango de celdas no esta en el excel de origen
        lista_check_config = []
        for id_pptx_iter, dicc_id_pptx in global_dicc_datos_id_pptx.items():

            #se actualiza el numero de slides
            id_pptx_path = dicc_id_pptx["ruta_plantilla_pptx"]
            tiempo_espera_max_apertura_pptx = dicc_id_pptx["tiempo_espera_max_apertura_pptx"]

            id_pptx_path_ajust = os.path.normpath(id_pptx_path)


            #lsi la ruta del excel existe se actualizan la lista de hojas sino se almacena
            #el warning en la variable global global_lista_dicc_warning
            if os.path.exists(id_pptx_path_ajust):

                pptx_objeto = Presentation(os.path.normpath(id_pptx_path_ajust))
                time.sleep(tiempo_espera_max_apertura_pptx)

                numero_total_slides = int(float(len([slide_pptx for slide_pptx in pptx_objeto.slides])))

                global_dicc_datos_id_pptx[id_pptx_iter]["numero_total_slides"] = numero_total_slides

            else:

                dicc_warning = dicc_modelo_warning_y_errores_procesos_app["WARNING"].copy()

                mensaje_log = f"La ruta del powerpoint de destino no existe (han eliminado el fichero o lo han movido de carpeta o le han cambiado el nombre):\n\n{id_pptx_path_ajust}"

                dicc_warning["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                dicc_warning["FASE_PROCESO"] = "Actualización del número total de slides de powerpoint asociado al id pptx."
                dicc_warning["ID_PPTX"] = id_pptx_iter
                dicc_warning["ID_XLS"] = None
                dicc_warning["HOJA_XLS"] = None
                dicc_warning["RANGO_CELDAS"] = None
                dicc_warning["SLIDE_PPTX"] = None
                dicc_warning["COMENTARIO"] = mensaje_log

                global_lista_dicc_warning.append(dicc_warning)


                #no se actualiza el numero total de slides (se conserva el que esta en memoria)
                numero_total_slides = global_dicc_datos_id_pptx[id_pptx_iter]["numero_total_slides"]






            #se actualiza la lista de hojas de cada excel de origen
            lista_dicc_id_xls = dicc_id_pptx["lista_dicc_id_xls"]

            if len(lista_dicc_id_xls) != 0:

                for ind, dicc_id_xls in enumerate(lista_dicc_id_xls):

                    id_xls = dicc_id_xls["id_xls"]
                    id_xls_path = dicc_id_xls["ruta_xls_origen"]
                    id_xls_path_ajust = os.path.normpath(id_xls_path)


                    #lsi la ruta del excel existe se actualizan la lista de hojas sino se almacena
                    #el warning en la variable global global_lista_dicc_warning
                    if os.path.exists(id_xls_path_ajust):
                        lista_hojas_xls = def_varios("LISTA_HOJAS_XLS_DESDE_PATH_XLS", id_xls_path = id_xls_path_ajust)
                        global_dicc_datos_id_pptx[id_pptx_iter]["lista_dicc_id_xls"][ind]["lista_hojas_xls"] = str(lista_hojas_xls)

                    else:

                        dicc_warning = dicc_modelo_warning_y_errores_procesos_app["WARNING"].copy()

                        mensaje_log = f"La ruta del excel de origen no existe (han eliminado el fichero o lo han movido de carpeta o le han cambiado el nombre):\n\n{id_pptx_path_ajust}"

                        dicc_warning["PROCESO"] = "Guardado de las configuraciones de rangos de celdas por id pptx en el sistema sqlite."
                        dicc_warning["FASE_PROCESO"] = "Actualización de la lista de hojas disponibles en el excel."
                        dicc_warning["ID_PPTX"] = id_pptx_iter
                        dicc_warning["ID_XLS"] = id_xls
                        dicc_warning["HOJA_XLS"] = None
                        dicc_warning["RANGO_CELDAS"] = None
                        dicc_warning["SLIDE_PPTX"] = None
                        dicc_warning["COMENTARIO"] = mensaje_log

                        global_lista_dicc_warning.append(dicc_warning)


                        #no se actualiza la lista de hojas (se conserva los que estan en memoria)
                        lista_hojas_xls = global_dicc_datos_id_pptx[id_pptx_iter]["lista_dicc_id_xls"][ind]["lista_hojas_xls"]



                    


                    #se informa la lista lista_check_config que permite localizar si hay desfases en los rangos de celdas 
                    #tras actualizar el numero total de slides de cada pptx y la lista de hojas de cada excel con rangos de celdas configurados
                    df_widget_treeview_rangos_celdas = dicc_id_xls["df_widget_treeview_rangos_celdas"]

                    if len(df_widget_treeview_rangos_celdas) != 0:
                        lista_rangos_celdas = df_widget_treeview_rangos_celdas[["HOJA_XLS", "RANGO_XLS", "SLIDE_PPTX"]].values.tolist()

                        for hoja_xls_iter, rango_celdas_iter, slide_pptx_iter in lista_rangos_celdas:

                            lista_check_config.append([id_pptx_iter
                                                        , numero_total_slides
                                                        , id_xls
                                                        , lista_hojas_xls
                                                        , hoja_xls_iter
                                                        , rango_celdas_iter
                                                        , slide_pptx_iter
                                                        ])


        #se localiza si hay que generar un log por configuraciones desfasadas iterando por la lista lista_check_config
        #y se almacenan en la variable global global_lista_dicc_warning
        for id_pptx_iter, numero_total_slides_iter, id_xls_iter, lista_hojas_xls_iter, hoja_xls_iter, rango_celdas_iter, slide_pptx_iter  in lista_check_config:

            mensaje_log_id_pptx = f"el id pptx tiene un número total de slides ({numero_total_slides_iter}) inferior a la slide de destino del rango de celdas ({slide_pptx_iter})." if slide_pptx_iter > numero_total_slides_iter else ""
            mensaje_log_id_xls = f"la hoja '{hoja_xls_iter}' configurada en el rango de celdas no esta en el excel de origen asociado al id xls {id_xls_iter}." if hoja_xls_iter not in lista_hojas_xls_iter else ""

            mensaje_log = (mensaje_log_id_pptx + "\n" + mensaje_log_id_xls
                           if len(mensaje_log_id_pptx) != 0 and len(mensaje_log_id_xls) != 0
                           else
                           mensaje_log_id_pptx
                           if len(mensaje_log_id_pptx) != 0 and len(mensaje_log_id_xls) == 0
                           else
                           mensaje_log_id_xls
                           if len(mensaje_log_id_pptx) == 0 and len(mensaje_log_id_xls) != 0
                           else
                           None
                            )
            
            if mensaje_log is not None:

                dicc_warning = dicc_modelo_warning_y_errores_procesos_app["WARNING"].copy()

                dicc_warning["PROCESO"] = "Desfase configuración rangos de celdas por refresh de los powerpoint y excel de origen."
                dicc_warning["ID_PPTX"] = id_pptx_iter
                dicc_warning["ID_XLS"] = id_xls_iter
                dicc_warning["HOJA_XLS"] = hoja_xls_iter
                dicc_warning["RANGO_CELDAS"] = rango_celdas_iter
                dicc_warning["SLIDE_PPTX"] = slide_pptx_iter
                dicc_warning["COMENTARIO"] = mensaje_log

                global_lista_dicc_warning.append(dicc_warning)





    ####################################################################################################################################################
    # CHECK_SI_PATH_ID_PPTX_YA_ASIGNADO
    ####################################################################################################################################################
    elif opcion == "CHECK_SI_PATH_ID_PPTX_YA_ASIGNADO":
        #genera mensaje de warning en caso de que la ruta ya esta asignada a otros id pptx
        #usa el parametro kwargs id_pptx_excluir_check_pathpara distinguir los casos:
        # --> el check se hace al crear un nuevo id pptx (no se informe id_pptx_excluir_check_path)
        # --> el check se hace al atualizar el path de un idpptx existente (se informa id_pptx_excluir_check_path para excluir del check el id pptx que se cambia)

        id_pptx_path_ajust = os.path.normpath(id_pptx_path)

        lista_path_id_pptx_ya_asignados_sin_exclusion = [id_pptx_iter for id_pptx_iter, dicc in global_dicc_datos_id_pptx.items()
                                                        if os.path.normpath(dicc["ruta_plantilla_pptx"]) == id_pptx_path_ajust]

        lista_path_id_pptx_ya_asignados_con_exclusion = [id_pptx_iter for id_pptx_iter, dicc in global_dicc_datos_id_pptx.items()
                                                        if id_pptx_excluir_check_path is not None 
                                                            and id_pptx_iter != id_pptx_excluir_check_path 
                                                            and  os.path.normpath(dicc["ruta_plantilla_pptx"]) == id_pptx_path_ajust
                                                        ]

        lista_path_id_pptx_ya_asignados = lista_path_id_pptx_ya_asignados_sin_exclusion if id_pptx_excluir_check_path is None else lista_path_id_pptx_ya_asignados_con_exclusion

        str_path_id_pptx_ya_asignados = ", ".join(lista_path_id_pptx_ya_asignados) if len(lista_path_id_pptx_ya_asignados) else None

        resultado_funcion = f"La ruta que has configurado ya está asignada a los id pptx siguientes: {str_path_id_pptx_ya_asignados}.\n\nDeseas continuar?" if str_path_id_pptx_ya_asignados is not None else None




    ####################################################################################################################################################
    # ASIGNAR_ID_PPTX_NUEVO
    ####################################################################################################################################################
    elif opcion == "ASIGNAR_ID_PPTX_NUEVO":
        #devuelve un id pptx nuevo correlativo con el ultimo creado
        #se usa una funcion lambda que normaliza el numero con 3 digitos
        #en el caso de que se hayan eliminado id pptx entre el 1ero y ultimo, el id pptx que se asigna es el 1er hueco que encuentra desde el 1ero
        #asimismo devuelve el número de slides que contiene el pptx pasado por parametro kwargs
        #(es nesecario pq el numero de slides se informa en la GUI)

        lambda_id_pptx = lambda cont: f"{cont:03}"

        lista_id_pptx_config = list(global_dicc_datos_id_pptx.keys())
        lista_id_pptx_config = sorted(lista_id_pptx_config)

        lista_id_pptx_config_huecos_por_cubrir = [f"{nombre_id_pptx}_{lambda_id_pptx(ind + 1)}" 
                                                  for ind, id_pptx_config in enumerate(lista_id_pptx_config)
                                                  if ind + 1 != int(id_pptx_config[-3:])
                                                    and f"{nombre_id_pptx}_{lambda_id_pptx(ind + 1)}" not in lista_id_pptx_config
                                                    ]

        #nuevo id pptx
        indice_correlativo = len(lista_id_pptx_config) + 1
        nuevo_id_pptx = f"{nombre_id_pptx}_{lambda_id_pptx(indice_correlativo)}" if len(lista_id_pptx_config_huecos_por_cubrir) == 0 else lista_id_pptx_config_huecos_por_cubrir[0]


        #numero_total_slides (Aqui no se hace un time.sleep despues de la apertura del pptx pq todavia no esta configurado)
        id_pptx_path_ajust = os.path.normpath(id_pptx_path)
        pptx_objeto = Presentation(os.path.normpath(id_pptx_path_ajust))

        numero_total_slides = int(float(len([slide_pptx for slide_pptx in pptx_objeto.slides])))


        resultado_funcion = nuevo_id_pptx, numero_total_slides



    ####################################################################################################################################################
    # CREAR_NUEVO_ID_PPTX
    ####################################################################################################################################################
    elif opcion == "CREAR_NUEVO_ID_PPTX":
        #crea un nuevo id pptx que se almacena en memoria en la variable global global_dicc_datos_id_pptx agregandole una key adicional
        #cuyo nombre es el id pptx (parametro kwargs id_pptx) y se completa el diccionario asociado:
        # --> id_pptx_desc  (kwargs id_pptx_desc)
        # --> ruta_plantilla_pptx (kwargs id_pptx_path)
        # --> tiempo_espera_max_apertura_pptx (kwargs id_pptx_tiempo_espera_max_apertura)
        # --> numero_total_slides (kwargs id_pptx_numero_total_slides)
        # --> df_widget_treeview_id_xls (dataframe vacio con 2 columnas: ID_XLS y DESC_ID_XLS)
        # --> lista_dicc_id_xls (lista vacia)


        id_pptx_path_ajust = os.path.normpath(id_pptx_path)

        #se convierte id_pptx_tiempo_espera_max_apertura en entero
        id_pptx_tiempo_espera_max_apertura_ajust = int(float(id_pptx_tiempo_espera_max_apertura))

        dicc_id_pptx_nuevo = {str(id_pptx):
                                {"id_pptx_desc": id_pptx_desc
                                , "ruta_plantilla_pptx": id_pptx_path_ajust
                                , "tiempo_espera_max_apertura_pptx": id_pptx_tiempo_espera_max_apertura_ajust
                                , "numero_total_slides": id_pptx_numero_total_slides
                                , "df_widget_treeview_id_xls": pd.DataFrame(columns = ["ID_XLS", "DESC_ID_XLS"])
                                , "lista_dicc_id_xls": []
                                }
                            }

        global_dicc_datos_id_pptx.update(dicc_id_pptx_nuevo)



    ####################################################################################################################################################
    # ELIMINAR_ID_PPTX
    ####################################################################################################################################################
    elif opcion == "ELIMINAR_ID_PPTX":
        #elimina el id_pptx (kwargs) en la variable global global_dicc_datos_id_pptx y genera la lista de opciones actualiza
        #para el combobox de seleccion de id pptx en la GUI

        global_dicc_datos_id_pptx = {key: valor for key, valor in global_dicc_datos_id_pptx.items() if key != id_pptx}

        resultado_funcion = def_varios_gui_config_id_pptx("COMBOBOX_LISTA_OPCIONES_ID_PPTX")



    ####################################################################################################################################################
    # UPDATE_PATH_ID_PPTX
    ####################################################################################################################################################
    elif opcion == "UPDATE_PATH_ID_PPTX":
        #actualiza el diccionario asociado al id_pptx (kwargs) en la variable global global_dicc_datos_id_pptx:
        # --> ruta_plantilla_pptx 
        # --> numero_total_slides
        #
        #funciona a su vez como funcion para poder alimentar la GUI y devuelve el path normalizado, el numero de slides del nuevo path
        #devuelve asimismo un mensaje de warning avisando que el cambio de path puede afectar la configuración de rangos de celdas
        #ya guardados con el antiguo path

        tiempo_espera_max_apertura_pptx = global_dicc_datos_id_pptx[id_pptx]["tiempo_espera_max_apertura_pptx"]

        id_pptx_path_ajust = os.path.normpath(id_pptx_path)
        pptx_objeto = Presentation(os.path.normpath(id_pptx_path_ajust))
        time.sleep(tiempo_espera_max_apertura_pptx)

        numero_total_slides = int(float(len([slide_pptx for slide_pptx in pptx_objeto.slides])))

        global_dicc_datos_id_pptx[id_pptx]["ruta_plantilla_pptx"] = id_pptx_path_ajust
        global_dicc_datos_id_pptx[id_pptx]["numero_total_slides"] = numero_total_slides

        mensaje_warning = "El cambio de pptx de destino puede alterar los rangos de celdas configurados con la anterior ubicación.\n\nSe deben resolver estos casos manualmente desde la GUI."

        resultado_funcion = mensaje_warning, id_pptx_path_ajust, numero_total_slides




    ####################################################################################################################################################
    # CHECK_SI_PATH_ID_XLS_YA_ASIGNADO
    ####################################################################################################################################################
    elif opcion == "CHECK_SI_PATH_ID_XLS_YA_ASIGNADO":
        #genera mensaje de warning en caso de que la ruta ya esta asignada a otros id xls para el mismo id pptx
        #usa el parametro kwargs id_xls_excluir_check_path para distinguir los casos:
        # --> el check se hace al crear un nuevo id xls (no se informa id_xls_excluir_check_path)
        # --> el check se hace al atualizar el path de un id xls existente (se informa id_xls_excluir_check_path para excluir del check el id xls que se cambia)
        #
        #a diferencia de la opcion CHECK_SI_PATH_ID_PPTX_YA_ASIGNADO aqui es warning es bloquante en la GUI

        id_xls_path_ajust = os.path.normpath(id_xls_path)

        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            lista_path_id_xls_ya_asignados_sin_exclusion = [dicc["id_xls"] for dicc in lista_dicc_id_xls if os.path.normpath(dicc["ruta_xls_origen"]) == id_xls_path_ajust]

            lista_path_id_xls_ya_asignados_con_exclusion = [dicc["id_xls"] for dicc in lista_dicc_id_xls
                                                            if id_xls_excluir_check_path is not None 
                                                                and dicc["id_xls"] != id_xls_excluir_check_path 
                                                                and os.path.normpath(dicc["ruta_xls_origen"]) == id_xls_path_ajust
                                                            ]

            lista_path_id_xls_ya_asignados = lista_path_id_xls_ya_asignados_sin_exclusion if id_xls_excluir_check_path is None else lista_path_id_xls_ya_asignados_con_exclusion

            str_path_id_xls_ya_asignados = ", ".join(lista_path_id_xls_ya_asignados) if len(lista_path_id_xls_ya_asignados) else None

            resultado_funcion = (f"La ruta que has configurado ya está asignada al id xls '{str_path_id_xls_ya_asignados}' (para el id pptx '{id_pptx}').\n\nNo se puede configurar más de una vez el mismo excel de origen para un mismo id pptx."
                                 if str_path_id_xls_ya_asignados is not None else None)



    ###################################################################################################################################################
    # ASIGNAR_ID_XLS_NUEVO
    ###################################################################################################################################################
    elif opcion == "ASIGNAR_ID_XLS_NUEVO":
        #devuelve un id pptx nuevo correlativo con el ultimo creado
        #se usa una funcion lambda que normaliza el numero con 3 digitos
        #en el caso de que se hayan eliminado id xls entre el 1ero y ultimo, el id xls que se asigna es el 1er hueco que encuentra desde el 1ero

        lambda_id_xls = lambda cont: f"{cont:03}"
        
        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]
        lista_dicc_id_xls_config = [dicc["id_xls"] for dicc in lista_dicc_id_xls] if len(lista_dicc_id_xls) != 0 else []
        lista_dicc_id_xls_config = sorted(lista_dicc_id_xls_config)


        lista_id_xls_config_huecos_por_cubrir = [f"{nombre_id_xls}_{lambda_id_xls(ind + 1)}" 
                                                for ind, id_xls_config in enumerate(lista_dicc_id_xls_config)
                                                if ind + 1 != int(id_xls_config[-3:])
                                                and f"{nombre_id_xls}_{lambda_id_xls(ind + 1)}" not in lista_dicc_id_xls_config
                                                ]
        
        
        indice_correlativo = len(lista_dicc_id_xls_config) + 1
        nuevo_id_xls = f"{nombre_id_xls}_{lambda_id_xls(indice_correlativo)}" if len(lista_id_xls_config_huecos_por_cubrir) == 0 else lista_id_xls_config_huecos_por_cubrir[0]


        resultado_funcion = nuevo_id_xls



    ###################################################################################################################################################
    # CREAR_NUEVO_ID_XLS
    ###################################################################################################################################################
    elif opcion == "CREAR_NUEVO_ID_XLS":
        #siempre y cuando el id_xls_path (kwargs) no este asignado ya a otro id_xls (para el mismo id_pptx), modifica el diccionario de la key id_pptx (kwargs) de la variable global global_dicc_datos_id_pptx:
        # --> df_widget_treeview_id_xls (recalcula el dataframe agregando el id_xls que se crea)
        #
        #siempre y cuando el id_xls_path (kwargs) no este asignado ya a otro id_xls (para el mismo id_pptx), agrega un diccionario en la lista almacenada en la key lista_dicc_id_xls:
        # --> id_xls (kwargs id_xls)
        # --> desc_id_xls (kwargs id_xls_desc)
        # --> ruta_xls_origen (kwargs id_xls_path)
        # --> tiempo_espera_max_apertura_xls (kwargs id_xls_tiempo_espera_max_apertura)
        # --> actualizar_vinculos_otros_xls (kwargs id_xls_actualizar_vinculos_otros_xls)
        # --> lista_hojas_xls (lista que se calcula en el proceso)
        # --> df_widget_treeview_rangos_celdas (dataframe vacio con 6 columnas: ID_XLS, HOJA_XLS, RANGO_XLS, SLIDE_PPTX, NOMBRE_SCREENSHOT, COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX y SCREENSHOT_PNG)
        #
        #funciona a su vez como funcion para poder alimentar la GUI y devuelve un df para poder actualizar el treeview de la GUI
        #devuelve tambien la lista lista_datos_items_seleccionados (sirve para seleccionar el item en el treview y para informar el atributo datos_items del objeto treeview)


        id_xls_path_ajust = os.path.normpath(id_xls_path)


        #se convierte id_xls_tiempo_espera_max_apertura en entero
        id_xls_tiempo_espera_max_apertura_ajust = int(float(id_xls_tiempo_espera_max_apertura))


        #se crea el df actualizado con el nuevo id xls (df_widget_treeview_id_xls_con_modif)
        #se crean lista_datos_items_seleccionados que sirve para el atributo datos_items del objeto treeview en la GUI
        #se actualiza la key df_widget_treeview_id_xls en global_dicc_datos_id_pptx

        df_widget_treeview_id_xls_anter = global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"]

        lista_para_df_widget_treeview_id_xls_nuevo = [[id_xls, id_xls_desc]]
        df_widget_treeview_id_xls_nuevo = pd.DataFrame(lista_para_df_widget_treeview_id_xls_nuevo, columns = ["ID_XLS", "DESC_ID_XLS"])

        df_widget_treeview_id_xls_con_modif = pd.concat([df_widget_treeview_id_xls_anter, df_widget_treeview_id_xls_nuevo])
        df_widget_treeview_id_xls_con_modif = df_widget_treeview_id_xls_con_modif[[col for col in df_widget_treeview_id_xls_con_modif.columns]].sort_values(["ID_XLS"], ascending = [True])
        df_widget_treeview_id_xls_con_modif.reset_index(drop = True, inplace = True)

        lista_datos_items_seleccionados = [[id_xls, id_xls_desc]]



        global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"] = df_widget_treeview_id_xls_con_modif





        #se agrega un diccionario a la key lista_dicc_id_xls (es lista por defecto) de global_dicc_datos_id_pptxdonde previamente se calcula:
        # --> lista_hojas_xls                          lista de hojas que contiene el excel
        # --> df_widget_treeview_rangos_celdas         df de rangos de celdas vacio

        lista_hojas_xls = def_varios("LISTA_HOJAS_XLS_DESDE_PATH_XLS", id_xls_path = id_xls_path_ajust)

        lista_campos_para_df_rangos_celdas = global_dicc_tablas_y_campos_sistema["PARAMETRICA"]["DICC_LISTA_CAMPOS"]["TREEVIEW_RANGOS_CELDAS"]
        df_widget_treeview_rangos_celdas = pd.DataFrame(columns = lista_campos_para_df_rangos_celdas)

        dicc_id_xls = {"id_xls": id_xls
                        , "desc_id_xls": id_xls_desc
                        , "ruta_xls_origen": id_xls_path_ajust
                        , "tiempo_espera_max_apertura_xls": id_xls_tiempo_espera_max_apertura_ajust
                        , "actualizar_vinculos_otros_xls": id_xls_actualizar_vinculos_otros_xls
                        , "lista_hojas_xls": lista_hojas_xls
                        , "df_widget_treeview_rangos_celdas": df_widget_treeview_rangos_celdas
                        }
        
        global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"].append(dicc_id_xls)
        

        #resultado de la funcion
        resultado_funcion = lista_datos_items_seleccionados, df_widget_treeview_id_xls_con_modif



    ###################################################################################################################################################
    # ELIMINAR_ID_XLS
    ###################################################################################################################################################
    elif opcion == "ELIMINAR_ID_XLS":
        #modifica el diccionario almacenado en la key id_pptx en la variable global global_dicc_datos_id_pptx
        # --> df_widget_treeview_id_xls (se recalcula el dataframe quitando el id_xls que se elimina)
        # --> lista_dicc_id_xls (se quita el diccionario asociado al id_xls que se elimina)
        #
        #funciona a su vez como funcion y genera el dataframe por id xls actualizado

        
        df_widget_treeview_id_xls = global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"]
        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        df_widget_treeview_id_xls_ajust = df_widget_treeview_id_xls.loc[df_widget_treeview_id_xls["ID_XLS"] != id_xls, [col for col in df_widget_treeview_id_xls.columns]]
        df_widget_treeview_id_xls_ajust.reset_index(drop = True, inplace = True)

        lista_dicc_id_xls_ajust = [dicc for dicc in lista_dicc_id_xls if dicc["id_xls"] != id_xls]


        global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"] = df_widget_treeview_id_xls_ajust
        global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"] = lista_dicc_id_xls_ajust

        resultado_funcion = df_widget_treeview_id_xls_ajust



    ###################################################################################################################################################
    # UPDATE_PATH_ID_XLS
    ###################################################################################################################################################
    elif opcion == "UPDATE_PATH_ID_XLS":
        #actualiza la key lista_dicc_id_xls del diccionario asociado al id_pptx (kwargs) en la variable global global_dicc_datos_id_pptx
        #se modifican las keys en el diccionario de esta lista asociado al id_xls (kwargs):
        # --> ruta_xls_origen
        # --> lista_hojas_xls
        #
        #funciona a su vez como funcion para poder alimentar la GUI
        #devuelve un mensaje warning (no bloquante) para informar en la GUI de que las configuraciones
        #de rangos de celdas xls asociadas al id xls pueden quedar con hojas xls desactualizadas
        #(se muestran en el mensaje que hojas ya no estan en el excel actualizado y que tienen rangos de celdas configurados)
        #devuelve tambien id_xls_path_ajust y la lista de hojas del excel

        id_xls_path_ajust = os.path.normpath(id_xls_path)

        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            for ind, dicc in enumerate(lista_dicc_id_xls):
                if dicc["id_xls"] == id_xls:

                    lista_hojas_xls = def_varios("LISTA_HOJAS_XLS_DESDE_PATH_XLS", id_xls_path = id_xls_path_ajust)

                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["ruta_xls_origen"] = id_xls_path_ajust
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["lista_hojas_xls"] = lista_hojas_xls

                    #resultado de la funcion
                    mensaje_warning = f"Las configuraciones de rangos de celdas asociadas al id_xls ('{id_xls}' pueden quedar desactualizadas con la hoja excel. Conviene revisarlas y modificarlas."
                    resultado_funcion = mensaje_warning, id_xls_path_ajust, lista_hojas_xls

                    break



    #############################################################################################################
    # UPDATE_EN_MEMORIA_DATOS_GUI_ANTES_HACER_CLICK_EN_OTRO_ID_XLS
    #############################################################################################################
    elif opcion == "UPDATE_EN_MEMORIA_DATOS_GUI_ANTES_HACER_CLICK_EN_OTRO_ID_XLS":
        #permite guardar los cambios realizados en la GUI de configuracion de id pptx cuando el usuario hace click en otro id xls
        #si el usuario ha modificado datos del id xls anterior al click en el item estos cambios se pierden
        #los cambios se guardan en la variable global global_dicc_datos_id_pptx

        #se guardan los datos del id xls seleccionado (si lo esta) en el treeview asociado al id pptx en pantalla
        #antes de actualizar el combobox de seleccion de id pptx y pulsar el boton VER
        #
        #funciona a su vez como funcion y devuelve un dataframe para poder actualizar el treeview por id xls y la lista (lista de listas)
        #lista_datos_items_seleccionados para poder informar el atributo datos_items del objeto treeview en la GUI
        lista_datos_items_seleccionados = None
        df_widget_treeview_id_xls_despues = None
        if id_xls is not None:


            df_widget_treeview_id_xls_antes = global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"]
            lista_df_widget_treeview_id_xls_antes = df_widget_treeview_id_xls_antes.values.tolist()


            lista_df_widget_treeview_id_xls_despues = []
            for id_xls_iter, id_xls_desc_iter in lista_df_widget_treeview_id_xls_antes:

                if id_xls_iter == id_xls:
                    lista_df_widget_treeview_id_xls_despues.append([id_xls, id_xls_desc])
                else:
                    lista_df_widget_treeview_id_xls_despues.append([id_xls_iter, id_xls_desc_iter])

            lista_col_df = [col for col in df_widget_treeview_id_xls_antes.columns]

            df_widget_treeview_id_xls_despues = pd.DataFrame(lista_df_widget_treeview_id_xls_despues, columns = lista_col_df)
            lista_datos_items_seleccionados = [[id_xls, id_xls_desc]]


            global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"] = df_widget_treeview_id_xls_despues


            for ind, dicc in enumerate(global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]):
                if dicc["id_xls"] == id_xls:


                    global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"] = df_widget_treeview_id_xls_despues
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["desc_id_xls"] = id_xls_desc
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["ruta_xls_origen"] = id_xls_path
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["tiempo_espera_max_apertura_xls"] = id_xls_tiempo_espera_max_apertura
                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["actualizar_vinculos_otros_xls"] = id_xls_actualizar_vinculos_otros_xls

                    break

        resultado_funcion = lista_datos_items_seleccionados, df_widget_treeview_id_xls_despues




        # --> id_pptx_desc                            descripcion del id pptx
        # --> ruta_plantilla_pptx                     ruta donde se ubica el fichero pptx en el cual se realizan als acciones pptx
        # --> tiempo_espera_max_apertura_pptx         tiempo de espera maximo que se espera para abrit el fichero pptx
        # --> numero_total_slides                     numero de slides que contiene el pptx
        # --> df_widget_treeview_id_xls               df con los id xls asociados al id pptx seleccionado (permite informar el treeview de la clase gui_config_id_pptx)
        #                                             (el df tiene las columnas siguientes ID_XLS y DESC_ID_XLS)
        #
        # --> lista_dicc_id_xls                       lista de diccionarios (cada diccionario corresponde a un id xls asociado al id pptx)
        #                                             cada diccionario contiene las keys:
        #                                                     --> id_xls                                     id xls asociado al id pptx
        #                                                     --> desc_id_xls                                descripcion del id xls
        #                                                     --> ruta_xls_origen                            ruta donde se ubica el fichero excel de origen correspondiente al id xls
        #                                                     --> tiempo_espera_max_apertura_xls             tiempo de espera maximo que se espera para abrit el fichero excel de origen
        #                                                     --> actualizar_vinculos_otros_xls              actualizar los vinculos hacia otros excels externos al abrir el fichero excel de origen
        #
        #                                                     --> lista_hojas_xls                            lista con las hojas disponibles en el excel asociado al id xls
        #                                                                                                    (se usa para actualizar el combobox de hojas xls en el frame de rangos de celdas)
        #
        #                                                     --> df_widget_treeview_rangos_celdas           df que sirve para informar el treeview por rangos de celda nada mas seleccionar un id xls en el treeview por id xls
        #                                                                                                    en la clase gui_config_id_pptx (sirve tambien para informar el treeview mencionado cuando se filtra el combobox 
        #                                                                                                    por la opcion TODOS para el treview por rangos de celda)
        #                                                                                                    (el df se crea con las columnas ID_XLS, HOJA_XLS, RANGO_XLS, SLIDE_PPTX, NOMBRE_SCREENSHOT,
        #                                                                                                     COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX y SCREENSHOT_PNG)












    ###################################################################################################################################################
    # CLICK_ITEM_TREEVIEW_ID_XLS
    ###################################################################################################################################################
    elif opcion == "CLICK_ITEM_TREEVIEW_ID_XLS":
        #funciona como funcion para poder alimentar la GUI cuando se se cliqua en un item del treeview por id xls
        #devuelve los valores:
        # --> id_xls_path
        # --> id_xls_desc
        # --> id_xls_tiempo_espera_max_apertura
        # --> id_xls_actualizar_vinculos_otros_xls
        # --> df_widget_treeview_rangos_celdas

        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            for dicc in lista_dicc_id_xls:
                if dicc["id_xls"] == id_xls:

                    id_xls_path = dicc["ruta_xls_origen"]
                    id_xls_desc = dicc["desc_id_xls"]
                    id_xls_tiempo_espera_max_apertura = dicc["tiempo_espera_max_apertura_xls"]
                    id_xls_actualizar_vinculos_otros_xls = dicc["actualizar_vinculos_otros_xls"]
                    df_widget_treeview_rangos_celdas = dicc["df_widget_treeview_rangos_celdas"]

                    id_xls_tiempo_espera_max_apertura_ajust = int(float(id_xls_tiempo_espera_max_apertura))

                    #resultado de la funcion
                    resultado_funcion = id_xls_path, id_xls_desc, id_xls_tiempo_espera_max_apertura_ajust, id_xls_actualizar_vinculos_otros_xls, df_widget_treeview_rangos_celdas

                    break


    ###################################################################################################################################################
    # DF_TREEVIEW_ID_XLS_TRAS_CAMBIOS_EN_DESCRIPCION
    ###################################################################################################################################################
    elif opcion == "DF_TREEVIEW_ID_XLS_TRAS_CAMBIOS_EN_DESCRIPCION":
        #en caso de que el usuario cambie en la GUI la descripcion de un id xls y luego clica en otro item del treeview de id xls
        #y luego vuelve al anterior la descripcion en el treeview no se actualiza pero en el widget de la descripcion si aparece el cambio
        #se recalcula el df para informar el treeviiew con la descripción actualizada
        #
        #funciona como funcion y devuelve el df actualizado para poder informar el treeview en la GUI

        id_xls_por_actualiz = lista_datos_items_selecc_update_desc[0][0]
        id_xls_desc_por_actualiz = lista_datos_items_selecc_update_desc[0][1]

        df_widget_treeview_id_xls = global_dicc_datos_id_pptx[id_pptx]["df_widget_treeview_id_xls"]
        lista_campos_df = [col for col in df_widget_treeview_id_xls.columns]

        lista_datos_widget_treeview_id_xls = df_widget_treeview_id_xls.values.tolist()

        if len(lista_datos_widget_treeview_id_xls) != 0:

            lista_datos_para_df_update = [[id_xls_por_actualiz, id_xls_desc_por_actualiz]
                                            if id_xls_por_actualiz == id_xls_iter
                                            else
                                            [id_xls_iter, id_xls_desc_iter]
                                            for id_xls_iter, id_xls_desc_iter in lista_datos_widget_treeview_id_xls
                                            ]
            

            resultado_funcion = pd.DataFrame(lista_datos_para_df_update, columns = lista_campos_df)
            




    ###################################################################################################################################################
    # CHECK_SI_RANGO_CELDAS_YA_ASIGNADO
    ###################################################################################################################################################
    elif opcion == "CHECK_SI_RANGO_CELDAS_YA_ASIGNADO":
        #funciona como funcion y permite chequear si el rango de celdas ya esta configurado

        slide_pptx_ajust = int(float(slide_pptx))

        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            for dicc in lista_dicc_id_xls:
                if dicc["id_xls"] == id_xls:

                    df_widget_treeview_rangos_celdas = dicc["df_widget_treeview_rangos_celdas"]

                    df_check = (df_widget_treeview_rangos_celdas.loc[(df_widget_treeview_rangos_celdas["HOJA_XLS"] == hoja_xls)
                                                                    & (df_widget_treeview_rangos_celdas["RANGO_XLS"] == rango_celdas)
                                                                    & (df_widget_treeview_rangos_celdas["SLIDE_PPTX"] == slide_pptx_ajust)
                                                                    , [col for col in df_widget_treeview_rangos_celdas.columns]
                                                                    ])

                    #resultado de la funcion
                    resultado_funcion = (f"El rango de celdas (misma hoja excel, mismo rango de celdas y misma slide pptx de destino) que intentas agregar ya está configurado en memoria para el id xls ({id_xls})"
                                        if len(df_check) != 0 else None)

                    break



    ###################################################################################################################################################
    # AGREGAR_RANGO_CELDAS
    ###################################################################################################################################################
    elif opcion == "AGREGAR_RANGO_CELDAS":
        #actualiza la key lista_dicc_id_xls del diccionario asociado al id_pptx (kwargs) en la variable global global_dicc_datos_id_pptx.
        #se modifica el diccionario de esta lista asociado al id_xls (kwargs):
        # --> df_widget_treeview_rangos_celdas (se agrega el nuevo rango de celdas en el dataframe siempre y cuando no exista previamente)
        #
        #funciona a su vez como función para poder alimentar la GUI y devuelve un dataframe para poder informar la GUI con los rangos de celdas actualizados
        #devuelve tambien la lista lista_datos_items_seleccionados
        #(sirve para seleccionar el item en el treview y para informar el atributo datos_items del objeto treeview)


        #se convierte el kwargs slide_pptx en entero
        slide_pptx_ajust = int(float(slide_pptx))

        #se crea la listas de lista lista_rangos_celdas_despues_update_selecc_treeview (tiene 1 sola sublista)
        #para poder seleccionar en la GUI el item actualizado
        lista_datos_items_seleccionados = [[hoja_xls, rango_celdas, slide_pptx_ajust]]


        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            for ind, dicc in enumerate(lista_dicc_id_xls):
                if dicc["id_xls"] == id_xls:

                    df_widget_treeview_rangos_celdas = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["df_widget_treeview_rangos_celdas"]


                    lista_para_df_rangos_celdas = [[id_xls
                                                    , hoja_xls
                                                    , rango_celdas
                                                    , slide_pptx_ajust
                                                    , None #NOMBRE_SCREENSHOT (aqui es None - se calcula al guardar en el sistemasqlite)
                                                    , None #COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX (aqui es None - se calcula cuando el susurio configura Paso 1 el pptx)
                                                    , None #SCREENSHOT_PNG (aqui es None - se calcula cuando el susurio configura Paso 1 el pptx)
                                                    ]]

                    lista_columnas_df = [col for col in df_widget_treeview_rangos_celdas.columns]

                    df_widget_treeview_rangos_celdas_nuevo_rango_celdas = pd.DataFrame(lista_para_df_rangos_celdas, columns = lista_columnas_df)

                    df_widget_treeview_rangos_celdas_nuevo = pd.concat([df_widget_treeview_rangos_celdas, df_widget_treeview_rangos_celdas_nuevo_rango_celdas])
                    df_widget_treeview_rangos_celdas_nuevo = df_widget_treeview_rangos_celdas_nuevo[lista_columnas_df].sort_values(["HOJA_XLS", "SLIDE_PPTX"], ascending = [True, True])
                    df_widget_treeview_rangos_celdas_nuevo.reset_index(drop = True, inplace = True)



                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["df_widget_treeview_rangos_celdas"] = df_widget_treeview_rangos_celdas_nuevo




                    #resultado de la funcion
                    resultado_funcion = lista_datos_items_seleccionados, df_widget_treeview_rangos_celdas_nuevo

                    break




    ###################################################################################################################################################
    # ELIMINAR_RANGO_CELDAS
    ###################################################################################################################################################
    elif opcion == "ELIMINAR_RANGO_CELDAS":
        #actualiza la key lista_dicc_id_xls del diccionario asociado al id_pptx (kwargs) en la variable global global_dicc_datos_id_pptx.
        #se modifica el diccionario de esta lista asociado al id_xls (kwargs):
        # --> df_widget_treeview_rangos_celdas (se quita el rango de celdas en el dataframe)
        #
        #funciona a su vez como funcion para poder alimentar la GUI y devuelve:
        # --> df_widget_treeview_rangos_celdas_nuevo

        #se convierte el kwargs slide_pptx en entero
        slide_pptx_ajust = int(float(slide_pptx))


        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            for ind, dicc in enumerate(lista_dicc_id_xls):
                if dicc["id_xls"] == id_xls:

                    df_widget_treeview_rangos_celdas = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["df_widget_treeview_rangos_celdas"]
                    lista_columnas_df = [col for col in df_widget_treeview_rangos_celdas.columns]


                    lista_datos_df_widget_treeview_rangos_celdas = df_widget_treeview_rangos_celdas.values.tolist()

                    lista_datos_df_widget_treeview_rangos_celdas_filtrado = [[id_xls_iter
                                                                              , hoja_xls_iter
                                                                              , rango_celdas_iter
                                                                              , slide_pptx_iter
                                                                              , nombre_screenshot_iter
                                                                              , screenshot_png_iter
                                                                              , coordenadas_screenshot_iter
                                                                              ]
                                                                            for id_xls_iter, hoja_xls_iter, rango_celdas_iter, slide_pptx_iter, nombre_screenshot_iter, screenshot_png_iter, coordenadas_screenshot_iter in lista_datos_df_widget_treeview_rangos_celdas
                                                                            if f"{hoja_xls_iter}_{rango_celdas_iter}_{int(float(slide_pptx_iter))}" != f"{hoja_xls}_{rango_celdas}_{slide_pptx_ajust}"
                                                                            ]
                    
                    df_widget_treeview_rangos_celdas_filtrado = pd.DataFrame(lista_datos_df_widget_treeview_rangos_celdas_filtrado, columns = lista_columnas_df)
                    df_widget_treeview_rangos_celdas_filtrado = df_widget_treeview_rangos_celdas_filtrado[lista_columnas_df].sort_values(["HOJA_XLS", "SLIDE_PPTX"], ascending = [True, True])
                    df_widget_treeview_rangos_celdas_filtrado.reset_index(drop = True, inplace = True)

                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["df_widget_treeview_rangos_celdas"] = df_widget_treeview_rangos_celdas_filtrado


                    #resultado de la funcion
                    resultado_funcion = df_widget_treeview_rangos_celdas_filtrado

                    break



    ###################################################################################################################################################
    # UPDATE_RANGO_CELDAS
    ###################################################################################################################################################
    elif opcion == "UPDATE_RANGO_CELDAS":
        #actualiza la variable global global_dicc_datos_id_pptx con el df df_widget_treeview_rangos_celdas asociado al id_pptx y id_xls pasados
        #con el rango de celdas actualizado
        #
        #funciona a su vez como funcion para poder alimentar la GUI y devuelve:
        # --> lista_rangos_celdas_despues_update_selecc_treeview
        # --> df_widget_treeview_rangos_celdas_nuevo


        hoja_xls_antes_update = lista_rangos_celdas_antes_update[0]
        rango_celdas_antes_update = lista_rangos_celdas_antes_update[1]
        slide_pptx_antes_update = int(float(lista_rangos_celdas_antes_update[2])) #se convierte en entero

        concat_items_rangos_celdas = f"{hoja_xls_antes_update}_{rango_celdas_antes_update}_{slide_pptx_antes_update}"

        hoja_xls_despues_update = lista_rangos_celdas_despues_update[0]
        rango_celdas_despues_update = lista_rangos_celdas_despues_update[1]
        slide_pptx_despues_update = int(float(lista_rangos_celdas_despues_update[2]))#se convierte en entero

        #se crea la listas de lista lista_rangos_celdas_despues_update_selecc_treeview (tiene 1 sola sublista)
        #para poder seleccionar en la GUI el item actualizado
        lista_rangos_celdas_despues_update_selecc_treeview = [[hoja_xls_despues_update, rango_celdas_despues_update, slide_pptx_despues_update]]

        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            for ind, dicc in enumerate(lista_dicc_id_xls):

                if dicc["id_xls"] == id_xls:

                    df_widget_treeview_rangos_celdas = dicc["df_widget_treeview_rangos_celdas"]
                    lista_columnas_df = [col for col in df_widget_treeview_rangos_celdas.columns]

                    lista_datos_antes_update = df_widget_treeview_rangos_celdas.values.tolist()

                    lista_datos_despues_update = [[id_xls_iter   #no varia
                                                   , hoja_xls_despues_update if f"{hoja_xls_iter}_{rango_xls_iter}_{slide_pptx_iter}" == concat_items_rangos_celdas else hoja_xls_iter
                                                   , rango_celdas_despues_update if f"{hoja_xls_iter}_{rango_xls_iter}_{slide_pptx_iter}" == concat_items_rangos_celdas else rango_xls_iter
                                                   , slide_pptx_despues_update if f"{hoja_xls_iter}_{rango_xls_iter}_{slide_pptx_iter}" == concat_items_rangos_celdas else slide_pptx_iter
                                                   , nombre_screenshot_iter   #no varia
                                                   , screenshot_png_iter      #no varia
                                                   , coordenadas_screenshot   #no varia
                                                   ]
                                                for id_xls_iter, hoja_xls_iter, rango_xls_iter, slide_pptx_iter, nombre_screenshot_iter, screenshot_png_iter, coordenadas_screenshot in lista_datos_antes_update
                                                ]
                    
                    df_widget_treeview_rangos_celdas_update = pd.DataFrame(lista_datos_despues_update, columns = lista_columnas_df)
                    df_widget_treeview_rangos_celdas_update = df_widget_treeview_rangos_celdas_update[lista_columnas_df].sort_values(["HOJA_XLS", "SLIDE_PPTX", "RANGO_XLS"], ascending = [True, True, True])
                    df_widget_treeview_rangos_celdas_update.reset_index(drop = True, inplace = True)

                    global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"][ind]["df_widget_treeview_rangos_celdas"] = df_widget_treeview_rangos_celdas_update

                    resultado_funcion = lista_rangos_celdas_despues_update_selecc_treeview, df_widget_treeview_rangos_celdas_update


                    break



    ###################################################################################################################################################
    # CLICK_ITEM_TREEVIEW_RANGOS_CELDAS
    ###################################################################################################################################################
    elif opcion == "CLICK_ITEM_TREEVIEW_RANGOS_CELDAS":
        #funciona como funcion para poder alimentar la GUI cuando se se cliqua en un item del treeview por rangos de celdas
        #devuelve los valores:
        # --> nombre_screenshot
        # --> coordenadas_screenshot

        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        slide_pptx_ajust = int(float(slide_pptx))

        if len(lista_dicc_id_xls) != 0:

            for dicc in lista_dicc_id_xls:
                if dicc["id_xls"] == id_xls:

                    df_widget_treeview_rangos_celdas = dicc["df_widget_treeview_rangos_celdas"]

                    df_temp = (df_widget_treeview_rangos_celdas.loc[(df_widget_treeview_rangos_celdas["HOJA_XLS"] == hoja_xls)
                                                                   & (df_widget_treeview_rangos_celdas["RANGO_XLS"] == rango_celdas)
                                                                   & (df_widget_treeview_rangos_celdas["SLIDE_PPTX"] == slide_pptx_ajust)
                                                                   , ["NOMBRE_SCREENSHOT", "COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX"]]
                                )
                    
                    df_temp.reset_index(drop = True, inplace = True)

                    nombre_screenshot = df_temp.iloc[0, 0]
                    coordenadas_screenshot = df_temp.iloc[0, 1]

                    #resultado de la funcion
                    resultado_funcion = nombre_screenshot, coordenadas_screenshot

                    break





    #############################################################################################################
    # SCREENSHOT_PNG_MUESTRA_PATH
    #############################################################################################################
    elif opcion == "SCREENSHOT_PNG_MUESTRA_PATH":
        #devuelve el path del screenshot si existe en la ruta guardada en la ruta almacenada en caso de que el path no existe devuevle None
        #se usa en la GUI de config id pptx para generar nueva ventana con la muestra del screenshot asociado al rango de celdas seleccionado en pantalla

        if nombre_screenshot is not None:
            ruta_descarga_screenshot_en_png = os.path.join(global_ruta_local_screenshots_png_muestra, f"{nombre_screenshot}.PNG")
            ruta_descarga_screenshot_en_png = os.path.normpath(ruta_descarga_screenshot_en_png)

            ruta_descarga_screenshot_en_png = ruta_descarga_screenshot_en_png if os.path.exists(ruta_descarga_screenshot_en_png) else None

            resultado_funcion = (ruta_descarga_screenshot_en_png 
                                if ruta_descarga_screenshot_en_png is not None and os.path.exists(ruta_descarga_screenshot_en_png)
                                else None)




    ####################################################################################################################################################
    # ABRIR_FICHERO
    ####################################################################################################################################################
    elif opcion == "ABRIR_FICHERO":
        #abre el fichero pptx o excel segun el caso
        #funciona a su vez como funcion y genera un mensaje warning en caso de que no se pueda abrir el fichero

        ruta_fichero_ajust = os.path.normpath(ruta_fichero)

        if os.path.exists(ruta_fichero_ajust):
            os.startfile(ruta_fichero_ajust)

        else:
            resultado_funcion = f"La ruta indicada no existe o no tienes acceso a la misma:\n\n{ruta_fichero_ajust}"



    ####################################################################################################################################################
    # COMBOBOX_LISTA_OPCIONES_ID_PPTX
    ####################################################################################################################################################
    elif opcion == "COMBOBOX_LISTA_OPCIONES_ID_PPTX":
        #devuelve una lista de id pptx para usarse como opciones en el combobox de la ventana para configurar id pptx

        resultado_funcion = sorted(list(global_dicc_datos_id_pptx.keys()))



    ####################################################################################################################################################
    # COMBOBOX_LISTA_OPCIONES_HOJAS_XLS
    ####################################################################################################################################################
    elif opcion == "COMBOBOX_LISTA_OPCIONES_HOJAS_XLS":
        #devuelve una lista con todas las hojas excel del idxls asociado al id_pptx pasados por parametros de la funcion
        #se usa para crear la lista de opciones del combobox de hojasxls (configuracion id pptx dentro del frame rangos celdas xls)

        lista_dicc_id_xls = global_dicc_datos_id_pptx[id_pptx]["lista_dicc_id_xls"]

        if len(lista_dicc_id_xls) != 0:

            for dicc in lista_dicc_id_xls:
                if dicc["id_xls"] == id_xls:

                    resultado_funcion = sorted(dicc["lista_hojas_xls"])

                    break


    ####################################################################################################################################################
    # COMBOBOX_LISTA_OPCIONES_SLIDES_PPTX
    ####################################################################################################################################################
    elif opcion == "COMBOBOX_LISTA_OPCIONES_SLIDES_PPTX":
        #devuelve una lista con todas las slides del id_pptx pasado por parametro de la funcion
        #se usa para crear la lista de opciones del combobox de slides pptx (configuracion id pptx dentro del frame rangos celdas xls)

        numero_total_slides = global_dicc_datos_id_pptx[id_pptx].get("numero_total_slides", None)
        resultado_funcion = [slide_pptx for slide_pptx in range(1, int(float(numero_total_slides)) + 1, 1)] if numero_total_slides is not None else []


    #resultado funcion
    return resultado_funcion



#################################################################################################################################
#################################################################################################################################
#################################################################################################################################
# RUTINAS / FUNCIONES - GUI VENTANA INICIO - ACCIONES PPTX
#################################################################################################################################
#################################################################################################################################
#################################################################################################################################


def def_acciones_screenshots_rangos_celdas_xls(opcion_acciones_screenshots, **kwargs):
    #rutina que ejecuta acciones poppler referente a la captura y guardado de screenshots de rangos de celdas excel a ficheros png
    #hay 2 opciones:
    # --> RECORTAR_BLANCOS_SCREENSHOT_PNG         permite quitar todo lo "blanco" en el png generado con Poppler usando la libreria open-cv
    #                                             (el png viene en formato pagina con el screenshot arriba a la izquierda
    #                                             y el resto de la pagina en blanco)
    #
    # --> SCREENSHOTS_XLS_A_PNG                   permite abrir los excels, exportar rangos a pdf, convertir a png y almacenar 
    #                                             los posibles errores en una variable global (se usa el metodo convert_from_path
    #                                             de la libreria pdf2image que funciona con el ejecutable de Poppler)


    global global_poppler_path
    global lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx
    global global_dicc_check_screenshots_si_todo_blanco
    global global_lista_dicc_errores



    #parametros kwargs
    id_pptx_selecc = kwargs.get("id_pptx_selecc", None)
    id_xls = kwargs.get("id_xls", None)
    hoja_xls = kwargs.get("hoja_xls", None)
    rango_celdas_xls = kwargs.get("rango_celdas_xls", None)
    slide_pptx = kwargs.get("slide_pptx", None)
    path_screenshot_png = kwargs.get("path_screenshot_png", None)
    nombre_screenshot = kwargs.get("nombre_screenshot", None)
    path_xls = kwargs.get("path_xls", None)
    tiempo_espera_max_apertura_excel = kwargs.get("tiempo_espera_max_apertura_excel", None)
    lista_hojas_con_rangos_celdas = kwargs.get("lista_hojas_con_rangos_celdas", None)
    directorio_screenshots = kwargs.get("directorio_screenshots", None)
    opcion_update_links = kwargs.get("opcion_update_links", False) #aqui es False
    opcion_interaccion_pptx = kwargs.get("opcion_interaccion_pptx", None)



    ##################################################################################################
    #    RECORTAR_BLANCOS_SCREENSHOT_PNG
    ##################################################################################################
    if opcion_acciones_screenshots == "RECORTAR_BLANCOS_SCREENSHOT_PNG":

        try:

            img = cv2.imread(path_screenshot_png)

            if img is None:
                return


            color_gris = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)


            #se convierte el png en blanco y negro, lo que queda en negro es el contenido real del screenshot
            #--> pixeles > opencv_recortar_png_casi_blancos --> se vuelven negro
            #--> pixeles <= opencv_recortar_png_casi_blancos --> se vuelven blanco
            _, thresh = cv2.threshold(color_gris, opencv_recortar_png_casi_blancos, 255, cv2.THRESH_BINARY_INV)

            #se localizan las coordenadas de todo lo que no es fondo blanco
            coordenadas = cv2.findNonZero(thresh)

            #por si todo es fondo blanco
            if coordenadas is None:
                return

            #se calcula el rectangulo mínimo que encierra todo el contenido negro
            # --> x = coordenadas de iquierda a derecha
            # --> y = coordenadas de arriba para abajo
            # --> width = ancho del rectangulo
            # --> height = alto del rectangulo
            x, y, width, height = cv2.boundingRect(coordenadas)
            rectangulo_screenshot_recortado = img[y:y + height, x:x + width]

            #se sobrescribe el png original por uno con el rectangulo del screenshot (ya con colores)
            cv2.imwrite(path_screenshot_png, rectangulo_screenshot_recortado)



            #se actualiza la variable global global_dicc_check_screenshots_si_todo_blanco
            #sus keys son los nombres de los screenshots de los rangos de celdas configurados para el id pptx
            #sus valores son dicionarios con las keys siguientes:
            # --> id_xls            
            # --> hoja_xls
            # --> rango_celdas
            # --> slide_pptx
            # --> todo_en_blanco       llegado a esta linea los screenshots son los que no son todo en blanco porlo que se informa a false
            global_dicc_check_screenshots_si_todo_blanco[nombre_screenshot]["todo_en_blanco"] = False





        except Exception as err_1:
            #se registra el posible error de recortes de "blancos" en el png
            #se pone "pass" para que se registren todos los errores no solo el 1ero con el que se topa

            traceback_error = traceback.extract_tb(err_1.__traceback__)
            modulo_python = os.path.basename(traceback_error[0].filename)
            rutina_python = traceback_error[0].name
            linea_error = traceback_error[0].lineno

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            dicc_error["PROCESO"] = opcion_interaccion_pptx
            dicc_error["FASE_PROCESO"] = "Recorte espacios 'blancos' del fichero png asociado al screenshot de rango de celdas."
            dicc_error["ID_PPTX"] = id_pptx_selecc
            dicc_error["ID_XLS"] = id_xls
            dicc_error["RUTA_FICHERO"] = path_xls
            dicc_error["HOJA_XLS"] = hoja_xls
            dicc_error["RANGO_CELDAS"] = rango_celdas_xls
            dicc_error["SLIDE_PPTX"] = slide_pptx
            dicc_error["RESUMEN_ERROR"] = "Fallo en el recorte de los espacios 'blancos' en el png resultante del rango de celdas."
            dicc_error["MODULO_PYTHON"] = modulo_python
            dicc_error["RUTINA_PYTHON"] = rutina_python
            dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
            dicc_error["TRACEBACK"] = str(err_1)

            global_lista_dicc_errores.append(dicc_error)

            pass



    ##################################################################################################
    #    SCREENSHOTS_XLS_A_PNG
    ##################################################################################################
    elif opcion_acciones_screenshots == "SCREENSHOTS_XLS_A_PNG":

        update_links_xls = False if opcion_interaccion_pptx == "CONFIG_PASO_1" else opcion_update_links

        lista_hojas_xls_config = [hoja_xls_iter for hoja_xls_iter, _ in lista_hojas_con_rangos_celdas]

        if not os.path.exists(path_xls):

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_log_1 = "No se ha conseguido abrir el fichero excel. Comprueba que la ruta existe o tienes acceso.\n"
            mensaje_log_2 = "No se realizan por lo tanto la captura de los screenshots de los rangos de celdas asociados a este excel \n"
            mensaje_log_3 = "ni la colocacion en el powerpoint de destino."

            dicc_error["PROCESO"] = opcion_interaccion_pptx
            dicc_error["FASE_PROCESO"] = "Generación del screenshot del rango de celdas excel via pdf usando Poppler"
            dicc_error["ID_PPTX"] = id_pptx_selecc
            dicc_error["ID_XLS"] = id_xls
            dicc_error["RUTA_FICHERO"] = path_xls
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2 + mensaje_log_3
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

            #se agrega el id xls cuyo path no se encuentra en la variable global lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx
            lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx.append(id_xls)



        else:

            app = xw.App(visible = False)
            wb = app.books.open(path_xls, update_links = update_links_xls)


            time.sleep(tiempo_espera_max_apertura_excel)

            #se des-ocultan las hojas del excel por si hay rangos que capturar en hojas ocultas
            for sheet in wb.sheets:

                try:
                    hoja_xls_iter_en_xls = sheet.name
                    sheet.visible = True


                except Exception as err_1:
                    #se registra el posible error de que el excel tiene hojas ocultas por codigo vba al abrir el excel
                    #siempre y cuando haya rangos de celdas configuradas en el app cuya hoja xls de origen
                    #sea este tipo de hoja oculta por vba

                    if hoja_xls_iter_en_xls in lista_hojas_xls_config:

                        traceback_error = traceback.extract_tb(err_1.__traceback__)
                        modulo_python = os.path.basename(traceback_error[0].filename)
                        rutina_python = traceback_error[0].name
                        linea_error = traceback_error[0].lineno

                        dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                        dicc_error["PROCESO"] = opcion_interaccion_pptx
                        dicc_error["FASE_PROCESO"] = "Hoja oculta en el excel, no se puede mostrar."
                        dicc_error["ID_PPTX"] = id_pptx_selecc
                        dicc_error["ID_XLS"] = id_xls
                        dicc_error["RUTA_FICHERO"] = path_xls
                        dicc_error["HOJA_XLS"] = hoja_xls_iter_en_xls
                        dicc_error["RANGO_CELDAS"] = None
                        dicc_error["SLIDE_PPTX"] = None
                        dicc_error["RESUMEN_ERROR"] = "La hoja no se puede mostrar posiblemente porque al abrir el excel hay una macro VBA que la oculta."
                        dicc_error["MODULO_PYTHON"] = modulo_python
                        dicc_error["RUTINA_PYTHON"] = rutina_python
                        dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                        dicc_error["TRACEBACK"] = str(err_1)

                        global_lista_dicc_errores.append(dicc_error)

                    pass


            #a modo recordatorio, el parametro kwargs lista_hojas_con_rangos_celdas es lista de listas donde cada sublista contiene:
            # --> hoja xls
            # --> lista de lista donde cada sublista contiene:
            #                                            --> slide pptx
            #                                            --> rango celdas xls
            #                                            --> nombre del screenshot
            for hoja_xls, lista_rangos_celdas in lista_hojas_con_rangos_celdas:

                for slide_pptx, rango_celdas_xls, nombre_screenshot in lista_rangos_celdas:

                    nombre_screenshot_pdf = os.path.join(directorio_screenshots, f"{nombre_screenshot}.pdf")
                    nombre_screenshot_png = os.path.join(directorio_screenshots, f"{nombre_screenshot}.png")

                    nombre_screenshot_pdf = os.path.normpath(nombre_screenshot_pdf)
                    nombre_screenshot_png = os.path.normpath(nombre_screenshot_png)


                    try:
                        ws = wb.sheets[hoja_xls]


                    except Exception as err_2:
                        #se registra el posible en la configuracion del area de impresion del rango de celdas excel
                        #se pone "pass" para que se registren todos los errores no solo el 1ero con el que se topa

                        traceback_error = traceback.extract_tb(err_2.__traceback__)
                        modulo_python = os.path.basename(traceback_error[0].filename)
                        rutina_python = traceback_error[0].name
                        linea_error = traceback_error[0].lineno

                        mensaje_log_1 = "No se puede acceder a la hoja excel. Comprueba que la hoja excel de origen del rango de celdas configurado existe en el excel de origen."

                        dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                        dicc_error["PROCESO"] = opcion_interaccion_pptx
                        dicc_error["FASE_PROCESO"] = "Generación del screenshot del rango de celdas excel via pdf usando Poppler"
                        dicc_error["ID_PPTX"] = id_pptx_selecc
                        dicc_error["ID_XLS"] = id_xls
                        dicc_error["RUTA_FICHERO"] = path_xls
                        dicc_error["HOJA_XLS"] = hoja_xls
                        dicc_error["RANGO_CELDAS"] = rango_celdas_xls
                        dicc_error["SLIDE_PPTX"] = slide_pptx
                        dicc_error["RESUMEN_ERROR"] = mensaje_log_1
                        dicc_error["MODULO_PYTHON"] = modulo_python
                        dicc_error["RUTINA_PYTHON"] = rutina_python
                        dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                        dicc_error["TRACEBACK"] = str(err_2)

                        global_lista_dicc_errores.append(dicc_error)

                        pass

                    else:

                        try:
                            ws.api.PageSetup.PrintArea = ""
                            rango_celdas_xls_objeto = ws.range(rango_celdas_xls)

                            #se configura el area de impresion del excel para que solo incluya el rango de celdas deseado
                            #y lo ajuste automaticamente a una sola página (1 pagina de ancho por 1 de alto)
                            #es necesario para que ExportAsFixedFormat genere un PDF limpio del rango
                            #independientemente de si el rango esta fuera de la zona visible en pantalla
                            ws.api.PageSetup.PrintArea = rango_celdas_xls_objeto.address
                            ws.api.PageSetup.Zoom = False
                            ws.api.PageSetup.FitToPagesWide = 1
                            ws.api.PageSetup.FitToPagesTall = 1



                        except Exception as err_3:
                            #se registra el posible en la configuracion del area de impresion del rango de celdas excel
                            #se pone "pass" para que se registren todos los errores no solo el 1ero con el que se topa

                            traceback_error = traceback.extract_tb(err_3.__traceback__)
                            modulo_python = os.path.basename(traceback_error[0].filename)
                            rutina_python = traceback_error[0].name
                            linea_error = traceback_error[0].lineno

                            mensaje_log_1 = "Error en la configuración del area de impresión del rango de celdas excel. "
                            mensaje_log_2 = f"Es posible que el motivo sea porque la apertura del excel supero el tiempo máximo de espera de apertura configurado ({tiempo_espera_max_apertura_excel}).\n"
                            mensaje_log_3 = "Las acciones siguientes se cancelan para este rango de celdas: exportacion a pdf, conversion a png via popper, recortes de los espacios 'blancos' en el png y colocación del screenshot en el pptx."

                            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                            dicc_error["PROCESO"] = opcion_interaccion_pptx
                            dicc_error["FASE_PROCESO"] = "Generación del screenshot del rango de celdas excel via pdf usando Poppler"
                            dicc_error["ID_PPTX"] = id_pptx_selecc
                            dicc_error["ID_XLS"] = id_xls
                            dicc_error["RUTA_FICHERO"] = path_xls
                            dicc_error["HOJA_XLS"] = hoja_xls
                            dicc_error["RANGO_CELDAS"] = rango_celdas_xls
                            dicc_error["SLIDE_PPTX"] = slide_pptx
                            dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2 + mensaje_log_3
                            dicc_error["MODULO_PYTHON"] = modulo_python
                            dicc_error["RUTINA_PYTHON"] = rutina_python
                            dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                            dicc_error["TRACEBACK"] = str(err_3)

                            global_lista_dicc_errores.append(dicc_error)

                            pass


                        else:
                            #si la configuracion del area de impresion del rango de celdas excel no da error se pasa a realizar la conversion de la hoja excel a pdf

                            try:
                                #se convierte la hoja excel que contiene el rango de celdas a pdf
                                time.sleep(time_sleep_prudencial_poppler)
                                ws.api.ExportAsFixedFormat(Type = 0, Filename = nombre_screenshot_pdf)
                                

                            except Exception as err_4:
                                #se registra el posible en la configuracion del area de impresion del rango de celdas excel
                                #se pone "pass" para que se registren todos los errores no solo el 1ero con el que se topa

                                traceback_error = traceback.extract_tb(err_4.__traceback__)
                                modulo_python = os.path.basename(traceback_error[0].filename)
                                rutina_python = traceback_error[0].name
                                linea_error = traceback_error[0].lineno

                                mensaje_log_1 = "Error en la conversión a pdf de la hoja excel que contiene el rango de celdas.\n"
                                mensaje_log_2 = "Las acciones siguientes se cancelan para este rango de celdas: conversion a png via popper, recortes de los espacios 'blancos' en el png y colocación del screenshot en el pptx."

                                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                                dicc_error["PROCESO"] = opcion_interaccion_pptx
                                dicc_error["FASE_PROCESO"] = "Generación del screenshot del rango de celdas excel via pdf usando Poppler"
                                dicc_error["ID_PPTX"] = id_pptx_selecc
                                dicc_error["ID_XLS"] = id_xls
                                dicc_error["RUTA_FICHERO"] = path_xls
                                dicc_error["HOJA_XLS"] = hoja_xls
                                dicc_error["RANGO_CELDAS"] = rango_celdas_xls
                                dicc_error["SLIDE_PPTX"] = slide_pptx
                                dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2
                                dicc_error["MODULO_PYTHON"] = modulo_python
                                dicc_error["RUTINA_PYTHON"] = rutina_python
                                dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                                dicc_error["TRACEBACK"] = str(err_4)

                                global_lista_dicc_errores.append(dicc_error)

                                pass

                            else:
                                #si la exportacion a pdf de la hoja excel que contiene el rango de celdas no da error
                                #se pasa a convertir el pdf a png pero solo el espacio que contiene el rango de celdas
                                #usando poppler

                                try:
                                    
                                    imagen_convert_from_path = convert_from_path(nombre_screenshot_pdf, dpi = dpi_pdf2image_convert_from_path, poppler_path = global_poppler_path)
                                    imagen_convert_from_path[0].save(nombre_screenshot_png)


                                except Exception as err_5:
                                    #se registran el posible error la captura del rango de celdas excel y su con su conversionb a png via pdf intermedio usando poppler
                                    #se pone "pass" para que se registren todos los errores no solo el 1ero con el que se topa

                                    traceback_error = traceback.extract_tb(err_5.__traceback__)
                                    modulo_python = os.path.basename(traceback_error[0].filename)
                                    rutina_python = traceback_error[0].name
                                    linea_error = traceback_error[0].lineno

                                    mensaje_log_1 = "Error en la conversión, mediante poppler, del pdf a png de la hoja excel que contiene el rango de celdas.\n"
                                    mensaje_log_2 = f"Comprueba que la ruta configurada de Poppler existe (puede que la hayan eliminado o la han movido de carpeta o le han cambiado el nombre):\n\n{global_poppler_path}\n\n"
                                    mensaje_log_3 = "Las acciones siguientes se cancelan para este rango de celdas: recortes de los espacios 'blancos' en el png y colocación del screenshot en el pptx."

                                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                                    dicc_error["PROCESO"] = opcion_interaccion_pptx
                                    dicc_error["FASE_PROCESO"] = "Generación del screenshot del rango de celdas excel via pdf usando Poppler"
                                    dicc_error["ID_PPTX"] = id_pptx_selecc
                                    dicc_error["ID_XLS"] = id_xls
                                    dicc_error["RUTA_FICHERO"] = path_xls
                                    dicc_error["HOJA_XLS"] = hoja_xls
                                    dicc_error["RANGO_CELDAS"] = rango_celdas_xls
                                    dicc_error["SLIDE_PPTX"] = slide_pptx
                                    dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2 + mensaje_log_3
                                    dicc_error["MODULO_PYTHON"] = modulo_python
                                    dicc_error["RUTINA_PYTHON"] = rutina_python
                                    dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                                    dicc_error["TRACEBACK"] = str(err_5)

                                    global_lista_dicc_errores.append(dicc_error)

                                    pass

                                else:
                                    #si la conversion mediante poppler del pdf a png no da error se pasa a recortar los espacios "blancos" en el png
                                    #para conservar tan solo el contenido del rango de celdas

                                    #la rutina def_acciones_screenshots_rangos_celdas_xls (opcion = RECORTAR_BLANCOS_SCREENSHOT_PNG) dispone de un bloque try except para
                                    #tambien registrar en la variable global global_lista_dicc_errores el posible error en los recortes de los espacios "blancos"

                                    #los parametros kwargs id_xls, hoja_xls y rango_celdas_xls se usan solo para registrar el posible error
                                    #el recorte se hace usando tan solo el parametro kwargs path_screenshot_png
                                    def_acciones_screenshots_rangos_celdas_xls("RECORTAR_BLANCOS_SCREENSHOT_PNG"
                                                                                , opcion_interaccion_pptx = opcion_interaccion_pptx
                                                                                , id_pptx_selecc = id_pptx_selecc
                                                                                , id_xls = id_xls
                                                                                , hoja_xls = hoja_xls
                                                                                , rango_celdas_xls = rango_celdas_xls
                                                                                , slide_pptx = slide_pptx
                                                                                , nombre_screenshot = nombre_screenshot
                                                                                , path_screenshot_png = nombre_screenshot_png)



            #se cierra el excel
            wb.close()
            app.quit()

            #se fuerza la salida del excel
            try:
                subprocess.run(["taskkill", "/PID", str(app.pid), "/F"], stdout = subprocess.DEVNULL, stderr = subprocess.DEVNULL)
            except:
                pass

            


def def_interaccion_pptx(opcion_pptx, **kwargs):
    #permite interactuar (configurar y generar) los pptx con los screenshots

    global global_proceso_en_ejecucion
    global global_poppler_path
    global global_ruta_local_config_sistema_sqlite
    global global_dicc_datos_id_pptx
    global global_dicc_check_screenshots_si_todo_blanco
    global lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx
    global global_lista_dicc_errores
    global global_lista_dicc_warning


    warnings.filterwarnings("ignore")

    resultado_funcion = None



    #parametros kwargs
    id_pptx_selecc = kwargs.get("id_pptx_selecc", None)
    directorio_pptx_destino = kwargs.get("directorio_pptx_destino", None)
    path_pptx_config = kwargs.get("path_pptx_config", None)
    opcion_interaccion_pptx = kwargs.get("opcion_interaccion_pptx", None)
    opcion_coordenadas_config_paso_1 = kwargs.get("opcion_coordenadas_config_paso_1", None)


    #AL INICIAR EL PROCESO: se inicializa la variable global global_lista_dicc_errores y global_lista_dicc_warning como listas vacias
    #se fija la variable global global_proceso_en_ejecucion para la rutina de los threads en la GUI y asi bloquear la ejecucion de cualquier otro proceso 
    #hasta que el actual en curso no haya finalizado
    global_lista_dicc_errores = []
    global_lista_dicc_warning = []
    global_proceso_en_ejecucion = "SI"



    #AL INICIAR EL PROCESO: se inicia la varaibel global lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx como lista vacia
    #cuando el aplicativo se conecta a los excel para realizar las capturas de rangos de celdas
    #si las rutas de los excels no se encuentran se almacena en esta lista (solo el id_xls)
    #esto permite cuando se realicen las colocaciones de los screenshots el pptx descartar los excels que estan en esta lista
    #si no se procede asi en el log de errores pueden salir errores en plan bola de nive cuando en realidad
    #el unico error el no entrar la ruta del excel
    lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx = []



    #########################################################################################
    # CONFIG_PASO_1 / EJECUCION
    #########################################################################################
    # --> se crea carpeta temporal en la misma ruta donde se ha seleccionado el pptx que se va a configurar (directorio_pptx_destino) para almacenar los screenshots
    # --> se abren los distintos excels configurados para el 'id_pptx' seleccionado aplicando el tiempo de espera max de espera
    #     para abrir el excel (en caso de que se supera se crea log de errores para informar de ello al usuario en la gui y asi poder adaptar sus configuraciones para el id_pptx)
    #            CONFIG_PASO_1 --> se abre cada excel sin actualizar vinculos
    #            EJECUCION     --> se abre cada excel segun que el usuario haya o no configurado actualizar vinculos
    #
    # --> se guardan los screenshot en .png (previamente en pdf) usando Poppler
    # --> se abre el fichero pptx (ruta_pptx_destino) y mediante loop sobre la tabla parametrica sqllite del id_pptx se colocan los screenshots en el pptx
    #            CONFIG_PASO_1 --> todos superpuestos en la esquina superior-izquierda en las slides configuradas
    #            EJECUCION     --> en las coordenadas y dimensiones configuradas por el usuario en el CONFIG_PASO_2
    #
    # --> se elimina la carpeta temporal mencionada
    # --> se guarda el pptx y se deja abierto para que el usuario realice sus ajustes manuales
    #            CONFIG_PASO_1 --> para poder re-dimensionar los screenshots copiados
    #            EJECUCION     --> para adaptar el pptx a los screenshots copiados y actualizados para el proceso de reporting
    #
    #el usuario tiene la opcion desde la GUI de optar por 2 tipos de colocaciones iniciales:
    # --> todos superpuestos en la esquina superior-izquierda en las slides configuradas (es la opción que escoje el aplicativo cuando se configura por 1era vez un powerpoint)
    #
    # --> si el powerpoint ya se configuro en otra sesión (CONFIFIG_PASO_2) y que las coordenadas y dimensiones de los pantallazos se guardaron en el sistema sqlite
    #     el usuario puede escojer colocar estos pantallazos con las coordenas y dimensiones que l mismo configuro en el powerpoint y todos los nuevos pantallazos
    #     se colocan en la esquina superior izquierda

    # se genera un log de errores en caso de que los hubiese

    if (opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_1"]["OPCION"]
        or opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["EJECUCION"]["OPCION"]):


        if not os.path.exists(global_ruta_local_config_sistema_sqlite):

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_log_1 = f"La ruta local siguiente no existe (han eliminado el directorio o lo han movido de carpeta o le han cambiado el nombre):\n\n{global_ruta_local_config_sistema_sqlite}\n\n"
            mensaje_log_2 = "Por lo tanto, el proceso se cancela."

            dicc_error["PROCESO"] = opcion_pptx
            dicc_error["FASE_PROCESO"] = "Acceso a la ruta del directorio local."
            dicc_error["ID_PPTX"] = None
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:

            now = dt.datetime.now()

            #se crea la carpeta temporal donde almacenar los screenshots (se usa la ruta almacenada en la variable global global_ruta_local_config_sistema_sqlite)
            directorio_screenshots = os.path.join(global_ruta_local_config_sistema_sqlite, nombre_carpeta_screenshots + "_" + str(re.sub("[^0-9a-zA-Z]+", "_", str(now))))
            directorio_screenshots = os.path.normpath(directorio_screenshots)

            if not os.path.exists(directorio_screenshots):
                os.makedirs(directorio_screenshots)


 
            #se recupera de la variable global global_dicc_datos_id_pptx el diccionario asociado al id pptx seleccionado para poder ejecutar el proceso
            dicc_datos_id_pptx = global_dicc_datos_id_pptx[id_pptx_selecc]

            ruta_plantilla_pptx = dicc_datos_id_pptx["ruta_plantilla_pptx"]
            tiempo_espera_max_apertura_pptx = dicc_datos_id_pptx["tiempo_espera_max_apertura_pptx"]
            lista_dicc_id_xls = dicc_datos_id_pptx["lista_dicc_id_xls"]




            #mediante iteracion sobre la lista lista_dicc_id_xls se ejecuta la rutina def_acciones_screenshots_rangos_celdas_xls (opcion --> SCREENSHOTS_XLS_A_PNG)
            #donde los screenshots de rangos de celdas se guardan en ficheros png en la carpeta local
            #los posibles errores se agregan a la variable global global_lista_dicc_errores para poder generar un log de errores en .txt
            #y avisar de ello en caso de que los hubiese en la GUI tras finalizar el proceso
            #
            #se ejecutan las acciones tan solo para los id xls asociados al id pptx seleccionado que tengan rangos de celdas configurados
            lista_columnas_df_rangos_celdas_filtrada = ["SLIDE_PPTX", "RANGO_XLS", "NOMBRE_SCREENSHOT"]

            lista_datos_screenshots_binarios = []

            if len(lista_dicc_id_xls) != 0:
                for dicc_id_xls in lista_dicc_id_xls:

                    id_xls = dicc_id_xls["id_xls"]
                    ruta_xls_origen = dicc_id_xls["ruta_xls_origen"]
                    tiempo_espera_max_apertura_xls = dicc_id_xls["tiempo_espera_max_apertura_xls"]
                    actualizar_vinculos_otros_xls = dicc_id_xls["actualizar_vinculos_otros_xls"]
                    df_widget_treeview_rangos_celdas = dicc_id_xls["df_widget_treeview_rangos_celdas"]

                    tiempo_espera_max_apertura_xls_ajust = int(float(tiempo_espera_max_apertura_xls))


                    #se crea la lista lista_hojas_xls_config que unifica todas las hojas xls configuradas para el id xls de la iteracion
                    lista_hojas_xls_config = ([df_widget_treeview_rangos_celdas.iloc[ind, df_widget_treeview_rangos_celdas.columns.get_loc("HOJA_XLS")] for ind in df_widget_treeview_rangos_celdas.index]
                                            if len(df_widget_treeview_rangos_celdas) != 0 else []
                                            )
                    
                    lista_hojas_xls_config = list(dict.fromkeys(lista_hojas_xls_config))

                    #se calcula la lista lista_hojas_con_rangos_celdas, es lista de listas donde cada sublista contiene:
                    # --> hoja xls
                    # --> lista de lista donde cada sublista contiene:
                    #                                            --> slide pptx
                    #                                            --> rango celdas xls
                    #                                            --> nombre del screenshot
                    lista_hojas_con_rangos_celdas = []
                    if len(lista_hojas_xls_config) != 0:
                        for hoja_xls_config in lista_hojas_xls_config:

                            df_rangos_celdas = df_widget_treeview_rangos_celdas.loc[df_widget_treeview_rangos_celdas["HOJA_XLS"] == hoja_xls_config, lista_columnas_df_rangos_celdas_filtrada]
                            lista_rangos_celdas = df_rangos_celdas.values.tolist()

                            if len(lista_rangos_celdas) != 0:
                                lista_hojas_con_rangos_celdas.append([hoja_xls_config, lista_rangos_celdas])



                    #se actualiza la variable global global_dicc_check_screenshots_si_todo_blanco
                    #sus keys son los nombres de los screenshots de los rangos de celdas configurados para el id pptx
                    #sus valores son dicionarios con las keys siguientes:
                    # --> id_xls            
                    # --> hoja_xls
                    # --> rango_celdas
                    # --> slide_pptx
                    # --> todo_en_blanco       se establece a True de inicio, pasa a false dentro de la rutina def_acciones_screenshots_rangos_celdas_xls
                    #                          (bloque siguiente) si el screenshot capurado es todo en blanco
                    for hoja_xls, lista_screenshots_iter in lista_hojas_con_rangos_celdas:
                        for slide_pptx, rango_celdas, nombre_screenshot in lista_screenshots_iter:

                            dicc_temp = {str(nombre_screenshot):
                                                                {"id_xls": id_xls
                                                                , "hoja_xls": hoja_xls
                                                                , "rango_celdas": rango_celdas
                                                                , "slide_pptx": slide_pptx
                                                                , "todo_en_blanco": True
                                                                }
                                        }
                            
                            global_dicc_check_screenshots_si_todo_blanco.update(dicc_temp)




                    #se realizan los screenshots de rangos xls usando poppler tan solo si lista_hojas_con_rangos_celdas no es vacia
                    if len(lista_hojas_con_rangos_celdas) != 0:

                        def_acciones_screenshots_rangos_celdas_xls("SCREENSHOTS_XLS_A_PNG"
                                                                    , opcion_interaccion_pptx = opcion_pptx
                                                                    , id_pptx_selecc = id_pptx_selecc
                                                                    , id_xls = id_xls
                                                                    , path_xls = ruta_xls_origen
                                                                    , tiempo_espera_max_apertura_excel = tiempo_espera_max_apertura_xls_ajust
                                                                    , opcion_update_links = actualizar_vinculos_otros_xls
                                                                    , lista_hojas_con_rangos_celdas = lista_hojas_con_rangos_celdas
                                                                    , directorio_screenshots = directorio_screenshots)
                        





                    #se almacenan los datos en la lista lista_datos_para_sitema_sqlite que permiten guardar en el sistema sqlite
                    #las imagenes binarias de los png creados (se usan en la GUI como muestra en la ventana de configuracion de los id pptx)
                    #(solo para el proceso CONFIG_PASO_1)
                    #se realiza el calculo tan solo si lista_hojas_con_rangos_celdas no es vacia
                    if len(lista_hojas_con_rangos_celdas) != 0:

                        if opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_1"]["OPCION"]:

                            for hoja_xls, lista_rangos_celdas in lista_hojas_con_rangos_celdas:

                                for slide_pptx, rango_celdas_xls, nombre_screenshot in lista_rangos_celdas:

                                    screenshot_path_png = os.path.join(directorio_screenshots, f"{nombre_screenshot}.PNG")
                                    screenshot_path_png = os.path.normpath(screenshot_path_png)

                                    if os.path.isfile(screenshot_path_png):

                                        with open(screenshot_path_png, "rb") as file:
                                            screenshot_binary = file.read()

                                            lista_datos_screenshots_binarios.append([id_pptx_selecc
                                                                                        , id_xls
                                                                                        , hoja_xls
                                                                                        , rango_celdas_xls
                                                                                        , slide_pptx
                                                                                        , nombre_screenshot
                                                                                        , screenshot_binary
                                                                                    ])
                                            


                    #se guardan las imagenes binarias de los screenshot en el sistema sqlite (solo para el proceso CONFIG_PASO_1)
                    #se realiza el calculo tan solo si lista_hojas_con_rangos_celdas no es vacia
                    if len(lista_hojas_con_rangos_celdas) != 0:

                        if opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_1"]["OPCION"]:

                            if len(lista_datos_screenshots_binarios) != 0:

                                def_config_sistema_sqlite("UPDATE_SCREENSHOTS_IMAGENES_BINARIAS"
                                                            , id_pptx_selecc = id_pptx_selecc
                                                            , lista_datos_screenshots_binarios = lista_datos_screenshots_binarios)





            #se informa la variable global global_lista_dicc_warning en caso de que hayan screenshots todo en blanco
            #sus keys son los nombres de los screenshots de los rangos de celdas configurados para el id pptx
            #sus valores son dicionarios con las keys siguientes:
            # --> id_xls            
            # --> hoja_xls
            # --> rango_celdas
            # --> slide_pptx
            # --> todo_en_blanco       se filtra tan solo los que pone true
            for nombre_screenshot, dicc_screenshot_iter in global_dicc_check_screenshots_si_todo_blanco.items():

                screenshot_todo_en_blanco = dicc_screenshot_iter["todo_en_blanco"]

                if screenshot_todo_en_blanco:

                    id_xls_iter_dicc = dicc_screenshot_iter["id_xls"]
                    hoja_xls_iter_dicc = dicc_screenshot_iter["hoja_xls"]
                    rango_celdas_iter_dicc = dicc_screenshot_iter["rango_celdas"]
                    slide_pptx_iter_dicc = dicc_screenshot_iter["slide_pptx"]

                    #si el id xls de la iteracion eesta en la lista id_xls_iter_dicc not in lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx
                    #no se registra el log (el error no es el que sale a continuacion para los rangos de celdas de este id xls sino pq no se encontro el path del excel
                    #y este log se registro en global_lista_dicc_warning al inicio del proceso)
                    if id_xls_iter_dicc not in lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx:

                        dicc_warning = dicc_modelo_warning_y_errores_procesos_app["WARNING"].copy()

                        mensaje_log_1 = f"El pantallazo ({nombre_screenshot}) ha salido todo en blanco\n.Comprueba en el excel asociado si el rango de celdas configurado:\n"
                        mensaje_log_2 = f"--> esta ubicado en filas y/o columnas ocultas.\n"
                        mensaje_log_3 = f"--> tiene su contenido todo en blanco.\n\n"
                        mensaje_log_4 = "El fallo puede estar debido a que el rango de celdas lo tienes configurado\nen el aplicativo con una slide powerpoint de destino que excede\n"
                        mensaje_log_5 = "el número total de slides que contiene dicho powerpoint\n(revisa las configuraciones de rangos de celdas)"

                        dicc_warning["PROCESO"] = opcion_pptx
                        dicc_warning["FASE_PROCESO"] = "Recorte espacios 'blancos' del fichero png asociado al screenshot de rango de celdas."
                        dicc_warning["ID_PPTX"] = id_pptx_selecc
                        dicc_warning["ID_XLS"] = id_xls_iter_dicc
                        dicc_warning["HOJA_XLS"] = hoja_xls_iter_dicc
                        dicc_warning["RANGO_CELDAS"] = rango_celdas_iter_dicc
                        dicc_warning["SLIDE_PPTX"] = slide_pptx_iter_dicc
                        dicc_warning["COMENTARIO"] = mensaje_log_1 + mensaje_log_2 + mensaje_log_3 + mensaje_log_4 + mensaje_log_5

                        global_lista_dicc_warning.append(dicc_warning)






            #se fusionan todos los df de las keys df_widget_treeview_rangos_celdas (dentro de la key id_xls)
            #con el df resultante se crea una lista de listas por slide pptx donde en cada una se asocian todos los nombres
            #de screenshots a copiar con sus respectivas coordenadas
            #en caso de que el proceso sea:
            # --> CONFIG_PASO_1      se asignan las coordenadas definidas en tupla_coordenadas_colocacion_pptx_por_defecto
            #                        (que colocan los swcreenshots en el pptx en la esquina superior-izquierda)   
            # --> EJECUCION          se usan las coordenadas guardadas en el sistema sqlite
            lista_columnas_extraccion_df = ["ID_XLS", "HOJA_XLS", "RANGO_XLS", "SLIDE_PPTX", "NOMBRE_SCREENSHOT", "COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX"]
            lista_columnas_extraccion_df_ajust = ["ID_XLS", "HOJA_XLS", "RANGO_XLS", "SLIDE_PPTX", "NOMBRE_SCREENSHOT", "COORDENADAS_POR_APLICAR"]
            lista_columnas_extraccion_df_ajust_sin_slide_pptx = ["ID_XLS", "HOJA_XLS", "RANGO_XLS", "NOMBRE_SCREENSHOT", "COORDENADAS_POR_APLICAR"]


            df_screenshots_por_id_pptx = pd.DataFrame(columns = lista_columnas_extraccion_df)
            for ind, dicc_id_xls in enumerate(lista_dicc_id_xls):

                if ind == 0:
                    df_screenshots_por_id_pptx = dicc_id_xls["df_widget_treeview_rangos_celdas"][lista_columnas_extraccion_df]
                else:
                    df_screenshots_por_id_pptx = pd.concat([df_screenshots_por_id_pptx, dicc_id_xls["df_widget_treeview_rangos_celdas"][lista_columnas_extraccion_df]])

            df_screenshots_por_id_pptx.reset_index(drop = True, inplace = True)



            #se crea el campo COORDENADAS_SCREENSHOT_POR_DEFECTO en el df df_screenshots_por_id_pptx
            #para almacenar las coordenadas a aplicar segun el proceso que se ejecuta en la rutina (CONFIG_PASO_1 o EJECUCION)
            str_tupla_coordenadas_por_defecto = str(tupla_coordenadas_colocacion_pptx_por_defecto)


            if opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_1"]["OPCION"]:

                if opcion_coordenadas_config_paso_1 == "COLOCACION_INICIAL":
                    df_screenshots_por_id_pptx["COORDENADAS_POR_APLICAR"] = df_screenshots_por_id_pptx.apply(lambda x: str_tupla_coordenadas_por_defecto, axis = 1)
                    

                
                elif opcion_coordenadas_config_paso_1 == "COLOCACION_HIBRIDA":
            
                    df_screenshots_por_id_pptx["COORDENADAS_POR_APLICAR"] = df_screenshots_por_id_pptx.apply(lambda x: x["COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX"] if x["COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX"] is not None
                                                                                                                        else
                                                                                                                        str_tupla_coordenadas_por_defecto
                                                                                                            , axis = 1)
                

            elif opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["EJECUCION"]["OPCION"]:
                df_screenshots_por_id_pptx["COORDENADAS_POR_APLICAR"] = df_screenshots_por_id_pptx.apply(lambda x: x["COORDENADAS_SCREENSHOT_EN_SLIDE_PPTX"], axis = 1)     
            


            df_screenshots_por_id_pptx = df_screenshots_por_id_pptx[lista_columnas_extraccion_df_ajust]



            #se crea la lista lista_screenshots_por_id_pptx, es lista de listas donde cada sublista contiene:
            # --> slide pptx
            # --> lista de listas donde cada sublista contiene:
            #        --> id xls
            #        --> id_xls_path
            #        --> hoja xls
            #        --> rango celdas
            #        --> nombre screenshot
            #        --> string de la tupla de coordenadas a aplicar
            lista_screenshots_por_id_pptx = [[df_screenshots_por_id_pptx.iloc[ind, df_screenshots_por_id_pptx.columns.get_loc("SLIDE_PPTX")]
                                            , None #se informa mas adelante (almacena una lista de listas donde cada sublista contiene el nombre del screenshot y sus coordenadas)
                                            ]
                                            for ind in df_screenshots_por_id_pptx.index
                                            ]
            
            lista_screenshots_por_id_pptx = [sublista for ind, sublista in enumerate(lista_screenshots_por_id_pptx) if sublista not in lista_screenshots_por_id_pptx[:ind]]


            for ind_lista, item_lista in enumerate(lista_screenshots_por_id_pptx):
                slide_pptx = item_lista[0]

                df_temp = df_screenshots_por_id_pptx.loc[df_screenshots_por_id_pptx["SLIDE_PPTX"] == int(slide_pptx), lista_columnas_extraccion_df_ajust_sin_slide_pptx]
                lista_screenshots_con_coordenadas = df_temp.values.tolist()

                lista_screenshots_por_id_pptx[ind_lista][1] = lista_screenshots_con_coordenadas




            #se realiza la accion de colocar los screenshots en el pptx
            #se registra en la variable global global_lista_dicc_errores si no se consigue acceso al pptx de destino
            ruta_plantilla_pptx_ajust = os.path.normpath(ruta_plantilla_pptx)

            if not os.path.exists(ruta_plantilla_pptx_ajust):

                #el error de conexion al fichero pptx no se hace mediante try except
                #pq no es un error del codigo simplemente que el usuario no tiene acceso al fichero
                dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                mensaje_log_1 = f"No se ha conseguido acceder al powerpoint de destino:\n\n{ruta_plantilla_pptx_ajust}.\n\n"
                mensjae_log_2 = "El fichero o bien se ha eliminado o lo han movido de carpeta o le han cambiado el nombre.\nLas acciones de colocación de los screenshots se cancelan.\n"
                mensaje_log_3 = "Tampoco se puede abrir el fichero powerpoint."

                dicc_error["PROCESO"] = opcion_pptx
                dicc_error["FASE_PROCESO"] = "Apertura del powerpoint de destino."
                dicc_error["ID_PPTX"] = id_pptx_selecc
                dicc_error["ID_XLS"] = None
                dicc_error["RUTA_FICHERO"] = None
                dicc_error["HOJA_XLS"] = None
                dicc_error["RANGO_CELDAS"] = None
                dicc_error["SLIDE_PPTX"] = None
                dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensjae_log_2 + mensaje_log_3
                dicc_error["MODULO_PYTHON"] = None
                dicc_error["RUTINA_PYTHON"] = None
                dicc_error["LINEA_CODIGO_PYTHON"] = None
                dicc_error["TRACEBACK"] = None

                global_lista_dicc_errores.append(dicc_error)


                #resulatado de la funcion
                resultado_funcion = None


            else:
                #se copia el fichero pptx de destino configurado en el sistema en la truta donde se quiere realizar el proceso
                nombre_fichero_pptx = nombre_pptx_tras_config if opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_1"]["OPCION"] else nombre_pptx_tras_ejecucion

                ruta_pptx_destino = os.path.join(directorio_pptx_destino, nombre_fichero_pptx + "_" + str(re.sub("[^0-9a-zA-Z]+", "_", str(now))) + ".pptx")
                ruta_pptx_destino = os.path.normpath(ruta_pptx_destino)



                shutil.copyfile(ruta_plantilla_pptx_ajust, ruta_pptx_destino)



                #se abre el pptx con el tiempo de espera configurado
                pptx_objeto = Presentation(ruta_pptx_destino)
                time.sleep(tiempo_espera_max_apertura_pptx)




                #se eliminan todos los shapes del pptx que contengan en su nombre el definido en la variable nomenclatura_nombres_screenshots
                #(en minusculas y trimeado)
                try:
                    for slide in pptx_objeto.slides:
                        lista_shapes_por_eliminar = [shape for shape in slide.shapes if nomenclatura_nombres_screenshots.strip().lower() in shape.name.strip().lower()]

                        if len(lista_shapes_por_eliminar) != 0:
                            for shape in lista_shapes_por_eliminar:
                                slide.shapes._spTree.remove(shape._element)


                except Exception as _:

                    #aunque el registro del error en la variable global global_lista_dicc_errores se haga mediante bloque try except
                    #aqui no se informa el traceback pq en realidad no es un error de codigo sino que no se ha podido acceder en memoria al objeto pptx
                    #en el tiempo maximo configurado, el bloque try except es el que permite localizar estos casos
                    mensaje_log_1 = f"No se ha conseguido acceder al fichero pptx pasado en el tiempo máximo de espera de apertura configurado ({tiempo_espera_max_apertura_pptx}).\n"
                    mensaje_log_2 = "Revisa la configuración de este tiempo de espera máximo aumentandolo."

                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                    dicc_error["PROCESO"] = opcion_pptx
                    dicc_error["FASE_PROCESO"] = "Conexión a los datos del powerpoint de destino."
                    dicc_error["ID_PPTX"] = id_pptx_selecc
                    dicc_error["ID_XLS"] = None
                    dicc_error["RUTA_FICHERO"] = None
                    dicc_error["HOJA_XLS"] = None
                    dicc_error["RANGO_CELDAS"] = None
                    dicc_error["SLIDE_PPTX"] = None
                    dicc_error["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2
                    dicc_error["MODULO_PYTHON"] = None
                    dicc_error["RUTINA_PYTHON"] = None
                    dicc_error["LINEA_CODIGO_PYTHON"] = None
                    dicc_error["TRACEBACK"] = None

                    global_lista_dicc_errores.append(dicc_error)

                    pass
            
                else:

                    #se recupera el numero total de slides que contiene el pptx de destino
                    #para poder almacenar en la variable global global_lista_dicc_errores
                    #un log de error si el numero de slide pptx configurado para el rango de celdas
                    #excede el numero total de slides
                    numero_total_slides = int(float(global_dicc_datos_id_pptx[id_pptx_selecc]["numero_total_slides"]))


                    #se copian los screenshots en el pptx en cada slide incluida en lista_screenshots_por_id_pptx
                    for slide_pptx_iter, lista_screenshots_iter in lista_screenshots_por_id_pptx:

                        #a modo recordatorio el 2ndo item (lista_screenshots_iter) de la lista lista_screenshots_por_id_pptx es lista que contiene:
                        # --> id xls
                        # --> hoja xls
                        # --> rango celdas
                        # --> nombre screenshot
                        # --> string de la tupla de coordenadas a aplicar
            
                        ###############################################################################################
                        #CASO 1 - el numero de slide pptx de la iteracion no excede el numero total de slides del pptx
                        ###############################################################################################
                        if slide_pptx_iter <= numero_total_slides:

                            slide_pptx_objeto = pptx_objeto.slides[int(slide_pptx_iter) - 1]

                            for id_xls_iter, hoja_xls_iter, rango_celdas_iter, nombre_screenshot_iter, coordenadas_iter in lista_screenshots_iter:

                                #se ejecuta la colocacion tan solo si el excel de origen existe
                                #al principio del proceso hay log de error que localiza los excel de origen existen
                                #y para los calculos a nivel de captura de los screneshots hasta conversion en png pero no frena el proceso actual de colocacion en un pptx
                                #NO HACE FALTA almacenar el log en la variablñevariable global global_lista_dicc_errores pq ya se hizo al principio del proceso
                                #se usa la variable global lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx
                                if id_xls_iter not in lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx:

                                    try:
                                        path_screenshot = os.path.join(directorio_screenshots, f"{nombre_screenshot_iter}.PNG")
                                        path_screenshot = os.path.normpath(path_screenshot)

                                        _, tupla_coordenadas = def_varios("CHECK_TIPO_DATO_PYTHON", valor_check_tipo_dato = coordenadas_iter)

                                        shape = slide_pptx_objeto.shapes.add_picture(path_screenshot
                                                                                    , left = tupla_coordenadas[0]
                                                                                    , top =  tupla_coordenadas[1]
                                                                                    , width =  tupla_coordenadas[2]
                                                                                    , height = tupla_coordenadas[3])
                                        
                                        shape.name = nombre_screenshot_iter

                                    except Exception as err_1:
                                        #se registra el posible error en la colococacion del screenshot

                                        traceback_error = traceback.extract_tb(err_1.__traceback__)
                                        modulo_python = os.path.basename(traceback_error[0].filename)
                                        rutina_python = traceback_error[0].name
                                        linea_error = traceback_error[0].lineno

                                        dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                                        mensaje_error_1 = f"No se ha podido colocar el pantallazo asociado al rango de celdas en la slide de destino del powerpoint.\n"

                                        mensaje_error_2_1 = f"La tupla de coordenadas que se intenta aplicar no existe.\n"
                                        mensaje_error_2_2 = f"Es posible que fuera que el rango de celdas cuando realizaste 'Configurar (paso 1)' saliese en blanco.\n"
                                        mensaje_error_2_3 = "Revisa el log de warning que se ejecuto entonces."              
                                        mensaje_error_2 = mensaje_error_2_1 + mensaje_error_2_2 + mensaje_error_2_3 if coordenadas_iter == "None" else ""#pdte corregir el pq sale None es string

                                        dicc_error["PROCESO"] = opcion_pptx
                                        dicc_error["FASE_PROCESO"] = "Colocación del pantallazo como shape en el powerpoint."
                                        dicc_error["ID_PPTX"] = id_pptx_selecc
                                        dicc_error["ID_XLS"] = id_xls_iter
                                        dicc_error["RUTA_FICHERO"] = None
                                        dicc_error["HOJA_XLS"] = hoja_xls_iter
                                        dicc_error["RANGO_CELDAS"] = rango_celdas_iter
                                        dicc_error["SLIDE_PPTX"] = rango_celdas_iter
                                        dicc_error["RESUMEN_ERROR"] = mensaje_error_1 + mensaje_error_2
                                        dicc_error["MODULO_PYTHON"] = modulo_python
                                        dicc_error["RUTINA_PYTHON"] = rutina_python
                                        dicc_error["LINEA_CODIGO_PYTHON"] = linea_error
                                        dicc_error["TRACEBACK"] = str(err_1)

                                        global_lista_dicc_errores.append(dicc_error)

                                        pass

                        ###############################################################################################
                        #CASO 2 - el numero de slide pptx de la iteracion SI excede el numero total de slides del pptx
                        #mediante bucle sobre la lista lista_screenshots_iter se agrega un diccionario en cada iteracion
                        #a la variable global global_lista_dicc_errores
                        #no se hace con bloque try except pq aqui no es error de codigo sino de configuracion del app
                        ###############################################################################################
                        elif slide_pptx_iter > numero_total_slides:

                            for id_xls_iter, hoja_xls_iter, rango_celdas_xls_iter, nombre_screenshot_iter, coordenadas_iter in lista_screenshots_iter:

                                #se ejecuta la colocacion tan solo si el excel de origen existe
                                #al principio del proceso hay log de error que localiza los excel de origen existen
                                #y para los calculos a nivel de captura de los screneshots hasta conversion en png pero no frena el proceso actual de colocacion en un pptx
                                #NO HACE FALTA almacenar el log en la variablñevariable global global_lista_dicc_errores pq ya se hizo al principio del proceso
                                #se usa la variable global lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx
                                if id_xls_iter not in lista_id_xls_con_path_inexistente_ejecucion_interaccion_pptx:

                                    dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

                                    dicc_error["PROCESO"] = opcion_pptx
                                    dicc_error["FASE_PROCESO"] = "Colocación del pantallazo como shape en el powerpoint."
                                    dicc_error["ID_PPTX"] = id_pptx_selecc
                                    dicc_error["ID_XLS"] = id_xls_iter
                                    dicc_error["RUTA_FICHERO"] = None
                                    dicc_error["HOJA_XLS"] = hoja_xls_iter
                                    dicc_error["RANGO_CELDAS"] = rango_celdas_iter
                                    dicc_error["SLIDE_PPTX"] = slide_pptx_iter
                                    dicc_error["RESUMEN_ERROR"] = f"La slide de destino ({slide_pptx_iter}) es superior al número de slides que contiene el pptx ({numero_total_slides})."
                                    dicc_error["MODULO_PYTHON"] = None
                                    dicc_error["RUTINA_PYTHON"] = None
                                    dicc_error["LINEA_CODIGO_PYTHON"] = None
                                    dicc_error["TRACEBACK"] = None

                                    global_lista_dicc_errores.append(dicc_error)


                    #se guarda el pptx
                    pptx_objeto.save(ruta_pptx_destino)


                    #resultado de la funcion
                    resultado_funcion = ruta_pptx_destino




       
            #se elimina la carpeta temporal con los screenshots
            if os.path.exists(directorio_screenshots):
                try:
                    shutil.rmtree(directorio_screenshots)
                except:
                    pass#por si alguien esta manipulando ficheros durante el proceso




    #########################################################################################
    # CONFIG_PASO_2
    #########################################################################################
    #se conecta al pptx donde se localizan todos los screenshots re-dimensionados por el usuario
    #y se guardan las coordenadas en forma de tupla convertida a string en el sistema sqlite

    elif opcion_pptx == dicc_gui_combobox_procesos["COMBOBOX_PPTX"]["CONFIG_PASO_2"]["OPCION"]:

        #se crea la lista lista_screenshots_con_coordenadas que lista de lista de listas donde cadasublista contiene para el idpptx seleccionado
        # --> nombre screenshot
        # --> id xls
        # --> hoja xls
        # --> rango celdas xls
        # --> slide pptx
        # --> coordenadas del screensshot en el pptx (tupla convertida a string)
        lista_columnas_df_widget_treeview_rangos_celdas_filtrada = ["NOMBRE_SCREENSHOT", "ID_XLS", "HOJA_XLS", "RANGO_XLS", "SLIDE_PPTX"]

        dicc_id_pptx = global_dicc_datos_id_pptx[id_pptx_selecc]
        lista_dicc_id_xls = dicc_id_pptx["lista_dicc_id_xls"]

        lista_screenshots_con_coordenadas = []
        if len(lista_dicc_id_xls) != 0:

            for dicc in lista_dicc_id_xls:
                df_widget_treeview_rangos_celdas = dicc["df_widget_treeview_rangos_celdas"]

                if len(df_widget_treeview_rangos_celdas) != 0:
                    lista_datos_df_widget_treeview_rangos_celdas = df_widget_treeview_rangos_celdas[lista_columnas_df_widget_treeview_rangos_celdas_filtrada].values.tolist()

                    for nombre_screenshot_iter, id_xls_iter, hoja_xls_iter, rango_celdas_xls_iter, slide_pptx_iter in lista_datos_df_widget_treeview_rangos_celdas:

                        lista_screenshots_con_coordenadas.append([nombre_screenshot_iter
                                                                    , id_xls_iter
                                                                    , hoja_xls_iter
                                                                    , rango_celdas_xls_iter
                                                                    , slide_pptx_iter
                                                                    , None #se usa para poner la tupla de coordenadas (convertida) a string - secalcula en el bloque siguiente
                                                                    ])
                        


        #se abre el pptx donde se configuraron los screenshots
        tiempo_espera_max_apertura_pptx = global_dicc_datos_id_pptx[id_pptx_selecc]["tiempo_espera_max_apertura_pptx"]


        if not os.path.exists(path_pptx_config):
        #se registra el error por si el pptx

            dicc_error = dicc_modelo_warning_y_errores_procesos_app["ERROR"].copy()

            mensaje_error = f"La ruta del powerpoint de destino no existe (han eliminado el fichero o lo han movido de carpeta o le han cambiado el nombre):\n\n{path_pptx_config}"

            dicc_error["PROCESO"] = opcion_interaccion_pptx
            dicc_error["FASE_PROCESO"] = "Conexión al powerpoint para recuperar las coordenadas y dimensiones de los pantallazos."
            dicc_error["ID_PPTX"] = id_pptx_selecc
            dicc_error["ID_XLS"] = None
            dicc_error["RUTA_FICHERO"] = None
            dicc_error["HOJA_XLS"] = None
            dicc_error["RANGO_CELDAS"] = None
            dicc_error["SLIDE_PPTX"] = None
            dicc_error["RESUMEN_ERROR"] = mensaje_error
            dicc_error["MODULO_PYTHON"] = None
            dicc_error["RUTINA_PYTHON"] = None
            dicc_error["LINEA_CODIGO_PYTHON"] = None
            dicc_error["TRACEBACK"] = None

            global_lista_dicc_errores.append(dicc_error)

        else:

            pptx_objeto = Presentation(path_pptx_config)
            time.sleep(tiempo_espera_max_apertura_pptx)


            #se informa la posicion 6 de las sublistas de lista_screenshots_con_coordenadas
            ind_pos_nombre_screenshot = 0
            ind_pos_coordenadas = 5
            for ind, item in enumerate(lista_screenshots_con_coordenadas):

                nombre_screenshot_ajust_iter = item[ind_pos_nombre_screenshot].lower().strip()

                for slide in pptx_objeto.slides:

                    for shape in slide.shapes:
                        nombre_screenshot_pptx_objeto = shape.name
                        nombre_screenshot_pptx_objeto_ajust = nombre_screenshot_pptx_objeto.lower().strip()

                        if nombre_screenshot_ajust_iter == nombre_screenshot_pptx_objeto_ajust:
                            lista_screenshots_con_coordenadas[ind][ind_pos_coordenadas] = str((shape.left, shape.top, shape.width, shape.height))

                            break

                
            #se almacena en la variable global global_lista_dicc_warning
            #los screenshots donde no se han localizado coordenadas
            for nombre_screenshot_iter, id_xls_iter, hoja_xls_iter, rango_celdas_xls_iter, slide_pptx_iter, coordenadas_iter in lista_screenshots_con_coordenadas:

                if coordenadas_iter is None:

                    dicc_warning = dicc_modelo_warning_y_errores_procesos_app["WARNING"].copy()

                    mensaje_log_1 = f"No se ha localizado ningún shape en el pptx con el nombre siguiente: {nombre_screenshot_iter}.\n"
                    mensaje_log_2 = "Lo más probable es que cuando ejecuste 'Configurar (paso 1)', los pantallazos salieron todo en blanco y se te informo en el log de warnings.\n"
                    mensaje_log_3 = "Tambien puede estar debido a que el rango de celdas lo tienes configurado en el aplicativo con una slide powerpoint\n"
                    mensaje_log_4 = "de destino que excede el número total de slides que contiene dicho powerpoint (revisa las configuraciones de rangos de celdas)"

                    dicc_warning["PROCESO"] = opcion_interaccion_pptx
                    dicc_warning["FASE_PROCESO"] = f"Localización de las coordenadas y dimensiones del shape en el pptx asociado al screenshot de rangos de celdas."
                    dicc_warning["ID_PPTX"] = id_pptx_selecc
                    dicc_warning["ID_XLS"] = id_xls_iter
                    dicc_warning["RUTA_FICHERO"] = None
                    dicc_warning["HOJA_XLS"] = hoja_xls_iter
                    dicc_warning["RANGO_CELDAS"] = rango_celdas_xls_iter
                    dicc_warning["SLIDE_PPTX"] = slide_pptx_iter
                    dicc_warning["RESUMEN_ERROR"] = mensaje_log_1 + mensaje_log_2 + mensaje_log_3 + mensaje_log_4
                    dicc_warning["MODULO_PYTHON"] = None
                    dicc_warning["RUTINA_PYTHON"] = None
                    dicc_warning["LINEA_CODIGO_PYTHON"] = None
                    dicc_warning["TRACEBACK"] = None

                    global_lista_dicc_warning.append(dicc_warning)   



            #se actualiza el sistema sqlite con las coordenadas de cada screenshot
            if len(lista_screenshots_con_coordenadas) != 0:
                def_config_sistema_sqlite("UPDATE_SCREENSHOTS_COORDENADAS"
                                        , id_pptx_selecc = id_pptx_selecc
                                        , lista_screenshots_con_coordenadas = lista_screenshots_con_coordenadas)


            #se importa de nuevo el sistema sqlite para poder tenerlo en memoria con los cambios de coordenadas
            def_config_sistema_sqlite("IMPORTAR_SISTEMA", path_sqlite = global_path_sistema_sqlite)
            def_varios("DICC_DATOS_ID_PPTX")





    #AL FINALIZAR EL PROCESO: se inicializa la variable global global_proceso_en_ejecucion que se usa
    #para impedir ejecutar en la GUI otro proceso hasta que acabe el que esta en curso
    #
    #NO HAY QUE REINICIAR ni global_lista_dicc_errores ni global_lista_dicc_warning
    #sino no se genera el log de posibles errores / warning en la GUI
    global_proceso_en_ejecucion = "NO"





    #resultado funcion
    return resultado_funcion


